"""多格式文档解析工具。

功能：
    compare  — 多解析器对比（默认）
    convert  — 文档转 md/txt/json
    batch    — 批量转换
    scan-dir — 批量质量扫描
    tables   — 表格提取（仅 PDF）
    chunk    — 文本分块
    render-pages — PDF 页面截图
    ocr      — 可选真实 OCR
    qa       — 基于知识包或文档的带引用抽取式问答
    diff-docs — 文档版本差异对比

支持格式：
    PDF  (.pdf)  — markitdown / pymupdf4llm / docling / pspdfkit / pymupdf / pypdf / pdfplumber / pdfminer / liteparse
    Word (.docx) — markitdown / python-docx
    PPT  (.pptx) — markitdown / python-pptx
    Excel(.xlsx) — markitdown / openpyxl
    HTML (.html) — markitdown / beautifulsoup4

用法：
    python parse_pdf_compare.py compare <文件> [选项]
    python parse_pdf_compare.py convert <文件> --parser pymupdf --format json -o output.json
    python parse_pdf_compare.py batch <目录> --parser pymupdf --output-dir ./out/
    python parse_pdf_compare.py tables <文件> --pages 1-5 --format csv
    python parse_pdf_compare.py <文件> [选项]  # 等同于 compare
"""
from __future__ import annotations

__version__ = "4.13.0"

import argparse
import csv
import difflib
import hashlib
import html
import importlib.util
import io
import json
import logging
import re
import shutil
import subprocess
import sys
import tempfile
import time
from collections import Counter
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import asdict, dataclass
from decimal import Decimal, InvalidOperation
from pathlib import Path
from typing import Callable

logger = logging.getLogger(__name__)

# 预编译正则
_RE_WHITESPACE = re.compile(r"[ \t]+")
_RE_MULTI_NEWLINE = re.compile(r"\n{3,}")
_RE_INLINE_WHITESPACE = re.compile(r"\s+")
_RE_PARAGRAPH_SPLIT = re.compile(r"\n\s*\n")
_RE_CHINESE = re.compile(r"[一-鿿㐀-䶿\U00020000-\U0002a6df]")
_RE_TOKEN = re.compile(r"[A-Za-z0-9_]+|[一-鿿]")

QUALITY_BAD_LABELS = {"empty", "bad"}
PAGE_CHUNK_PARSERS = {"pymupdf", "pypdf", "pdfplumber", "pdfminer"}
PDF_SUBSET_TARGET_PARSERS = {"markitdown", "pymupdf4llm", "docling", "pspdfkit"}
RAG_SNIPPET_CHARS = 500
PARSER_HEALTH_CACHE_TTL_SECONDS = 24 * 60 * 60
PARSER_HEALTH_CACHE_STATUSES = {"failed", "skipped", "timeout"}
CUSTOMER_BEST_SCHEMA_VERSION = "1.0"
PARSER_RESULT_CACHE_SCHEMA_VERSION = "1.0"
BUSINESS_STRUCTURED_PROFILES = {"contract", "bank_statement", "quotation", "purchase_order", "report", "annual_report"}
CUSTOMER_TEXT_TYPES = {"customer_invoice_delivery", "customer_structured_delivery", "customer_best_text"}
GOLDEN_REPORT_NAMES = {"golden_report.json", "vote_report.json", "probe_report.json", "manifest.json", "metadata.json", "index.json"}
GOLDEN_DEFAULT_PARSERS = "pymupdf,pdfplumber,liteparse"

# ---------------------------------------------------------------------------
# 数据结构
# ---------------------------------------------------------------------------

@dataclass
class ParseResult:
    parser: str
    status: str
    seconds: float
    output_file: str | None = None
    error: str | None = None
    chars: int = 0
    non_space_chars: int = 0
    lines: int = 0
    paragraphs: int = 0
    headings: int = 0
    chinese_chars: int = 0
    cid_markers: int = 0
    control_chars: int = 0
    quality_score: float = 0.0
    quality_label: str = "unknown"
    preview: str = ""
    cache_status: str | None = None


# ---------------------------------------------------------------------------
# 工具函数
# ---------------------------------------------------------------------------

def module_exists(module_name: str) -> bool:
    try:
        return importlib.util.find_spec(module_name) is not None
    except ModuleNotFoundError:
        return False


def command_exists(*command_names: str) -> bool:
    """Return whether at least one command name is available on PATH."""
    return any(shutil.which(command_name) is not None for command_name in command_names)


def find_command(*command_names: str) -> str | None:
    """Return the first available command name/path from PATH."""
    for command_name in command_names:
        found = shutil.which(command_name)
        if found:
            return found
    return None


def normalize_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = _RE_WHITESPACE.sub(" ", text)
    text = _RE_MULTI_NEWLINE.sub("\n\n", text)
    return text.strip()


def compact_text(text: str) -> str:
    """Return text with all whitespace removed for regexes across layout breaks."""
    return re.sub(r"\s+", "", text or "")


def decimal_from_text(value: str | None) -> float | None:
    """Parse a currency/number field into float without raising."""
    if not value:
        return None
    cleaned = value.replace(",", "").replace("¥", "").replace("￥", "").strip()
    try:
        return float(Decimal(cleaned))
    except (InvalidOperation, ValueError):
        return None


def quality_from_text(
    text: str,
    non_space_chars: int | None = None,
    cid_markers: int | None = None,
    control_chars: int | None = None,
) -> tuple[float, str]:
    """Estimate extraction quality from text-only signals."""
    if non_space_chars is None:
        non_space_chars = len(_RE_INLINE_WHITESPACE.sub("", text))
    if cid_markers is None:
        cid_markers = text.count("(cid:")
    if control_chars is None:
        control_chars = sum(1 for ch in text if ord(ch) < 32 and ch not in "\n\t")

    if non_space_chars == 0:
        return 0.0, "empty"

    chinese_chars = len(_RE_CHINESE.findall(text))
    cid_penalty = min(0.65, cid_markers / max(1, non_space_chars) * 10)
    control_penalty = min(0.65, control_chars / max(1, non_space_chars) * 8)
    length_penalty = 0.25 if non_space_chars < 80 else 0.1 if non_space_chars < 250 else 0.0
    language_bonus = min(0.2, chinese_chars / max(1, non_space_chars) * 0.5)

    score = max(0.0, min(1.0, 0.85 + language_bonus - cid_penalty - control_penalty - length_penalty))
    if score >= 0.75:
        label = "high"
    elif score >= 0.5:
        label = "medium"
    elif score >= 0.25:
        label = "low"
    else:
        label = "bad"
    return round(score, 3), label


def parse_extensions(ext_str: str | None) -> set[str]:
    """Parse comma-separated extensions into normalized dotted suffixes."""
    raw = ext_str or ".pdf"
    extensions = {
        e.strip().lower() if e.strip().startswith(".") else f".{e.strip().lower()}"
        for e in raw.split(",")
        if e.strip()
    }
    return extensions or {".pdf"}


def quality_gate_status(
    result: ParseResult | None,
    min_quality: float,
    fail_on_bad: bool,
) -> dict[str, object]:
    """Return a reusable quality gate payload for compare and scan commands."""
    threshold = max(0.0, min(1.0, min_quality))
    if result is None:
        passed = False
        reason = "no_successful_output"
    elif result.quality_score < threshold:
        passed = False
        reason = "below_min_quality"
    elif fail_on_bad and result.quality_label in QUALITY_BAD_LABELS:
        passed = False
        reason = "bad_quality_label"
    else:
        passed = True
        reason = "ok"
    return {
        "enabled": bool(threshold > 0 or fail_on_bad),
        "passed": passed,
        "min_quality": threshold,
        "fail_on_bad": fail_on_bad,
        "reason": reason,
        "best_parser": result.parser if result else None,
        "best_quality_score": result.quality_score if result else 0.0,
        "best_quality_label": result.quality_label if result else "none",
    }


def compute_page_range(page_count: int, start_page: int, max_pages: int | None) -> tuple[int, int]:
    """根据总页数、起始页、最大页数计算实际的 (start_index, end_index)。"""
    start_index = min(start_page - 1, page_count - 1)
    if max_pages is None:
        end_index = page_count - 1
    else:
        end_index = min(start_index + max_pages - 1, page_count - 1)
    return start_index, end_index


def get_pdf_page_count(pdf_path: Path) -> int:
    """获取 PDF 总页数，优先用 fitz，后备用 pypdf。"""
    if module_exists("fitz"):
        import fitz  # type: ignore
        doc = fitz.open(pdf_path)
        count = doc.page_count
        doc.close()
        return count
    if module_exists("pypdf"):
        from pypdf import PdfReader  # type: ignore
        return len(PdfReader(str(pdf_path)).pages)
    raise ImportError("需要 fitz 或 pypdf 来获取 PDF 页数")


def selected_pdf_page_indices(pdf_path: Path, start_page: int, max_pages: int | None) -> list[int]:
    """Return selected PDF page indices using 1-based CLI page arguments."""
    page_count = get_pdf_page_count(pdf_path)
    start_index, end_index = compute_page_range(page_count, start_page, max_pages)
    return list(range(start_index, end_index + 1))


def text_stats(parser: str, status: str, seconds: float, text: str, output_file: Path | None) -> ParseResult:
    """从已标准化的文本生成统计数据。"""
    lines = text.splitlines()
    paragraphs = [p for p in _RE_PARAGRAPH_SPLIT.split(text) if p.strip()]
    non_space_chars = len(_RE_INLINE_WHITESPACE.sub("", text))
    cid_markers = text.count("(cid:")
    control_chars = sum(1 for ch in text if ord(ch) < 32 and ch not in "\n\t")
    quality_score, quality_label = quality_from_text(text, non_space_chars, cid_markers, control_chars)
    return ParseResult(
        parser=parser,
        status=status,
        seconds=round(seconds, 3),
        output_file=str(output_file) if output_file else None,
        chars=len(text),
        non_space_chars=non_space_chars,
        lines=len(lines),
        paragraphs=len(paragraphs),
        headings=sum(1 for line in lines if line.lstrip().startswith("#")),
        chinese_chars=len(_RE_CHINESE.findall(text)),
        cid_markers=cid_markers,
        control_chars=control_chars,
        quality_score=quality_score,
        quality_label=quality_label,
        preview=text[:600],
    )


def result_from_text(parser: str, text: str, seconds: float = 0.0) -> ParseResult:
    """Build a ParseResult from already-normalized text."""
    return text_stats(parser, "ok", seconds, normalize_text(text), None)


TEXT_OUTPUT_FORMATS = {"md", "txt", "json", "all"}


def output_suffixes(output_format: str) -> list[str]:
    """Return concrete output suffixes for a requested text export format."""
    if output_format == "all":
        return ["md", "txt", "json"]
    return [output_format]


def default_output_path(source_file: Path, output_format: str) -> Path:
    """Build the default convert output path for md/txt/json/all."""
    if output_format == "all":
        return source_file.parent / f"{source_file.stem}_converted"
    return source_file.parent / f"{source_file.stem}.{output_format}"


def write_extracted_outputs(
    out_dir: Path,
    parser_name: str,
    text: str,
    output_format: str = "md",
) -> tuple[Path, str]:
    """Normalize text, write requested formats, and return the primary output path."""
    normalized = normalize_text(text)
    out_dir.mkdir(parents=True, exist_ok=True)
    paths: dict[str, Path] = {}

    if output_format in ("md", "all"):
        md_path = out_dir / f"{parser_name}.md"
        md_path.write_text(normalized + "\n", encoding="utf-8")
        paths["md"] = md_path

    if output_format in ("txt", "all"):
        txt_path = out_dir / f"{parser_name}.txt"
        txt_path.write_text(normalized + "\n", encoding="utf-8")
        paths["txt"] = txt_path

    if output_format in ("json", "all"):
        json_path = out_dir / f"{parser_name}.json"
        payload = {
            "version": __version__,
            "parser": parser_name,
            "format": "text",
            "chars": len(normalized),
            "text": normalized,
        }
        json_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
        paths["json"] = json_path

    primary_suffix = "md" if output_format == "all" else output_format
    return paths[primary_suffix], normalized


def write_single_conversion_output(
    out_path: Path,
    parser_name: str,
    source_file: Path,
    text: str,
    output_format: str,
) -> list[Path]:
    """Write one conversion result as md, txt, json, or all formats."""
    normalized = normalize_text(text)
    written: list[Path] = []
    suffixes = output_suffixes(output_format)

    if output_format == "all":
        out_path.mkdir(parents=True, exist_ok=True)
        targets = {suffix: out_path / f"{source_file.stem}.{suffix}" for suffix in suffixes}
    else:
        out_path.parent.mkdir(parents=True, exist_ok=True)
        targets = {output_format: out_path}

    for suffix, target in targets.items():
        if suffix in ("md", "txt"):
            target.write_text(normalized + "\n", encoding="utf-8")
        elif suffix == "json":
            payload = {
                "version": __version__,
                "source_file": str(source_file),
                "parser": parser_name,
                "format": "text",
                "chars": len(normalized),
                "text": normalized,
            }
            target.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
        written.append(target)

    return written


def write_text(out_dir: Path, parser_name: str, text: str) -> tuple[Path, str]:
    """Backward-compatible Markdown writer."""
    return write_extracted_outputs(out_dir, parser_name, text, "md")


def default_parser_health_cache_path(out_dir: Path) -> Path:
    """Default parser health cache path colocated with the current run outputs."""
    return out_dir / ".parser_health_cache.json"


def load_parser_health_cache(path: Path | None) -> dict[str, object]:
    """Load parser health cache if present."""
    if path is None or not path.exists():
        return {"version": __version__, "entries": {}}
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return {"version": __version__, "entries": {}}
    if not isinstance(data, dict):
        return {"version": __version__, "entries": {}}
    if not isinstance(data.get("entries"), dict):
        data["entries"] = {}
    return data


def write_parser_health_cache(path: Path | None, cache: dict[str, object]) -> None:
    """Persist parser health cache without failing the parse workflow."""
    if path is None:
        return
    try:
        path.parent.mkdir(parents=True, exist_ok=True)
        cache["version"] = __version__
        cache["updated_at"] = time.strftime("%Y-%m-%d %H:%M:%S")
        path.write_text(json.dumps(cache, ensure_ascii=False, indent=2), encoding="utf-8")
    except Exception as exc:
        logger.debug("failed to write parser health cache %s: %r", path, exc)


def parser_health_cache_key(file_path: Path, fmt: str, parser_name: str, start_page: int, max_pages: int | None) -> str:
    """Build a conservative cache key for parser runtime health."""
    return "|".join([fmt, parser_name, str(file_path.resolve()), str(start_page), str(max_pages)])


def should_skip_parser_from_health_cache(
    cache: dict[str, object],
    key: str,
    ttl_seconds: int = PARSER_HEALTH_CACHE_TTL_SECONDS,
) -> tuple[bool, str | None]:
    """Return whether a recent unhealthy parser entry should be skipped."""
    entries = cache.get("entries") if isinstance(cache.get("entries"), dict) else {}
    entry = entries.get(key) if isinstance(entries, dict) else None
    if not isinstance(entry, dict):
        return False, None
    status = str(entry.get("status") or "")
    timestamp = float(entry.get("timestamp") or 0)
    if status not in PARSER_HEALTH_CACHE_STATUSES:
        return False, None
    if ttl_seconds > 0 and time.time() - timestamp > ttl_seconds:
        return False, None
    reason = str(entry.get("error") or entry.get("reason") or status)
    return True, reason


def update_parser_health_cache(
    cache: dict[str, object],
    key: str,
    parser_name: str,
    result: ParseResult,
    timeout_seconds: float | None = None,
) -> None:
    """Record parser runtime health for later runs."""
    entries = cache.setdefault("entries", {})
    if not isinstance(entries, dict):
        entries = {}
        cache["entries"] = entries
    entries[key] = {
        "parser": parser_name,
        "status": result.status,
        "quality_label": result.quality_label,
        "quality_score": result.quality_score,
        "seconds": result.seconds,
        "error": result.error,
        "timeout_seconds": timeout_seconds,
        "timestamp": time.time(),
        "updated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
    }


def default_parser_result_cache_dir(out_dir: Path) -> Path:
    """Default content cache directory colocated with the current run outputs."""
    return out_dir / ".parser_result_cache"


def file_cache_fingerprint(file_path: Path) -> dict[str, object]:
    """Build a cheap invalidation fingerprint for parser result caching."""
    stat = file_path.stat()
    return {
        "path": str(file_path.resolve()),
        "size_bytes": stat.st_size,
        "modified_ns": getattr(stat, "st_mtime_ns", int(stat.st_mtime * 1_000_000_000)),
    }


def parser_result_cache_key(
    file_path: Path,
    fmt: str,
    parser_name: str,
    start_page: int,
    max_pages: int | None,
    output_format: str,
) -> str:
    """Build a stable key for same-file same-parser extraction output."""
    payload = {
        "schema": PARSER_RESULT_CACHE_SCHEMA_VERSION,
        "version": __version__,
        "file": file_cache_fingerprint(file_path),
        "format": fmt,
        "parser": parser_name,
        "start_page": start_page,
        "max_pages": max_pages,
        "output_format": output_format,
    }
    raw = json.dumps(payload, ensure_ascii=False, sort_keys=True)
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()


def parser_result_cache_path(cache_dir: Path, key: str) -> Path:
    """Return the JSON file path for a parser result cache key."""
    return cache_dir / f"{key}.json"


def load_parser_result_cache(cache_dir: Path | None, key: str) -> dict[str, object] | None:
    """Load cached parser text if available and valid."""
    if cache_dir is None:
        return None
    path = parser_result_cache_path(cache_dir, key)
    if not path.exists():
        return None
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return None
    if not isinstance(data, dict):
        return None
    if data.get("schema_version") != PARSER_RESULT_CACHE_SCHEMA_VERSION or data.get("key") != key:
        return None
    if not isinstance(data.get("normalized"), str):
        return None
    return data


def write_parser_result_cache(
    cache_dir: Path | None,
    key: str,
    parser_name: str,
    result: ParseResult,
    normalized: str,
) -> None:
    """Persist parser text cache without failing the main workflow."""
    if cache_dir is None or result.status != "ok" or not normalized:
        return
    try:
        cache_dir.mkdir(parents=True, exist_ok=True)
        payload = {
            "schema_version": PARSER_RESULT_CACHE_SCHEMA_VERSION,
            "key": key,
            "version": __version__,
            "parser": parser_name,
            "status": result.status,
            "quality_score": result.quality_score,
            "quality_label": result.quality_label,
            "chars": result.chars,
            "non_space_chars": result.non_space_chars,
            "normalized": normalized,
            "created_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        }
        parser_result_cache_path(cache_dir, key).write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")
    except Exception as exc:
        logger.debug("failed to write parser result cache %s: %r", cache_dir, exc)


def extract_one_subprocess(
    file_path: Path,
    fmt: str,
    parser_name: str,
    out_dir: Path,
    output_format: str,
    start_page: int,
    max_pages: int | None,
    timeout_seconds: float,
) -> tuple[ParseResult, str]:
    """Run one parser in a child process so slow parsers can be killed on timeout."""
    temp_dir = Path(tempfile.mkdtemp(prefix="parse_one_"))
    payload_path = temp_dir / "payload.json"
    result_path = temp_dir / "result.json"
    payload = {
        "file": str(file_path),
        "fmt": fmt,
        "parser": parser_name,
        "out_dir": str(out_dir),
        "output_format": output_format,
        "start_page": start_page,
        "max_pages": max_pages,
        "result_path": str(result_path),
    }
    payload_path.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")
    start = time.perf_counter()
    try:
        completed = subprocess.run(
            [sys.executable, str(Path(__file__).resolve()), "__extract-one", str(payload_path)],
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=timeout_seconds,
        )
    except subprocess.TimeoutExpired:
        seconds = round(time.perf_counter() - start, 3)
        shutil.rmtree(temp_dir, ignore_errors=True)
        return ParseResult(
            parser=parser_name,
            status="timeout",
            seconds=seconds,
            error=f"timeout_after_{timeout_seconds}s",
        ), ""

    try:
        if completed.returncode != 0:
            seconds = round(time.perf_counter() - start, 3)
            err = (completed.stderr or completed.stdout or "").strip()
            return ParseResult(parser=parser_name, status="failed", seconds=seconds, error=err[:2000] or f"exit_{completed.returncode}"), ""
        data = json.loads(result_path.read_text(encoding="utf-8"))
        result_data = data.get("result") if isinstance(data.get("result"), dict) else {}
        result = ParseResult(**result_data)
        return result, str(data.get("normalized") or "")
    except Exception as exc:
        seconds = round(time.perf_counter() - start, 3)
        return ParseResult(parser=parser_name, status="failed", seconds=seconds, error=f"invalid_child_result: {exc!r}"), ""
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


def cmd_extract_one_internal(payload_path: Path) -> int:
    """Internal helper for timed parser subprocesses."""
    payload = json.loads(payload_path.read_text(encoding="utf-8"))
    file_path = Path(str(payload["file"]))
    fmt = str(payload["fmt"])
    parser_name = str(payload["parser"])
    out_dir = Path(str(payload["out_dir"]))
    output_format = str(payload["output_format"])
    start_page = int(payload["start_page"])
    max_pages = payload.get("max_pages")
    max_pages = int(max_pages) if max_pages is not None else None
    result_path = Path(str(payload["result_path"]))
    extract_func = get_extractor(fmt, parser_name)
    if extract_func is None:
        raise RuntimeError(f"未注册的解析器：{parser_name}")
    start = time.perf_counter()
    raw_text = extract_func(file_path, start_page, max_pages)
    out_file, normalized = write_extracted_outputs(out_dir, parser_name, raw_text, output_format)
    result = text_stats(parser_name, "ok", time.perf_counter() - start, normalized, out_file)
    write_json_file(result_path, {"result": asdict(result), "normalized": normalized})
    return 0


def parse_page_spec(spec: str, total_pages: int) -> list[int]:
    """解析页码规格，如 '1-5' 或 '1,3,5' 或 '1-3,7,10-12'。返回 0-based 索引列表。"""
    pages: set[int] = set()
    for part in spec.split(","):
        part = part.strip()
        if "-" in part:
            start_s, end_s = part.split("-", 1)
            start = max(1, int(start_s))
            end = min(total_pages, int(end_s))
            pages.update(range(start - 1, end))
        else:
            p = int(part)
            if 1 <= p <= total_pages:
                pages.add(p - 1)
    return sorted(pages)


def resolve_page_indices(spec: str | None, total_pages: int, default_all: bool = False) -> list[int]:
    """Resolve CLI page specs; supports None, 'all', '*', and parse_page_spec syntax."""
    if total_pages <= 0:
        return []
    if spec is None or not spec.strip():
        return list(range(total_pages)) if default_all else [0]
    normalized = spec.strip().lower()
    if normalized in {"all", "*"}:
        return list(range(total_pages))
    pages = parse_page_spec(spec, total_pages)
    if not pages:
        raise ValueError(f"页码范围无效：{spec}")
    return pages


# ---------------------------------------------------------------------------
# 格式检测
# ---------------------------------------------------------------------------

EXTENSION_FORMAT_MAP: dict[str, str] = {
    ".pdf": "pdf",
    ".docx": "word",
    ".doc": "word",
    ".pptx": "ppt",
    ".ppt": "ppt",
    ".xlsx": "excel",
    ".xls": "excel",
    ".html": "html",
    ".htm": "html",
    # 文本格式
    ".txt": "text",
    ".csv": "text",
    ".json": "text",
    ".xml": "text",
    ".md": "text",
    ".markdown": "text",
    # 图片格式
    ".jpg": "image",
    ".jpeg": "image",
    ".png": "image",
    ".gif": "image",
    ".bmp": "image",
    ".tiff": "image",
    ".tif": "image",
    ".webp": "image",
    # 音频格式
    ".mp3": "audio",
    ".wav": "audio",
    # 电子书
    ".epub": "epub",
    # 压缩包
    ".zip": "archive",
}

FORMAT_PARSERS: dict[str, list[str]] = {
    "pdf":     [
        "markitdown",
        "pymupdf4llm",
        "docling",
        "pspdfkit",
        "pymupdf",
        "pypdf",
        "pdfplumber",
        "pdfminer",
        "liteparse",
        "opendataloader",
        "ocr-tesseract",
    ],
    "word":    ["markitdown", "python-docx"],
    "ppt":     ["markitdown", "python-pptx"],
    "excel":   ["markitdown", "openpyxl"],
    "html":    ["markitdown", "beautifulsoup4"],
    "text":    ["markitdown"],
    "image":   ["markitdown"],
    "audio":   ["markitdown"],
    "epub":    ["markitdown"],
    "archive": ["markitdown"],
}


PARSER_MODULES: dict[str, list[str]] = {
    "markitdown": ["markitdown"],
    "pymupdf4llm": ["pymupdf4llm"],
    "docling": ["docling.document_converter"],
    "pymupdf": ["fitz"],
    "pypdf": ["pypdf", "PyPDF2"],
    "pdfplumber": ["pdfplumber"],
    "pdfminer": ["pdfminer.high_level"],
    "liteparse": ["liteparse"],
    "opendataloader": ["opendataloader_pdf"],
    "ocr-tesseract": ["fitz", "pytesseract", "PIL"],
    "python-docx": ["docx"],
    "python-pptx": ["pptx"],
    "openpyxl": ["openpyxl"],
    "beautifulsoup4": ["bs4"],
}


PARSER_COMMANDS: dict[str, list[str]] = {
    "pspdfkit": ["pdf-to-markdown", "pspdfkit-pdf-to-markdown"],
}


PARSER_PACKAGES: dict[str, str] = {
    "markitdown": "markitdown",
    "pymupdf4llm": "pymupdf4llm",
    "docling": "docling",
    "pspdfkit": "@pspdfkit/pdf-to-markdown (pdf-to-markdown CLI)",
    "pymupdf": "pymupdf",
    "pypdf": "pypdf",
    "pdfplumber": "pdfplumber",
    "pdfminer": "pdfminer.six",
    "liteparse": "liteparse",
    "opendataloader": "opendataloader-pdf",
    "ocr-tesseract": "pymupdf pytesseract pillow + system Tesseract OCR",
    "python-docx": "python-docx",
    "python-pptx": "python-pptx",
    "openpyxl": "openpyxl",
    "beautifulsoup4": "beautifulsoup4",
}


OCR_DEPENDENCIES: list[dict[str, object]] = [
    {
        "name": "pytesseract",
        "available": module_exists("pytesseract"),
        "kind": "python",
        "package": "pytesseract",
    },
    {
        "name": "Pillow",
        "available": module_exists("PIL"),
        "kind": "python",
        "package": "pillow",
    },
    {
        "name": "tesseract",
        "available": shutil.which("tesseract") is not None,
        "kind": "system",
        "package": "tesseract-ocr",
    },
]

OPENLOADER_DEPENDENCIES: list[dict[str, object]] = [
    {
        "name": "opendataloader_pdf",
        "available": module_exists("opendataloader_pdf"),
        "kind": "python",
        "package": "opendataloader-pdf",
    },
    {
        "name": "java",
        "available": shutil.which("java") is not None,
        "kind": "system",
        "package": "JDK/JRE (java command)",
    },
]


def detect_format(file_path: Path) -> str | None:
    """根据文件扩展名检测格式。"""
    return EXTENSION_FORMAT_MAP.get(file_path.suffix.lower())


def get_parsers_for_format(fmt: str) -> list[str]:
    """获取指定格式的所有可用解析器。"""
    return FORMAT_PARSERS.get(fmt, [])


def parser_available(parser_name: str) -> bool:
    """Return whether a parser's Python module or external command is available."""
    commands = PARSER_COMMANDS.get(parser_name)
    if commands:
        return command_exists(*commands)
    if parser_name == "ocr-tesseract":
        return module_exists("fitz") and module_exists("pytesseract") and module_exists("PIL") and shutil.which("tesseract") is not None
    modules = PARSER_MODULES.get(parser_name, [parser_name])
    return any(module_exists(module_name) for module_name in modules)


def parser_dependency_rows(fmt_filter: str | None = None) -> list[dict[str, object]]:
    """Build parser dependency rows for doctor output."""
    formats = [fmt_filter] if fmt_filter else sorted(FORMAT_PARSERS)
    rows: list[dict[str, object]] = []
    for fmt in formats:
        for parser_name in FORMAT_PARSERS.get(fmt, []):
            commands = PARSER_COMMANDS.get(parser_name, [])
            modules = [] if commands else PARSER_MODULES.get(parser_name, [parser_name])
            package = PARSER_PACKAGES.get(parser_name, parser_name)
            available = parser_available(parser_name)
            kind = "command" if commands else "ocr" if parser_name == "ocr-tesseract" else "python"
            rows.append(
                {
                    "format": fmt,
                    "parser": parser_name,
                    "available": available,
                    "status": "ok" if available else "missing",
                    "kind": kind,
                    "modules": modules,
                    "commands": commands + (["tesseract"] if parser_name == "ocr-tesseract" else []),
                    "package": package,
                }
            )
    return rows


def extract_text_with_parser(
    file_path: Path,
    parser_name: str,
    start_page: int = 1,
    max_pages: int | None = None,
) -> str:
    """Extract text with one parser and handle whole-PDF parser page slicing."""
    fmt = detect_format(file_path)
    if fmt is None:
        raise ValueError(f"不支持的文件格式：{file_path.suffix}")
    parser_name = resolve_parser_name(parser_name)
    if parser_name not in get_parsers_for_format(fmt):
        raise ValueError(f"{fmt} 格式不支持解析器 {parser_name}")
    extract_func = get_extractor(fmt, parser_name)
    if extract_func is None:
        raise ValueError(f"未找到解析器 {parser_name}")

    temp_dir_obj: tempfile.TemporaryDirectory[str] | None = None
    try:
        target_path = file_path
        if fmt == "pdf" and parser_name in PDF_SUBSET_TARGET_PARSERS:
            temp_dir_obj = tempfile.TemporaryDirectory(prefix="parse_extract_")
            target_path, skip_reason = prepare_pdf_subset_target(
                file_path, start_page, max_pages, Path(temp_dir_obj.name)
            )
            if skip_reason:
                raise RuntimeError(skip_reason)
        return extract_func(target_path, start_page, max_pages)
    finally:
        if temp_dir_obj:
            temp_dir_obj.cleanup()


def extract_best_text(
    file_path: Path,
    parser_name: str = "auto",
    start_page: int = 1,
    max_pages: int | None = None,
) -> tuple[str, ParseResult]:
    """Extract text with an explicit parser or choose the best available parser."""
    fmt = detect_format(file_path)
    if fmt is None:
        raise ValueError(f"不支持的文件格式：{file_path.suffix}")

    if parser_name != "auto":
        start = time.perf_counter()
        raw_text = extract_text_with_parser(file_path, parser_name, start_page, max_pages)
        normalized = normalize_text(raw_text)
        return normalized, text_stats(resolve_parser_name(parser_name), "ok", time.perf_counter() - start, normalized, None)

    results: list[tuple[str, ParseResult]] = []
    errors: list[str] = []
    for candidate in get_parsers_for_format(fmt):
        start = time.perf_counter()
        try:
            raw_text = extract_text_with_parser(file_path, candidate, start_page, max_pages)
            normalized = normalize_text(raw_text)
            results.append((normalized, text_stats(candidate, "ok", time.perf_counter() - start, normalized, None)))
        except Exception as exc:
            errors.append(f"{candidate}: {exc}")

    recommended = choose_recommended_result([result for _, result in results])
    if recommended is None:
        raise RuntimeError("没有成功解析结果：" + "; ".join(errors))
    for text, result in results:
        if result.parser == recommended.parser:
            return text, result
    raise RuntimeError("无法匹配推荐解析器结果")


def file_metadata(file_path: Path) -> dict[str, object]:
    """Collect file-level metadata; add PDF details when possible."""
    fmt = detect_format(file_path)
    stat = file_path.stat()
    payload: dict[str, object] = {
        "version": __version__,
        "source_file": str(file_path),
        "file_name": file_path.name,
        "extension": file_path.suffix.lower(),
        "format": fmt,
        "size_bytes": stat.st_size,
        "modified_at": time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(stat.st_mtime)),
    }

    if fmt == "pdf":
        payload["pdf"] = pdf_metadata(file_path)

    return payload


def pdf_metadata(pdf_path: Path) -> dict[str, object]:
    """Collect PDF metadata using PyMuPDF when available, falling back to pypdf."""
    if module_exists("fitz"):
        import fitz  # type: ignore

        doc = fitz.open(pdf_path)
        try:
            sample_pages = min(5, doc.page_count)
            text_chars_by_page: list[int] = []
            for index in range(sample_pages):
                text_chars_by_page.append(len((doc.load_page(index).get_text("text") or "").strip()))
            toc = doc.get_toc(simple=True)
            return {
                "page_count": doc.page_count,
                "metadata": dict(doc.metadata or {}),
                "toc_count": len(toc),
                "toc_preview": [{"level": item[0], "title": item[1], "page": item[2]} for item in toc[:30]],
                "sample_text_chars_by_page": text_chars_by_page,
                "has_sample_text_layer": any(count > 0 for count in text_chars_by_page),
                "parser": "pymupdf",
            }
        finally:
            doc.close()

    if module_exists("pypdf"):
        from pypdf import PdfReader  # type: ignore

        reader = PdfReader(str(pdf_path))
        metadata = {}
        if reader.metadata:
            metadata = {str(k): str(v) for k, v in dict(reader.metadata).items()}
        return {
            "page_count": len(reader.pages),
            "metadata": metadata,
            "toc_count": 0,
            "toc_preview": [],
            "has_sample_text_layer": None,
            "parser": "pypdf",
        }

    raise ImportError("需要 PyMuPDF/fitz 或 pypdf 才能读取 PDF 元数据")


def metadata_to_markdown(metadata: dict[str, object]) -> str:
    """Render metadata as Markdown."""
    lines = [
        "# Document Metadata",
        "",
        f"- File: `{metadata.get('source_file')}`",
        f"- Format: `{metadata.get('format')}`",
        f"- Size bytes: `{metadata.get('size_bytes')}`",
        f"- Modified: `{metadata.get('modified_at')}`",
    ]

    pdf = metadata.get("pdf")
    if isinstance(pdf, dict):
        lines.extend(
            [
                "",
                "## PDF",
                "",
                f"- Page count: `{pdf.get('page_count')}`",
                f"- Metadata parser: `{pdf.get('parser')}`",
                f"- TOC entries: `{pdf.get('toc_count')}`",
                f"- Has sample text layer: `{pdf.get('has_sample_text_layer')}`",
                f"- Sample text chars by page: `{pdf.get('sample_text_chars_by_page')}`",
            ]
        )
        doc_meta = pdf.get("metadata")
        if isinstance(doc_meta, dict) and doc_meta:
            lines.extend(["", "## Embedded Metadata", ""])
            for key, value in doc_meta.items():
                if value:
                    lines.append(f"- `{key}`: {value}")
        toc_preview = pdf.get("toc_preview")
        if isinstance(toc_preview, list) and toc_preview:
            lines.extend(["", "## TOC Preview", ""])
            for item in toc_preview:
                if isinstance(item, dict):
                    lines.append(f"- L{item.get('level')} p{item.get('page')}: {item.get('title')}")

    return "\n".join(lines) + "\n"


def split_text_into_chunks(text: str, chunk_size: int, overlap: int) -> list[dict[str, object]]:
    """Split text into overlapping character chunks."""
    normalized = normalize_text(text)
    if chunk_size <= 0:
        raise ValueError("chunk_size must be positive")
    if overlap < 0:
        raise ValueError("overlap must be >= 0")
    if overlap >= chunk_size:
        raise ValueError("overlap must be smaller than chunk_size")

    chunks: list[dict[str, object]] = []
    start = 0
    chunk_id = 1
    while start < len(normalized):
        end = min(start + chunk_size, len(normalized))
        chunk_text = normalized[start:end]
        chunks.append(
            {
                "chunk_id": chunk_id,
                "start": start,
                "end": end,
                "chars": len(chunk_text),
                "text": chunk_text,
            }
        )
        if end == len(normalized):
            break
        start = end - overlap
        chunk_id += 1
    return chunks


def split_pages_into_chunks(pages: list[dict[str, object]]) -> list[dict[str, object]]:
    """Build one chunk per source page."""
    chunks: list[dict[str, object]] = []
    for chunk_id, page in enumerate(pages, start=1):
        text = normalize_text(str(page.get("text", "")))
        quality_score, quality_label = quality_from_text(text)
        chunks.append(
            {
                "chunk_id": chunk_id,
                "chunk_by": "page",
                "page": page.get("page"),
                "start": 0,
                "end": len(text),
                "chars": len(text),
                "quality_score": quality_score,
                "quality_label": quality_label,
                "text": text,
            }
        )
    return chunks


def write_chunks(chunks: list[dict[str, object]], out_path: Path, output_format: str, metadata: dict[str, object]) -> list[Path]:
    """Write chunks as jsonl, json, md, txt, or all."""
    if output_format == "all":
        out_path.mkdir(parents=True, exist_ok=True)
        targets = {
            "jsonl": out_path / "chunks.jsonl",
            "json": out_path / "chunks.json",
            "md": out_path / "chunks.md",
            "txt": out_path / "chunks.txt",
        }
    else:
        out_path.parent.mkdir(parents=True, exist_ok=True)
        targets = {output_format: out_path}

    written: list[Path] = []
    for suffix, target in targets.items():
        if suffix == "jsonl":
            lines = [json.dumps({**metadata, **chunk}, ensure_ascii=False) for chunk in chunks]
            target.write_text("\n".join(lines) + ("\n" if lines else ""), encoding="utf-8")
        elif suffix == "json":
            payload = {"metadata": metadata, "chunk_count": len(chunks), "chunks": chunks}
            target.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
        elif suffix == "md":
            parts = ["# Document Chunks", ""]
            for chunk in chunks:
                parts.extend([f"## Chunk {chunk['chunk_id']}", "", str(chunk["text"]), ""])
            target.write_text("\n".join(parts), encoding="utf-8")
        elif suffix == "txt":
            parts = []
            for chunk in chunks:
                page_suffix = f" page {chunk['page']}" if "page" in chunk else ""
                parts.append(f"--- chunk {chunk['chunk_id']}{page_suffix} [{chunk['start']}:{chunk['end']}] ---")
                parts.append(str(chunk["text"]))
            target.write_text("\n\n".join(parts) + ("\n" if parts else ""), encoding="utf-8")
        written.append(target)
    return written


# ---------------------------------------------------------------------------
# PDF 解析器
# ---------------------------------------------------------------------------

def _extract_pdf_markitdown(pdf_path: Path, start_page: int, max_pages: int | None) -> str:
    if not module_exists("markitdown"):
        raise ImportError("未安装 markitdown")
    from markitdown import MarkItDown  # type: ignore
    result = MarkItDown().convert(str(pdf_path))
    return getattr(result, "text_content", "") or ""


def _extract_pdf_pymupdf4llm(pdf_path: Path, start_page: int, max_pages: int | None) -> str:
    """使用 pymupdf4llm 生成面向 LLM/RAG 的 Markdown。"""
    if not module_exists("pymupdf4llm"):
        raise ImportError("未安装 pymupdf4llm，运行: pip install pymupdf4llm")
    import pymupdf4llm  # type: ignore

    return pymupdf4llm.to_markdown(str(pdf_path)) or ""


def _extract_pdf_docling(pdf_path: Path, start_page: int, max_pages: int | None) -> str:
    """使用 Docling 解析 PDF 并导出 Markdown。"""
    if not module_exists("docling.document_converter"):
        raise ImportError("未安装 docling，运行: pip install docling")
    from docling.document_converter import DocumentConverter  # type: ignore

    result = DocumentConverter().convert(str(pdf_path))
    document = getattr(result, "document", None)
    if document is None or not hasattr(document, "export_to_markdown"):
        raise RuntimeError("docling 未返回可导出 Markdown 的 document")
    return document.export_to_markdown() or ""


def _extract_pdf_pspdfkit(pdf_path: Path, start_page: int, max_pages: int | None) -> str:
    """使用 PSPDFKit/Nutrient pdf-to-markdown CLI 解析 PDF。"""
    command = find_command("pdf-to-markdown", "pspdfkit-pdf-to-markdown")
    if command is None:
        raise ImportError(
            "未找到 pdf-to-markdown CLI。安装 @pspdfkit/pdf-to-markdown 后确保 pdf-to-markdown 在 PATH 中；"
            "该工具当前主要支持 macOS/Linux，Windows 环境可能不可用。"
        )

    with tempfile.TemporaryDirectory(prefix="pspdfkit_pdfmd_") as tmpdir:
        out_dir = Path(tmpdir)
        out_file = out_dir / f"{pdf_path.stem}.md"
        candidates = [
            [command, str(pdf_path)],
            [command, str(pdf_path), str(out_file)],
        ]
        errors: list[str] = []
        for cmd in candidates:
            try:
                completed = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    encoding="utf-8",
                    errors="replace",
                    check=False,
                )
            except OSError as exc:
                errors.append(repr(exc))
                continue

            stdout = completed.stdout.strip()
            if completed.returncode == 0 and stdout:
                return stdout

            if completed.returncode == 0 and out_file.exists():
                return out_file.read_text(encoding="utf-8", errors="replace")

            md_files = sorted(out_dir.glob("*.md"))
            if completed.returncode == 0 and md_files:
                return "\n\n".join(path.read_text(encoding="utf-8", errors="replace") for path in md_files)

            stderr = completed.stderr.strip()
            errors.append(stderr or stdout or f"exit code {completed.returncode}")

        raise RuntimeError("pdf-to-markdown CLI 解析失败：" + " | ".join(errors))


def prepare_pdf_subset_target(
    pdf_path: Path,
    start_page: int,
    max_pages: int | None,
    tmp_dir: Path,
) -> tuple[Path, str | None]:
    """Return a page-limited PDF for whole-document parsers when possible."""
    if max_pages is None:
        return pdf_path, None

    subset = create_subset_pdf(pdf_path, start_page, max_pages, tmp_dir)
    if subset:
        return subset, None

    return pdf_path, "设置了 --max-pages，但当前环境无法生成临时子 PDF；已跳过需要整份 PDF 输入的解析器"


def _extract_pdf_pymupdf(pdf_path: Path, start_page: int, max_pages: int | None) -> str:
    if not module_exists("fitz"):
        raise ImportError("未安装 PyMuPDF/fitz")
    import fitz  # type: ignore

    parts: list[str] = []
    doc = fitz.open(pdf_path)
    try:
        start_index, end_index = compute_page_range(doc.page_count, start_page, max_pages)
        for index in range(start_index, end_index + 1):
            page = doc.load_page(index)
            parts.append(f"\n\n<!-- page {index + 1} -->\n\n")
            parts.append(page.get_text("text"))
    finally:
        doc.close()
    return "".join(parts)


def _extract_pdf_pages_pymupdf(pdf_path: Path, start_page: int, max_pages: int | None) -> list[dict[str, object]]:
    if not module_exists("fitz"):
        raise ImportError("未安装 PyMuPDF/fitz")
    import fitz  # type: ignore

    pages: list[dict[str, object]] = []
    doc = fitz.open(pdf_path)
    try:
        start_index, end_index = compute_page_range(doc.page_count, start_page, max_pages)
        for index in range(start_index, end_index + 1):
            page = doc.load_page(index)
            pages.append({"page": index + 1, "text": page.get_text("text") or ""})
    finally:
        doc.close()
    return pages


def _extract_pdf_pypdf(pdf_path: Path, start_page: int, max_pages: int | None) -> str:
    reader_cls = None
    if module_exists("pypdf"):
        from pypdf import PdfReader  # type: ignore
        reader_cls = PdfReader
    elif module_exists("PyPDF2"):
        from PyPDF2 import PdfReader  # type: ignore
        reader_cls = PdfReader
    else:
        raise ImportError("未安装 pypdf 或 PyPDF2")

    reader = reader_cls(str(pdf_path))
    start_index, end_index = compute_page_range(len(reader.pages), start_page, max_pages)
    parts: list[str] = []
    for index in range(start_index, end_index + 1):
        parts.append(f"\n\n<!-- page {index + 1} -->\n\n")
        parts.append(reader.pages[index].extract_text() or "")
    return "".join(parts)


def _extract_pdf_pages_pypdf(pdf_path: Path, start_page: int, max_pages: int | None) -> list[dict[str, object]]:
    reader_cls = None
    if module_exists("pypdf"):
        from pypdf import PdfReader  # type: ignore
        reader_cls = PdfReader
    elif module_exists("PyPDF2"):
        from PyPDF2 import PdfReader  # type: ignore
        reader_cls = PdfReader
    else:
        raise ImportError("未安装 pypdf 或 PyPDF2")

    reader = reader_cls(str(pdf_path))
    start_index, end_index = compute_page_range(len(reader.pages), start_page, max_pages)
    return [
        {"page": index + 1, "text": reader.pages[index].extract_text() or ""}
        for index in range(start_index, end_index + 1)
    ]


def _extract_pdf_pdfplumber(pdf_path: Path, start_page: int, max_pages: int | None) -> str:
    if not module_exists("pdfplumber"):
        raise ImportError("未安装 pdfplumber")
    import pdfplumber  # type: ignore

    parts: list[str] = []
    with pdfplumber.open(str(pdf_path)) as pdf:
        start_index, end_index = compute_page_range(len(pdf.pages), start_page, max_pages)
        for index in range(start_index, end_index + 1):
            parts.append(f"\n\n<!-- page {index + 1} -->\n\n")
            parts.append(pdf.pages[index].extract_text() or "")
    return "".join(parts)


def _extract_pdf_pages_pdfplumber(pdf_path: Path, start_page: int, max_pages: int | None) -> list[dict[str, object]]:
    if not module_exists("pdfplumber"):
        raise ImportError("未安装 pdfplumber")
    import pdfplumber  # type: ignore

    pages: list[dict[str, object]] = []
    with pdfplumber.open(str(pdf_path)) as pdf:
        start_index, end_index = compute_page_range(len(pdf.pages), start_page, max_pages)
        for index in range(start_index, end_index + 1):
            pages.append({"page": index + 1, "text": pdf.pages[index].extract_text() or ""})
    return pages


def _extract_pdf_pdfminer(pdf_path: Path, start_page: int, max_pages: int | None) -> str:
    if not module_exists("pdfminer.high_level"):
        raise ImportError("未安装 pdfminer.six")
    from pdfminer.high_level import extract_text  # type: ignore

    page_count = get_pdf_page_count(pdf_path)
    start_index, end_index = compute_page_range(page_count, start_page, max_pages)
    text = extract_text(str(pdf_path), page_numbers=list(range(start_index, end_index + 1)))
    markers = "".join(
        f"\n\n<!-- page {i + 1} -->\n\n"
        for i in range(start_index, end_index + 1)
    )
    return markers + text


def _extract_pdf_pages_pdfminer(pdf_path: Path, start_page: int, max_pages: int | None) -> list[dict[str, object]]:
    if not module_exists("pdfminer.high_level"):
        raise ImportError("未安装 pdfminer.six")
    from pdfminer.high_level import extract_text  # type: ignore

    page_count = get_pdf_page_count(pdf_path)
    start_index, end_index = compute_page_range(page_count, start_page, max_pages)
    pages: list[dict[str, object]] = []
    for index in range(start_index, end_index + 1):
        text = extract_text(str(pdf_path), page_numbers=[index])
        pages.append({"page": index + 1, "text": text or ""})
    return pages


def _extract_pdf_liteparse(pdf_path: Path, start_page: int, max_pages: int | None) -> str:
    """使用 liteparse（Rust 原生绑定）解析 PDF。"""
    if not module_exists("liteparse"):
        raise ImportError("未安装 liteparse，运行: pip install liteparse")
    from liteparse import LiteParse  # type: ignore

    target_pages = None
    if max_pages is not None:
        end_page = start_page + max_pages - 1
        target_pages = f"{start_page}-{end_page}"

    parser = LiteParse(target_pages=target_pages, quiet=True, ocr_enabled=False)
    result = parser.parse(str(pdf_path))

    parts: list[str] = []
    for page in result.pages:
        parts.append(f"\n\n<!-- page {page.page_num} -->\n\n")
        parts.append(page.text or "")
    return "".join(parts)


def _extract_pdf_opendataloader(pdf_path: Path, start_page: int, max_pages: int | None) -> str:
    """使用 opendataloader-pdf（Java JAR）解析 PDF。"""
    if not module_exists("opendataloader_pdf"):
        raise ImportError("未安装 opendataloader-pdf，运行: pip install opendataloader-pdf（需要系统 Java 环境）")
    from opendataloader_pdf import convert  # type: ignore

    logger.info("opendataloader-pdf: 开始解析 %s (start_page=%s, max_pages=%s)", pdf_path.name, start_page, max_pages)

    with tempfile.TemporaryDirectory(prefix="opendataloader_") as tmpdir:
        # 构造页码范围参数
        pages_spec = None
        if max_pages is not None:
            end_page = start_page + max_pages - 1
            pages_spec = f"{start_page}-{end_page}"
            logger.debug("opendataloader-pdf: 页码范围 %s", pages_spec)
        elif start_page > 1:
            pages_spec = f"{start_page}-"
            logger.debug("opendataloader-pdf: 从第 %s 页到末尾", start_page)

        try:
            convert(
                str(pdf_path),
                output_dir=tmpdir,
                format="text",
                quiet=True,
                **({"pages": pages_spec} if pages_spec else {}),
            )
        except Exception as e:
            logger.error("opendataloader-pdf: 解析失败 - %s", e)
            raise

        # 读取输出文件（文件名包含原 PDF 名）
        output_files = list(Path(tmpdir).glob("*.txt"))
        if not output_files:
            logger.warning("opendataloader-pdf: 未生成输出文件")
            return ""

        output_file = output_files[0]
        logger.info("opendataloader-pdf: 输出文件 %s", output_file.name)
        text = output_file.read_text(encoding="utf-8", errors="replace")
        logger.info("opendataloader-pdf: 提取 %d 字符", len(text))
        return text


# ---------------------------------------------------------------------------
# Word 解析器
# ---------------------------------------------------------------------------

def _extract_word_markitdown(file_path: Path, start_page: int, max_pages: int | None) -> str:
    if not module_exists("markitdown"):
        raise ImportError("未安装 markitdown")
    from markitdown import MarkItDown  # type: ignore
    result = MarkItDown().convert(str(file_path))
    return getattr(result, "text_content", "") or ""


def _extract_word_docx(file_path: Path, start_page: int, max_pages: int | None) -> str:
    if not module_exists("docx"):
        raise ImportError("未安装 python-docx")
    from docx import Document  # type: ignore

    doc = Document(str(file_path))
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            style_name = para.style.name if para.style else ""
            if style_name.startswith("Heading"):
                level = style_name.replace("Heading", "").strip()
                parts.append(f"\n{'#' * int(level or 1)} {text}\n")
            else:
                parts.append(text + "\n")

    for table in doc.tables:
        parts.append("\n")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append("| " + " | ".join(cells) + " |")
        parts.append("\n")

    return "\n".join(parts)


# ---------------------------------------------------------------------------
# PPT 解析器
# ---------------------------------------------------------------------------

def _extract_ppt_markitdown(file_path: Path, start_page: int, max_pages: int | None) -> str:
    if not module_exists("markitdown"):
        raise ImportError("未安装 markitdown")
    from markitdown import MarkItDown  # type: ignore
    result = MarkItDown().convert(str(file_path))
    return getattr(result, "text_content", "") or ""


def _extract_ppt_pptx(file_path: Path, start_page: int, max_pages: int | None) -> str:
    if not module_exists("pptx"):
        raise ImportError("未安装 python-pptx")
    from pptx import Presentation  # type: ignore

    prs = Presentation(str(file_path))
    parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        parts.append(f"\n\n<!-- slide {slide_num} -->\n\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append("| " + " | ".join(cells) + " |")

    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Excel 解析器
# ---------------------------------------------------------------------------

def _extract_excel_markitdown(file_path: Path, start_page: int, max_pages: int | None) -> str:
    if not module_exists("markitdown"):
        raise ImportError("未安装 markitdown")
    from markitdown import MarkItDown  # type: ignore
    result = MarkItDown().convert(str(file_path))
    return getattr(result, "text_content", "") or ""


def _extract_excel_openpyxl(file_path: Path, start_page: int, max_pages: int | None) -> str:
    if not module_exists("openpyxl"):
        raise ImportError("未安装 openpyxl")
    from openpyxl import load_workbook  # type: ignore

    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"\n\n## {sheet_name}\n\n")
        for row in ws.iter_rows(values_only=True):
            cells = [str(cell) if cell is not None else "" for cell in row]
            parts.append("| " + " | ".join(cells) + " |")

    wb.close()
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# HTML 解析器
# ---------------------------------------------------------------------------

def _extract_html_markitdown(file_path: Path, start_page: int, max_pages: int | None) -> str:
    if not module_exists("markitdown"):
        raise ImportError("未安装 markitdown")
    from markitdown import MarkItDown  # type: ignore
    result = MarkItDown().convert(str(file_path))
    return getattr(result, "text_content", "") or ""


def _extract_html_bs4(file_path: Path, start_page: int, max_pages: int | None) -> str:
    if not module_exists("bs4"):
        raise ImportError("未安装 beautifulsoup4")
    from bs4 import BeautifulSoup  # type: ignore

    html = file_path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "html.parser")

    for tag in soup(["script", "style"]):
        tag.decompose()

    parts: list[str] = []
    for element in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "td", "th", "pre", "blockquote"]):
        text = element.get_text(strip=True)
        if not text:
            continue
        tag = element.name
        if tag.startswith("h"):
            level = int(tag[1])
            parts.append(f"\n{'#' * level} {text}\n")
        elif tag == "li":
            parts.append(f"- {text}")
        elif tag in ("td", "th"):
            parts.append(f"| {text}", )
        elif tag == "pre":
            parts.append(f"\n```\n{text}\n```\n")
        elif tag == "blockquote":
            parts.append(f"> {text}")
        else:
            parts.append(text + "\n")

    return "\n".join(parts)


# ---------------------------------------------------------------------------
# 文本格式解析器（TXT, CSV, JSON, XML, MD）
# ---------------------------------------------------------------------------

def _extract_text_markitdown(file_path: Path, start_page: int, max_pages: int | None) -> str:
    """使用 markitdown 解析文本格式（TXT, CSV, JSON, XML, MD 等）。"""
    if not module_exists("markitdown"):
        raise ImportError("未安装 markitdown")
    from markitdown import MarkItDown  # type: ignore

    logger.info("markitdown: 解析文本文件 %s", file_path.name)
    result = MarkItDown().convert(str(file_path))
    text = getattr(result, "text_content", "") or ""
    logger.info("markitdown: 提取 %d 字符", len(text))
    return text


# ---------------------------------------------------------------------------
# 图片解析器（JPG, PNG, GIF, BMP, TIFF, WebP）
# ---------------------------------------------------------------------------

def _extract_image_markitdown(file_path: Path, start_page: int, max_pages: int | None) -> str:
    """使用 markitdown 解析图片（EXIF 元数据 + OCR 文字识别）。"""
    if not module_exists("markitdown"):
        raise ImportError("未安装 markitdown")
    from markitdown import MarkItDown  # type: ignore

    logger.info("markitdown: 解析图片 %s (EXIF + OCR)", file_path.name)
    result = MarkItDown().convert(str(file_path))
    text = getattr(result, "text_content", "") or ""
    logger.info("markitdown: 提取 %d 字符", len(text))
    return text


# ---------------------------------------------------------------------------
# 音频解析器（MP3, WAV）
# ---------------------------------------------------------------------------

def _extract_audio_markitdown(file_path: Path, start_page: int, max_pages: int | None) -> str:
    """使用 markitdown 解析音频（EXIF 元数据 + 语音转录）。"""
    if not module_exists("markitdown"):
        raise ImportError("未安装 markitdown")
    from markitdown import MarkItDown  # type: ignore

    logger.info("markitdown: 解析音频 %s (EXIF + 语音转录)", file_path.name)
    result = MarkItDown().convert(str(file_path))
    text = getattr(result, "text_content", "") or ""
    logger.info("markitdown: 提取 %d 字符", len(text))
    return text


# ---------------------------------------------------------------------------
# 电子书解析器（EPub）
# ---------------------------------------------------------------------------

def _extract_epub_markitdown(file_path: Path, start_page: int, max_pages: int | None) -> str:
    """使用 markitdown 解析 EPub 电子书。"""
    if not module_exists("markitdown"):
        raise ImportError("未安装 markitdown")
    from markitdown import MarkItDown  # type: ignore

    logger.info("markitdown: 解析 EPub %s", file_path.name)
    result = MarkItDown().convert(str(file_path))
    text = getattr(result, "text_content", "") or ""
    logger.info("markitdown: 提取 %d 字符", len(text))
    return text


# ---------------------------------------------------------------------------
# 压缩包解析器（ZIP）
# ---------------------------------------------------------------------------

def _extract_archive_markitdown(file_path: Path, start_page: int, max_pages: int | None) -> str:
    """使用 markitdown 解析 ZIP 压缩包（遍历内部内容）。"""
    if not module_exists("markitdown"):
        raise ImportError("未安装 markitdown")
    from markitdown import MarkItDown  # type: ignore

    logger.info("markitdown: 解析压缩包 %s (遍历内部文件)", file_path.name)
    result = MarkItDown().convert(str(file_path))
    text = getattr(result, "text_content", "") or ""
    logger.info("markitdown: 提取 %d 字符", len(text))
    return text


def _extract_pdf_ocr_tesseract(file_path: Path, start_page: int, max_pages: int | None) -> str:
    """Render PDF pages and run Tesseract OCR so OCR can join parser voting."""
    page_count = get_pdf_page_count(file_path)
    first = max(0, start_page - 1)
    if first >= page_count:
        return ""
    last = page_count if max_pages is None else min(page_count, first + max(0, max_pages))
    page_indices = list(range(first, last))
    pages = ocr_pdf_pages(file_path, page_indices, "chi_sim+eng", 200)
    return ocr_pages_to_text(pages)


# ---------------------------------------------------------------------------
# 表格提取（PDF 专用）
# ---------------------------------------------------------------------------

def normalize_table_cells(table: list[list[object | None]]) -> tuple[list[str], list[list[str]]]:
    """Normalize a raw table into headers and rows."""
    if not table:
        return [], []
    headers = [str(h).strip() if h is not None else "" for h in table[0]]
    rows: list[list[str]] = []
    for row in table[1:]:
        rows.append([str(cell).strip() if cell is not None else "" for cell in row])
    return headers, rows


def table_quality_metrics(table: dict[str, object]) -> dict[str, object]:
    """Score a table by density, shape consistency, and header usefulness."""
    headers = table.get("headers") if isinstance(table.get("headers"), list) else []
    rows = table.get("rows") if isinstance(table.get("rows"), list) else []
    column_count = max(len(headers), max((len(row) for row in rows if isinstance(row, list)), default=0))
    row_count = len(rows)
    cells = list(headers)
    for row in rows:
        if isinstance(row, list):
            cells.extend(row)
    total_cells = len(cells)
    non_empty_cells = sum(1 for cell in cells if str(cell or "").strip())
    density = non_empty_cells / max(1, total_cells)
    expected = column_count or 1
    consistent_rows = sum(1 for row in rows if isinstance(row, list) and len(row) == expected)
    consistency = consistent_rows / max(1, row_count)
    header_non_empty = sum(1 for cell in headers if str(cell or "").strip())
    header_score = header_non_empty / max(1, len(headers))
    size_score = min(1.0, (row_count * max(1, column_count)) / 24)
    score = (
        density * 0.36
        + consistency * 0.24
        + header_score * 0.20
        + size_score * 0.20
    )
    if row_count == 0 or column_count < 2:
        score *= 0.45
    return {
        "score": round(max(0.0, min(1.0, score)), 4),
        "row_count": row_count,
        "column_count": column_count,
        "non_empty_cells": non_empty_cells,
        "density": round(density, 4),
        "shape_consistency": round(consistency, 4),
        "header_score": round(header_score, 4),
        "size_score": round(size_score, 4),
    }


def annotate_tables(tables: list[dict[str, object]], method: str) -> list[dict[str, object]]:
    """Attach method and quality metrics to extracted tables."""
    annotated: list[dict[str, object]] = []
    for index, table in enumerate(tables, 1):
        item = dict(table)
        item.setdefault("table_index", index)
        item["method"] = method
        item["quality"] = table_quality_metrics(item)
        annotated.append(item)
    return annotated


def extract_tables_pdfplumber(pdf_path: Path, page_indices: list[int], strategy: str = "default") -> list[dict[str, object]]:
    """Extract PDF tables with pdfplumber using default or text-oriented settings."""
    if not module_exists("pdfplumber"):
        raise ImportError("表格提取需要 pdfplumber，运行: pip install pdfplumber")
    import pdfplumber  # type: ignore

    settings = None
    if strategy == "text":
        settings = {
            "vertical_strategy": "text",
            "horizontal_strategy": "text",
            "intersection_tolerance": 8,
            "snap_tolerance": 4,
            "join_tolerance": 4,
            "min_words_vertical": 2,
            "min_words_horizontal": 1,
        }

    tables: list[dict[str, object]] = []
    with pdfplumber.open(str(pdf_path)) as pdf:
        for page_idx in page_indices:
            if page_idx >= len(pdf.pages):
                continue
            page = pdf.pages[page_idx]
            page_tables = page.extract_tables(table_settings=settings) if settings else page.extract_tables()
            for tbl_idx, tbl in enumerate(page_tables):
                if not tbl:
                    continue
                headers, rows = normalize_table_cells(tbl)
                tables.append({
                    "page": page_idx + 1,
                    "table_index": tbl_idx + 1,
                    "headers": headers,
                    "rows": rows,
                })
    return tables


def split_text_columns(line: str) -> list[str]:
    """Split a text line into likely table columns."""
    if "\t" in line:
        parts = [part.strip() for part in line.split("\t")]
    elif "|" in line:
        parts = [part.strip() for part in line.strip("|").split("|")]
    else:
        parts = [part.strip() for part in re.split(r"\s{2,}", line.strip())]
    return [part for part in parts if part]


def extract_tables_pymupdf_text(pdf_path: Path, page_indices: list[int]) -> list[dict[str, object]]:
    """Heuristic table extraction from PyMuPDF text blocks for borderless tables."""
    if not module_exists("fitz"):
        raise ImportError("pymupdf-text 表格候选需要 PyMuPDF/fitz，运行: pip install pymupdf")
    import fitz  # type: ignore

    tables: list[dict[str, object]] = []
    doc = fitz.open(pdf_path)
    try:
        for page_idx in page_indices:
            if page_idx < 0 or page_idx >= doc.page_count:
                continue
            text = doc.load_page(page_idx).get_text("text")
            current: list[list[str]] = []
            table_index = 1
            for raw_line in text.splitlines() + [""]:
                cells = split_text_columns(raw_line)
                if len(cells) >= 2:
                    current.append(cells)
                    continue
                if len(current) >= 2:
                    width = max(len(row) for row in current)
                    rectangular = [row + [""] * (width - len(row)) for row in current]
                    headers = rectangular[0]
                    rows = rectangular[1:]
                    tables.append({
                        "page": page_idx + 1,
                        "table_index": table_index,
                        "headers": headers,
                        "rows": rows,
                    })
                    table_index += 1
                current = []
    finally:
        doc.close()
    return tables


def extract_tables_from_pdf(pdf_path: Path, page_indices: list[int]) -> list[dict[str, object]]:
    """Extract PDF tables with the default high-precision parser."""
    return annotate_tables(extract_tables_pdfplumber(pdf_path, page_indices, "default"), "pdfplumber")


def extract_tables_by_method(pdf_path: Path, page_indices: list[int], method: str) -> list[dict[str, object]]:
    """Extract tables with a named table parser/method."""
    if method == "pdfplumber":
        return annotate_tables(extract_tables_pdfplumber(pdf_path, page_indices, "default"), method)
    if method == "pdfplumber-text":
        return annotate_tables(extract_tables_pdfplumber(pdf_path, page_indices, "text"), method)
    if method == "pymupdf-text":
        return annotate_tables(extract_tables_pymupdf_text(pdf_path, page_indices), method)
    raise ValueError(f"unknown table parser: {method}")


def table_signature(table: dict[str, object]) -> tuple[int, int, int]:
    """Build a rough signature for table consensus."""
    quality = table.get("quality") if isinstance(table.get("quality"), dict) else table_quality_metrics(table)
    return (
        int(table.get("page") or 0),
        int(quality.get("row_count") or 0),
        int(quality.get("column_count") or 0),
    )


def build_table_vote_payload(
    pdf_path: Path,
    page_indices: list[int],
    methods: list[str],
) -> dict[str, object]:
    """Run table extraction methods and choose the best table set."""
    method_rows: list[dict[str, object]] = []
    signature_counts: Counter[tuple[int, int, int]] = Counter()
    for method in methods:
        try:
            start = time.perf_counter()
            tables = extract_tables_by_method(pdf_path, page_indices, method)
            seconds = round(time.perf_counter() - start, 3)
            for table in tables:
                signature_counts[table_signature(table)] += 1
            quality_scores = [
                float((table.get("quality") if isinstance(table.get("quality"), dict) else {}).get("score") or 0.0)
                for table in tables
            ]
            avg_quality = round(sum(quality_scores) / max(1, len(quality_scores)), 4)
            coverage = min(1.0, len(tables) / max(1, len(page_indices)))
            method_rows.append({
                "method": method,
                "status": "ok",
                "seconds": seconds,
                "table_count": len(tables),
                "avg_quality": avg_quality,
                "coverage_score": round(coverage, 4),
                "tables": tables,
                "error": None,
            })
        except Exception as exc:
            method_rows.append({
                "method": method,
                "status": "failed",
                "seconds": 0,
                "table_count": 0,
                "avg_quality": 0.0,
                "coverage_score": 0.0,
                "tables": [],
                "error": repr(exc),
            })

    for row in method_rows:
        consensus_hits = 0
        tables = row.get("tables") if isinstance(row.get("tables"), list) else []
        for table in tables:
            if isinstance(table, dict) and signature_counts[table_signature(table)] > 1:
                consensus_hits += 1
        consensus_score = consensus_hits / max(1, len(tables))
        row["consensus_score"] = round(consensus_score, 4)
        row["vote_score"] = round(
            min(
                1.0,
                float(row.get("avg_quality") or 0.0) * 0.48
                + consensus_score * 0.26
                + float(row.get("coverage_score") or 0.0) * 0.18
                + min(1.0, int(row.get("table_count") or 0) / 5) * 0.08,
            ),
            4,
        ) if row.get("status") == "ok" else 0.0

    winners = [row for row in method_rows if row.get("status") == "ok" and int(row.get("table_count") or 0) > 0]
    winners.sort(key=lambda row: (float(row.get("vote_score") or 0.0), float(row.get("avg_quality") or 0.0), int(row.get("table_count") or 0)), reverse=True)
    winner = winners[0] if winners else None
    diagnostics = []
    if winner:
        diagnostics.append(f"Selected table parser {winner.get('method')} by quality, coverage, and cross-method consensus.")
    else:
        diagnostics.append("No table parser produced usable tables.")
    if any(row.get("status") == "failed" for row in method_rows):
        diagnostics.append("One or more table parsers failed; see method rows for dependency/runtime errors.")
    return {
        "version": __version__,
        "type": "table-vote",
        "source_file": str(pdf_path),
        "generated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        "pages": [index + 1 for index in page_indices],
        "methods": methods,
        "winner": {
            key: winner.get(key)
            for key in ["method", "vote_score", "avg_quality", "consensus_score", "coverage_score", "table_count", "seconds"]
        } if winner else None,
        "diagnostics": diagnostics,
        "rows": method_rows,
        "best_tables": winner.get("tables") if winner else [],
    }


def table_vote_to_markdown(payload: dict[str, object]) -> str:
    """Render table vote payload as Markdown."""
    winner = payload.get("winner") if isinstance(payload.get("winner"), dict) else None
    lines = [
        "# Table Vote Report",
        "",
        f"- Source: `{payload.get('source_file')}`",
        f"- Pages: `{payload.get('pages')}`",
        f"- Winner: `{winner.get('method') if winner else None}`",
        f"- Vote score: `{winner.get('vote_score') if winner else 0}`",
        "",
        "## Methods",
        "",
        "| Method | Status | Vote | Avg quality | Consensus | Coverage | Tables | Time(s) | Error |",
        "|---|---|---:|---:|---:|---:|---:|---:|---|",
    ]
    rows = payload.get("rows") if isinstance(payload.get("rows"), list) else []
    for row in rows:
        if not isinstance(row, dict):
            continue
        error = str(row.get("error") or "").replace("|", "\\|")
        lines.append(
            f"| {row.get('method')} | {row.get('status')} | {row.get('vote_score')} | "
            f"{row.get('avg_quality')} | {row.get('consensus_score')} | {row.get('coverage_score')} | "
            f"{row.get('table_count')} | {row.get('seconds')} | {error[:180]} |"
        )
    diagnostics = payload.get("diagnostics") if isinstance(payload.get("diagnostics"), list) else []
    lines.extend(["", "## Diagnostics", ""])
    for item in diagnostics:
        lines.append(f"- {item}")
    best_tables = payload.get("best_tables") if isinstance(payload.get("best_tables"), list) else []
    if best_tables:
        lines.extend(["", "## Best Tables", "", tables_to_markdown(best_tables)])
    return "\n".join(lines) + "\n"


def write_table_vote_outputs(payload: dict[str, object], out_dir: Path, output_format: str) -> list[Path]:
    """Write table vote reports and best table outputs."""
    out_dir.mkdir(parents=True, exist_ok=True)
    written: list[Path] = []
    report_json = write_json_file(out_dir / "table_vote_report.json", payload)
    report_md = out_dir / "table_vote_report.md"
    report_md.write_text(table_vote_to_markdown(payload), encoding="utf-8")
    written.extend([report_json, report_md])
    tables = payload.get("best_tables") if isinstance(payload.get("best_tables"), list) else []
    if not tables:
        return written
    suffixes = ["md", "csv", "json"] if output_format == "all" else [output_format]
    for suffix in suffixes:
        path = out_dir / f"best_tables.{suffix}"
        if suffix == "md":
            path.write_text(tables_to_markdown(tables), encoding="utf-8")
        elif suffix == "txt":
            path.write_text(tables_to_markdown(tables), encoding="utf-8")
        elif suffix == "csv":
            path.write_text(tables_to_csv(tables), encoding="utf-8")
        elif suffix == "json":
            path.write_text(tables_to_json(tables), encoding="utf-8")
        written.append(path)
    return written


def tables_to_markdown(tables: list[dict]) -> str:
    """将表格列表转为 Markdown 格式。"""
    parts: list[str] = []
    for tbl in tables:
        parts.append(f"\n### Page {tbl['page']} - Table {tbl['table_index']}\n")
        parts.append("| " + " | ".join(tbl["headers"]) + " |")
        parts.append("|" + "|".join(["---"] * len(tbl["headers"])) + "|")
        for row in tbl["rows"]:
            parts.append("| " + " | ".join(row) + " |")
        parts.append("")
    return "\n".join(parts)


def tables_to_csv(tables: list[dict]) -> str:
    """将表格列表转为 CSV 格式（多个表格用空行分隔）。"""
    output = io.StringIO()
    writer = csv.writer(output)
    for i, tbl in enumerate(tables):
        if i > 0:
            writer.writerow([])  # 空行分隔
        writer.writerow([f"# Page {tbl['page']} Table {tbl['table_index']}"])
        writer.writerow(tbl["headers"])
        for row in tbl["rows"]:
            writer.writerow(row)
    return output.getvalue()


def tables_to_json(tables: list[dict]) -> str:
    """Render extracted tables as JSON."""
    return json.dumps(
        {
            "version": __version__,
            "format": "tables",
            "table_count": len(tables),
            "tables": tables,
        },
        ensure_ascii=False,
        indent=2,
    )


# ---------------------------------------------------------------------------
# 解析器注册表
# ---------------------------------------------------------------------------

_EXTRACTORS: dict[tuple[str, str], Callable[[Path, int, int | None], str]] = {
    # PDF
    ("pdf", "markitdown"):   _extract_pdf_markitdown,
    ("pdf", "pymupdf4llm"):  _extract_pdf_pymupdf4llm,
    ("pdf", "docling"):      _extract_pdf_docling,
    ("pdf", "pspdfkit"):     _extract_pdf_pspdfkit,
    ("pdf", "pymupdf"):      _extract_pdf_pymupdf,
    ("pdf", "pypdf"):        _extract_pdf_pypdf,
    ("pdf", "pdfplumber"):   _extract_pdf_pdfplumber,
    ("pdf", "pdfminer"):     _extract_pdf_pdfminer,
    ("pdf", "liteparse"):    _extract_pdf_liteparse,
    ("pdf", "opendataloader"): _extract_pdf_opendataloader,
    ("pdf", "ocr-tesseract"): _extract_pdf_ocr_tesseract,
    # Word
    ("word", "markitdown"):  _extract_word_markitdown,
    ("word", "python-docx"): _extract_word_docx,
    # PPT
    ("ppt", "markitdown"):   _extract_ppt_markitdown,
    ("ppt", "python-pptx"):  _extract_ppt_pptx,
    # Excel
    ("excel", "markitdown"): _extract_excel_markitdown,
    ("excel", "openpyxl"):   _extract_excel_openpyxl,
    # HTML
    ("html", "markitdown"):      _extract_html_markitdown,
    ("html", "beautifulsoup4"):  _extract_html_bs4,
    # 文本格式 (TXT, CSV, JSON, XML, MD)
    ("text", "markitdown"):  _extract_text_markitdown,
    # 图片格式 (JPG, PNG, GIF, BMP, TIFF, WebP)
    ("image", "markitdown"): _extract_image_markitdown,
    # 音频格式 (MP3, WAV)
    ("audio", "markitdown"): _extract_audio_markitdown,
    # 电子书 (EPub)
    ("epub", "markitdown"):  _extract_epub_markitdown,
    # 压缩包 (ZIP)
    ("archive", "markitdown"): _extract_archive_markitdown,
}

_PDF_PAGE_EXTRACTORS: dict[str, Callable[[Path, int, int | None], list[dict[str, object]]]] = {
    "pymupdf": _extract_pdf_pages_pymupdf,
    "pypdf": _extract_pdf_pages_pypdf,
    "pdfplumber": _extract_pdf_pages_pdfplumber,
    "pdfminer": _extract_pdf_pages_pdfminer,
}

_PARSER_ALIASES: dict[str, str] = {
    "fitz": "pymupdf",
    "pypdf2": "pypdf",
    "docx": "python-docx",
    "pptx": "python-pptx",
    "bs4": "beautifulsoup4",
    "llama-parse": "liteparse",
    "odl": "opendataloader",
    "opendataloader-pdf": "opendataloader",
    "pdf-to-markdown": "pspdfkit",
    "pspdfkit-pdf-to-markdown": "pspdfkit",
    "nutrient": "pspdfkit",
}


def resolve_parser_name(name: str) -> str:
    """解析器名称标准化，处理别名。"""
    return _PARSER_ALIASES.get(name, name)


def get_extractor(fmt: str, parser_name: str) -> Callable[[Path, int, int | None], str] | None:
    """获取指定格式和解析器的提取函数。"""
    return _EXTRACTORS.get((fmt, parser_name))


def get_pdf_page_extractor(parser_name: str) -> Callable[[Path, int, int | None], list[dict[str, object]]] | None:
    """Return page-by-page PDF extractor when available."""
    return _PDF_PAGE_EXTRACTORS.get(parser_name)


def list_all_parsers() -> list[str]:
    """列出所有已注册的解析器名称（去重）。"""
    names: set[str] = set()
    for _, parser_name in _EXTRACTORS:
        names.add(parser_name)
    return sorted(names)


# ---------------------------------------------------------------------------
# PDF 子集生成（仅 PDF 格式需要）
# ---------------------------------------------------------------------------

def create_subset_pdf_with_fitz(pdf_path: Path, start_page: int, max_pages: int, tmp_dir: Path) -> Path | None:
    if not module_exists("fitz"):
        return None
    import fitz  # type: ignore

    source = fitz.open(pdf_path)
    try:
        start_index, end_index = compute_page_range(source.page_count, start_page, max_pages)
        subset = fitz.open()
        subset.insert_pdf(source, from_page=start_index, to_page=end_index)
        subset_path = tmp_dir / f"{pdf_path.stem}.pages_{start_index + 1}_to_{end_index + 1}.pdf"
        subset.save(subset_path)
        subset.close()
        return subset_path
    finally:
        source.close()


def create_subset_pdf_with_pypdf(pdf_path: Path, start_page: int, max_pages: int, tmp_dir: Path) -> Path | None:
    reader_cls = None
    writer_cls = None

    if module_exists("pypdf"):
        from pypdf import PdfReader, PdfWriter  # type: ignore
        reader_cls = PdfReader
        writer_cls = PdfWriter
    elif module_exists("PyPDF2"):
        from PyPDF2 import PdfReader, PdfWriter  # type: ignore
        reader_cls = PdfReader
        writer_cls = PdfWriter
    else:
        return None

    reader = reader_cls(str(pdf_path))
    writer = writer_cls()
    start_index, end_index = compute_page_range(len(reader.pages), start_page, max_pages)
    for index in range(start_index, end_index + 1):
        writer.add_page(reader.pages[index])

    subset_path = tmp_dir / f"{pdf_path.stem}.pages_{start_index + 1}_to_{end_index + 1}.pdf"
    with subset_path.open("wb") as fh:
        writer.write(fh)
    return subset_path


def create_subset_pdf(pdf_path: Path, start_page: int, max_pages: int, tmp_dir: Path) -> Path | None:
    for creator_name, creator in [
        ("fitz", create_subset_pdf_with_fitz),
        ("pypdf", create_subset_pdf_with_pypdf),
    ]:
        try:
            subset = creator(pdf_path, start_page, max_pages, tmp_dir)
            if subset:
                return subset
        except Exception as exc:
            logger.debug("使用 %s 创建子 PDF 失败: %s", creator_name, exc)
    return None


# ---------------------------------------------------------------------------
# 相似度计算
# ---------------------------------------------------------------------------

def similarity_score(a: str, b: str, sample_chars: int) -> float:
    a_sample = a[:sample_chars]
    b_sample = b[:sample_chars]
    if not a_sample and not b_sample:
        return 1.0
    if not a_sample or not b_sample:
        return 0.0
    return round(difflib.SequenceMatcher(None, a_sample, b_sample, autojunk=True).ratio(), 4)


def choose_recommended_result(results: list[ParseResult]) -> ParseResult | None:
    """Choose a likely best parser result using quality first, then text volume and speed."""
    ok_results = [r for r in results if r.status == "ok" and r.non_space_chars > 0]
    if not ok_results:
        return None
    return max(ok_results, key=lambda r: (r.quality_score, r.non_space_chars, -r.seconds))


def structure_score(text: str) -> float:
    """Estimate whether output preserved useful Markdown/table/list structure."""
    lines = text.splitlines()
    if not lines:
        return 0.0
    headings = sum(1 for line in lines if line.lstrip().startswith("#"))
    table_lines = sum(1 for line in lines if line.strip().startswith("|") and line.strip().endswith("|"))
    list_lines = sum(1 for line in lines if re.match(r"^\s*(?:[-*+]|\d+[.)])\s+", line))
    page_markers = text.count("<!-- page ")
    raw = headings * 0.04 + table_lines * 0.025 + list_lines * 0.015 + page_markers * 0.02
    return round(min(1.0, raw), 4)


def repetition_metrics(text: str) -> dict[str, object]:
    """Estimate repeated content that can inflate coverage without adding value."""
    meaningful_lines = [
        re.sub(r"\s+", " ", line).strip()
        for line in text.splitlines()
        if len(re.sub(r"\s+", "", line)) >= 8
    ]
    if meaningful_lines:
        line_counts = Counter(meaningful_lines)
        duplicate_line_count = sum(count - 1 for count in line_counts.values() if count > 1)
        duplicate_line_ratio = duplicate_line_count / len(meaningful_lines)
        max_line_repeat = max(line_counts.values())
    else:
        duplicate_line_ratio = 0.0
        max_line_repeat = 0

    tokens = _RE_TOKEN.findall(compact_text(text))
    ngram_ratio = 0.0
    max_ngram_repeat = 0
    if len(tokens) >= 36:
        ngram_size = 6
        ngrams = ["".join(tokens[index:index + ngram_size]) for index in range(0, len(tokens) - ngram_size + 1)]
        ngram_counts = Counter(ngrams)
        duplicate_ngram_count = sum(count - 1 for count in ngram_counts.values() if count > 1)
        ngram_ratio = duplicate_ngram_count / max(1, len(ngrams))
        max_ngram_repeat = max(ngram_counts.values())

    penalty = 0.0
    if duplicate_line_ratio > 0.08:
        penalty += min(0.22, duplicate_line_ratio * 0.55)
    if ngram_ratio > 0.2:
        penalty += min(0.18, (ngram_ratio - 0.2) * 0.35)
    if max_line_repeat >= 3:
        penalty += min(0.08, (max_line_repeat - 2) * 0.02)

    score = max(0.0, 1.0 - min(0.45, penalty) / 0.45)
    return {
        "score": round(score, 4),
        "penalty": round(min(0.45, penalty), 4),
        "duplicate_line_ratio": round(duplicate_line_ratio, 4),
        "repeated_ngram_ratio": round(ngram_ratio, 4),
        "max_line_repeat": max_line_repeat,
        "max_ngram_repeat": max_ngram_repeat,
    }


def detect_vote_profile(
    requested_profile: str,
    fmt: str,
    normalized_texts: dict[str, str],
) -> dict[str, object]:
    """Resolve auto profile selection for vote workflows."""
    if requested_profile in {"contract", "bank_statement", "quotation", "purchase_order", "report", "annual_report"}:
        return {
            "requested": requested_profile,
            "resolved": requested_profile,
            "reason": f"explicit_{requested_profile}_profile",
            "candidates": [],
        }
    if requested_profile == "invoice":
        return {
            "requested": requested_profile,
            "resolved": "invoice",
            "reason": "explicit_invoice_profile",
            "candidates": [],
        }
    if requested_profile == "none":
        return {
            "requested": requested_profile,
            "resolved": "none",
            "reason": "explicit_no_profile",
            "candidates": [],
        }

    candidates: list[dict[str, object]] = []
    for parser_name, text in normalized_texts.items():
        if not text:
            continue
        classification = classify_document(text[:50000], {"format": fmt})
        signals = classification.get("signals") if isinstance(classification.get("signals"), dict) else {}
        candidates.append({
            "parser": parser_name,
            "profile": classification.get("profile"),
            "confidence": classification.get("confidence"),
            "invoice_hits": signals.get("invoice_hits"),
            "contract_hits": signals.get("contract_hits"),
            "bank_statement_hits": signals.get("bank_statement_hits"),
            "quotation_hits": signals.get("quotation_hits"),
            "purchase_order_hits": signals.get("purchase_order_hits"),
            "report_hits": signals.get("report_hits"),
            "annual_report_hits": signals.get("annual_report_hits"),
        })

    invoice_candidates = [
        item
        for item in candidates
        if item.get("profile") == "invoice" or int(item.get("invoice_hits") or 0) >= 3
    ]
    if invoice_candidates:
        return {
            "requested": requested_profile,
            "resolved": "invoice",
            "reason": "auto_detected_invoice",
            "candidates": candidates,
        }
    structured_profiles = {"contract", "bank_statement", "quotation", "purchase_order", "report", "annual_report"}
    profile_votes = Counter(
        str(item.get("profile"))
        for item in candidates
        if item.get("profile") in structured_profiles
    )
    if profile_votes:
        profile_name, _ = profile_votes.most_common(1)[0]
        return {
            "requested": requested_profile,
            "resolved": profile_name,
            "reason": f"auto_detected_{profile_name}",
            "candidates": candidates,
        }
    return {
        "requested": requested_profile,
        "resolved": "none",
        "reason": "auto_no_structured_profile",
        "candidates": candidates,
    }


def invoice_vote_metrics(fields: dict[str, object]) -> dict[str, object]:
    """Score structured invoice extraction for parser voting."""
    max_field_score = 23
    field_score = score_invoice_fields(fields)
    completeness = min(1.0, field_score / max_field_score)
    validation = fields.get("validation") if isinstance(fields.get("validation"), dict) else {}
    status = validation.get("status") if isinstance(validation, dict) else "unknown"
    missing = validation.get("missing_fields") if isinstance(validation, dict) else []
    missing_count = len(missing) if isinstance(missing, list) else 0
    checks = validation.get("checks") if isinstance(validation, dict) else []
    failed_checks = [
        check for check in checks
        if isinstance(check, dict) and not bool(check.get("passed"))
    ] if isinstance(checks, list) else []
    item_metrics = invoice_item_repetition_metrics(fields.get("items"))
    item_count = len(fields.get("items")) if isinstance(fields.get("items"), list) else 0

    validation_adjustment = 0.12 if status == "ok" else -0.12 if status == "failed_checks" else -0.08
    no_item_penalty = 0.08 if item_count == 0 else 0.0
    failed_check_penalty = min(0.2, len(failed_checks) * 0.04)
    missing_penalty = min(0.16, missing_count * 0.025)
    duplicate_item_penalty = float(item_metrics.get("penalty") or 0.0)
    score = max(
        0.0,
        min(
            1.0,
            completeness
            + validation_adjustment
            - no_item_penalty
            - failed_check_penalty
            - missing_penalty
            - duplicate_item_penalty,
        ),
    )
    return {
        "score": round(score, 4),
        "field_score": field_score,
        "completeness_score": round(completeness, 4),
        "validation_status": status,
        "missing_fields": missing if isinstance(missing, list) else [],
        "failed_check_count": len(failed_checks),
        "item_count": item_count,
        "duplicate_item_ratio": item_metrics.get("duplicate_item_ratio", 0.0),
        "duplicate_item_penalty": round(duplicate_item_penalty, 4),
        "penalty": round(no_item_penalty + failed_check_penalty + missing_penalty + duplicate_item_penalty, 4),
    }


def first_nonempty(*values: object) -> object | None:
    """Return the first non-empty value."""
    for value in values:
        if value is None:
            continue
        if isinstance(value, str) and not value.strip():
            continue
        return value
    return None


def amount_first(text: str) -> float | None:
    """Extract the first currency-looking amount from text."""
    value = regex_first(r"(?:金额|合计|总价|总金额|价税合计|合同价款|报价总计|订单总额)[:：]?[¥￥]?\s*([0-9][0-9,]*(?:\.[0-9]{1,2})?)", text)
    if value:
        return decimal_from_text(value)
    fallback = regex_first(r"[¥￥]\s*([0-9][0-9,]*(?:\.[0-9]{1,2})?)", text)
    return decimal_from_text(fallback)


def date_first(text: str) -> str | None:
    """Extract the first common date string."""
    return regex_first(r"([0-9]{4}年[0-9]{1,2}月[0-9]{1,2}日|[0-9]{4}[-/][0-9]{1,2}[-/][0-9]{1,2})", text)


def keyword_lines(text: str, keywords: list[str], limit: int = 8) -> list[str]:
    """Return concise source lines that contain business keywords."""
    lines: list[str] = []
    seen: set[str] = set()
    for raw in text.splitlines():
        line = normalize_text(raw)
        if not line or line in seen:
            continue
        if any(keyword in line for keyword in keywords):
            seen.add(line)
            lines.append(line[:240])
        if len(lines) >= limit:
            break
    return lines


def business_vote_metrics(fields: dict[str, object]) -> dict[str, object]:
    """Score lightweight business extraction completeness."""
    ignored = {
        "version",
        "schema_version",
        "profile",
        "source_file",
        "parser",
        "quality_score",
        "quality_label",
        "raw_text_chars",
        "validation",
    }
    present_keys: list[str] = []
    for key, value in fields.items():
        if key in ignored:
            continue
        if value is None:
            continue
        if isinstance(value, str) and not value.strip():
            continue
        if isinstance(value, list) and not value:
            continue
        present_keys.append(key)
    keyword_like = sum(
        1
        for key in present_keys
        if key.endswith("terms") or key.endswith("clauses") or key.startswith("key_") or key.endswith("keywords")
    )
    field_count = len(present_keys)
    score = min(1.0, field_count / 8)
    if keyword_like:
        score = min(1.0, score + min(0.12, keyword_like * 0.03))
    return {
        "score": round(score, 4),
        "status": "ok" if field_count >= 3 else "partial" if field_count else "empty",
        "field_count": field_count,
        "present_keys": present_keys,
    }


def extract_business_fields_from_text(
    profile: str,
    text: str,
    source_file: Path | None = None,
    parser_name: str = "unknown",
    result: ParseResult | None = None,
) -> dict[str, object]:
    """Extract lightweight structured fields for non-invoice business profiles."""
    normalized = normalize_text(text)
    fields: dict[str, object] = {
        "version": __version__,
        "schema_version": CUSTOMER_BEST_SCHEMA_VERSION,
        "profile": profile,
        "source_file": str(source_file) if source_file else None,
        "parser": parser_name,
        "quality_score": result.quality_score if result else None,
        "quality_label": result.quality_label if result else None,
        "raw_text_chars": len(normalized),
    }

    if profile == "contract":
        fields.update({
            "contract_number": regex_first(r"(?:合同编号|协议编号|编号)[:：]?\s*([A-Za-z0-9\-_/（）()]{3,60})", normalized),
            "party_a": regex_first(r"(?:甲方|委托方|买方)[:：]?\s*([^\n\r]{2,80})", normalized),
            "party_b": regex_first(r"(?:乙方|受托方|卖方)[:：]?\s*([^\n\r]{2,80})", normalized),
            "signing_date": first_nonempty(regex_first(r"(?:签订日期|签署日期|合同日期)[:：]?\s*([0-9]{4}年[0-9]{1,2}月[0-9]{1,2}日|[0-9]{4}[-/][0-9]{1,2}[-/][0-9]{1,2})", normalized), date_first(normalized)),
            "amount": amount_first(normalized),
            "term": regex_first(r"(?:履行期限|服务期限|合同期限|有效期)[:：]?\s*([^\n\r。；;]{2,120})", normalized),
            "payment_terms": regex_first(r"(?:付款方式|付款条件|支付方式|结算方式)[:：]?\s*([^\n\r。；;]{2,160})", normalized),
            "dispute_resolution": regex_first(r"(?:争议解决|管辖法院|仲裁)[:：]?\s*([^\n\r。；;]{2,160})", normalized),
            "key_clauses": keyword_lines(normalized, ["甲方", "乙方", "合同金额", "付款", "违约", "争议解决", "履行期限"]),
        })
    elif profile == "bank_statement":
        fields.update({
            "account_name": regex_first(r"(?:账户名称|户名|客户名称)[:：]?\s*([^\n\r]{2,80})", normalized),
            "account_number": regex_first(r"(?:账号|账户|卡号)[:：]?\s*([0-9\s*]{8,40})", normalized),
            "bank_name": regex_first(r"(?:开户行|银行名称|银行)[:：]?\s*([^\n\r]{2,80})", normalized),
            "statement_period": regex_first(r"(?:期间|起止日期|交易日期)[:：]?\s*([0-9年月日/\-\s至~]{8,60})", normalized),
            "currency": regex_first(r"(?:币种|货币)[:：]?\s*([A-Z]{3}|人民币|美元|欧元|港币)", normalized),
            "opening_balance": decimal_from_text(regex_first(r"(?:期初余额|上期余额)[:：]?[¥￥]?\s*([0-9][0-9,]*(?:\.[0-9]{1,2})?)", normalized)),
            "closing_balance": decimal_from_text(regex_first(r"(?:期末余额|余额)[:：]?[¥￥]?\s*([0-9][0-9,]*(?:\.[0-9]{1,2})?)", normalized)),
            "transaction_keywords": keyword_lines(normalized, ["交易日期", "借方", "贷方", "余额", "交易金额", "对方户名", "摘要"]),
        })
    elif profile == "quotation":
        fields.update({
            "quotation_number": regex_first(r"(?:报价单号|报价编号|编号)[:：]?\s*([A-Za-z0-9\-_/]{3,60})", normalized),
            "quotation_date": first_nonempty(regex_first(r"(?:报价日期|日期)[:：]?\s*([0-9]{4}年[0-9]{1,2}月[0-9]{1,2}日|[0-9]{4}[-/][0-9]{1,2}[-/][0-9]{1,2})", normalized), date_first(normalized)),
            "customer_name": regex_first(r"(?:客户名称|客户|需方|采购方)[:：]?\s*([^\n\r]{2,80})", normalized),
            "supplier_name": regex_first(r"(?:供应商|报价方|供方|销售方)[:：]?\s*([^\n\r]{2,80})", normalized),
            "total_amount": amount_first(normalized),
            "valid_until": regex_first(r"(?:有效期|报价有效期|有效至)[:：]?\s*([^\n\r。；;]{2,80})", normalized),
            "payment_terms": regex_first(r"(?:付款方式|付款条件|支付方式)[:：]?\s*([^\n\r。；;]{2,120})", normalized),
            "contact": regex_first(r"(?:联系人|报价人)[:：]?\s*([^\s\n\r]{2,40})", normalized),
            "price_terms": keyword_lines(normalized, ["报价", "单价", "金额", "含税", "税率", "有效期", "付款"]),
        })
    elif profile == "purchase_order":
        fields.update({
            "purchase_order_number": regex_first(r"(?:采购订单号|订单编号|采购单号|PO)[:：]?\s*([A-Za-z0-9\-_/]{3,60})", normalized),
            "order_date": first_nonempty(regex_first(r"(?:订单日期|下单日期|日期)[:：]?\s*([0-9]{4}年[0-9]{1,2}月[0-9]{1,2}日|[0-9]{4}[-/][0-9]{1,2}[-/][0-9]{1,2})", normalized), date_first(normalized)),
            "buyer_name": regex_first(r"(?:采购方|买方|需方)[:：]?\s*([^\n\r]{2,80})", normalized),
            "supplier_name": regex_first(r"(?:供应商|卖方|供方)[:：]?\s*([^\n\r]{2,80})", normalized),
            "delivery_date": regex_first(r"(?:交货日期|交付日期|到货日期)[:：]?\s*([0-9]{4}年[0-9]{1,2}月[0-9]{1,2}日|[0-9]{4}[-/][0-9]{1,2}[-/][0-9]{1,2}|[^\n\r。；;]{2,60})", normalized),
            "total_amount": amount_first(normalized),
            "payment_terms": regex_first(r"(?:付款方式|付款条件|支付方式|结算方式)[:：]?\s*([^\n\r。；;]{2,120})", normalized),
            "order_terms": keyword_lines(normalized, ["采购", "供应商", "交货", "数量", "单价", "金额", "付款"]),
        })
    elif profile == "annual_report":
        fields.update({
            "company_name": regex_first(r"([^\n\r]{2,80}(?:股份有限公司|有限公司|集团有限公司))", normalized),
            "report_year": regex_first(r"([12][0-9]{3})\s*(?:年|年度报告|年报)", normalized),
            "report_date": date_first(normalized),
            "total_assets": decimal_from_text(regex_first(r"(?:资产总计|总资产)[:：]?\s*([0-9][0-9,]*(?:\.[0-9]{1,2})?)", normalized)),
            "revenue": decimal_from_text(regex_first(r"(?:营业收入|主营业务收入)[:：]?\s*([0-9][0-9,]*(?:\.[0-9]{1,2})?)", normalized)),
            "net_profit": decimal_from_text(regex_first(r"(?:净利润|归属于.*?净利润)[:：]?\s*([0-9][0-9,]*(?:\.[0-9]{1,2})?)", normalized)),
            "auditor": regex_first(r"(?:会计师事务所|审计机构)[:：]?\s*([^\n\r]{2,80})", normalized),
            "key_sections": keyword_lines(normalized, ["年度报告", "资产负债表", "利润表", "现金流量表", "董事会", "审计报告"]),
        })
    else:
        fields.update({
            "title": first_nonempty(regex_first(r"^#\s*(.+)$", normalized, re.M), regex_first(r"^([^\n\r]{4,80})$", normalized, re.M)),
            "report_date": date_first(normalized),
            "organization": regex_first(r"(?:编制单位|报告单位|公司|机构)[:：]?\s*([^\n\r]{2,80})", normalized),
            "summary": regex_first(r"(?:摘要|概述)[:：]?\s*([^\n\r。]{10,220})", normalized),
            "conclusion": regex_first(r"(?:结论|建议)[:：]?\s*([^\n\r。]{10,220})", normalized),
            "amount": amount_first(normalized),
            "key_findings": keyword_lines(normalized, ["摘要", "结论", "指标", "同比", "环比", "风险", "建议", "附录"]),
        })

    metrics = business_vote_metrics(fields)
    fields["validation"] = {
        "status": "ok" if metrics.get("field_count", 0) >= 3 else "partial",
        "field_count": metrics.get("field_count", 0),
        "present_keys": metrics.get("present_keys", []),
    }
    return fields


def collect_parser_outputs(
    file_path: Path,
    fmt: str,
    parser_names: list[str],
    out_dir: Path,
    output_format: str,
    start_page: int,
    max_pages: int | None,
    parallel: bool = False,
    progress: bool = False,
    timeout_seconds: float | None = None,
    health_cache_path: Path | None = None,
    use_health_cache: bool = False,
    result_cache_dir: Path | None = None,
    use_result_cache: bool = False,
) -> tuple[list[ParseResult], dict[str, str], Path]:
    """Run parser candidates and write their raw outputs."""
    out_dir.mkdir(parents=True, exist_ok=True)
    if use_health_cache and health_cache_path is None:
        health_cache_path = default_parser_health_cache_path(out_dir)
    health_cache = load_parser_health_cache(health_cache_path) if use_health_cache else {"version": __version__, "entries": {}}
    if use_result_cache and result_cache_dir is None:
        result_cache_dir = default_parser_result_cache_dir(out_dir)
    actual_path = file_path
    temp_dir_obj: tempfile.TemporaryDirectory[str] | None = None
    subset_skip_reason: str | None = None
    if fmt == "pdf" and any(name in PDF_SUBSET_TARGET_PARSERS for name in parser_names):
        temp_dir_obj = tempfile.TemporaryDirectory(prefix="parse_candidates_")
        actual_path, subset_skip_reason = prepare_pdf_subset_target(
            file_path, start_page, max_pages, Path(temp_dir_obj.name)
        )

    def run_parser(name: str) -> tuple[ParseResult, str]:
        cache_key = parser_health_cache_key(file_path, fmt, name, start_page, max_pages)
        if use_health_cache:
            skip, reason = should_skip_parser_from_health_cache(health_cache, cache_key)
            if skip:
                return ParseResult(
                    parser=name,
                    status="skipped",
                    seconds=0,
                    error=f"parser_health_cache: {reason}",
                ), ""

        if name in PDF_SUBSET_TARGET_PARSERS and subset_skip_reason:
            return ParseResult(parser=name, status="skipped", seconds=0, error=subset_skip_reason), ""

        extract_func = get_extractor(fmt, name)
        if extract_func is None:
            return ParseResult(parser=name, status="skipped", seconds=0, error=f"未注册的解析器：{name}"), ""

        result_cache_key = parser_result_cache_key(file_path, fmt, name, start_page, max_pages, output_format)
        if use_result_cache:
            cached = load_parser_result_cache(result_cache_dir, result_cache_key)
            if cached is not None:
                start = time.perf_counter()
                normalized = str(cached.get("normalized") or "")
                out_file, normalized = write_extracted_outputs(out_dir, name, normalized, output_format)
                result = text_stats(name, "ok", time.perf_counter() - start, normalized, out_file)
                result.cache_status = "hit"
                if use_health_cache:
                    update_parser_health_cache(health_cache, cache_key, name, result, timeout_seconds)
                return result, normalized

        target = actual_path if (name in PDF_SUBSET_TARGET_PARSERS and fmt == "pdf") else file_path
        start = time.perf_counter()
        try:
            if timeout_seconds and timeout_seconds > 0:
                result, normalized = extract_one_subprocess(
                    target,
                    fmt,
                    name,
                    out_dir,
                    output_format,
                    start_page,
                    max_pages,
                    timeout_seconds,
                )
            else:
                raw_text = extract_func(target, start_page, max_pages)
                out_file, normalized = write_extracted_outputs(out_dir, name, raw_text, output_format)
                result = text_stats(name, "ok", time.perf_counter() - start, normalized, out_file)
            result.cache_status = "miss" if use_result_cache and result.status == "ok" else None
            if use_result_cache:
                write_parser_result_cache(result_cache_dir, result_cache_key, name, result, normalized)
            if use_health_cache:
                update_parser_health_cache(health_cache, cache_key, name, result, timeout_seconds)
            return result, normalized
        except Exception as exc:
            result = ParseResult(
                parser=name,
                status="failed",
                seconds=round(time.perf_counter() - start, 3),
                error=repr(exc),
            )
            if use_health_cache:
                update_parser_health_cache(health_cache, cache_key, name, result, timeout_seconds)
            return result, ""

    results: list[ParseResult] = []
    normalized_texts: dict[str, str] = {}
    try:
        if parallel:
            if progress:
                print("并行解析中...")
            with ThreadPoolExecutor(max_workers=min(len(parser_names), 4)) as pool:
                futures = {pool.submit(run_parser, name): name for name in parser_names}
                for future in as_completed(futures):
                    result, text = future.result()
                    results.append(result)
                    if text:
                        normalized_texts[result.parser] = text
                    if progress:
                        icon = "OK" if result.status == "ok" else "FAIL" if result.status == "failed" else "SKIP"
                        print(f"  [{icon}] {result.parser}: {result.chars} chars, {result.seconds}s")
        else:
            for name in parser_names:
                if progress:
                    print(f"解析中: {name}...", end=" ", flush=True)
                result, text = run_parser(name)
                results.append(result)
                if text:
                    normalized_texts[result.parser] = text
                if progress:
                    if result.status == "ok":
                        print(f"OK {result.chars} chars, {result.seconds}s")
                    elif result.status == "failed":
                        print(f"FAIL {result.error}")
                    else:
                        print(f"SKIP {result.error}")
    finally:
        if use_health_cache:
            write_parser_health_cache(health_cache_path, health_cache)
        if temp_dir_obj:
            temp_dir_obj.cleanup()

    return results, normalized_texts, actual_path


def validate_parser_selection(fmt: str, parsers: str | None) -> tuple[list[str], list[str]]:
    """Return supported parser names and unsupported names from an optional comma list."""
    available_parsers = get_parsers_for_format(fmt)
    if parsers:
        parser_names = [resolve_parser_name(n.strip().lower()) for n in parsers.split(",") if n.strip()]
    else:
        parser_names = available_parsers
    valid_names = set(available_parsers)
    unknown = [n for n in parser_names if n not in valid_names]
    parser_names = [n for n in parser_names if n in valid_names]
    return parser_names, unknown


def build_vote_payload(
    file_path: Path,
    actual_path: Path,
    fmt: str,
    results: list[ParseResult],
    normalized_texts: dict[str, str],
    similarity_chars: int,
    min_quality: float,
    fail_on_bad: bool,
    profile: str = "auto",
    preflight_probe: dict[str, object] | None = None,
) -> dict[str, object]:
    """Score parser outputs by quality, consensus, coverage, structure, and profile checks."""
    ok_results = [r for r in results if r.status == "ok" and r.non_space_chars > 0]
    max_chars = max((r.non_space_chars for r in ok_results), default=0)
    profile_info = detect_vote_profile(profile, fmt, normalized_texts)
    resolved_profile = str(profile_info.get("resolved") or "none")
    invoice_enabled = resolved_profile == "invoice"
    business_enabled = resolved_profile in BUSINESS_STRUCTURED_PROFILES
    weights = {
        "quality_score": 0.34 if invoice_enabled else 0.36 if business_enabled else 0.42,
        "consensus_score": 0.18 if invoice_enabled else 0.20 if business_enabled else 0.24,
        "coverage_score": 0.10 if invoice_enabled else 0.12 if business_enabled else 0.16,
        "structure_score": 0.08 if invoice_enabled else 0.10,
        "repetition_score": 0.07 if invoice_enabled else 0.07 if business_enabled else 0.08,
    }
    if invoice_enabled:
        weights["invoice_score"] = 0.23
    if business_enabled:
        weights["business_score"] = 0.15
    votes: list[dict[str, object]] = []

    for result in results:
        text = normalized_texts.get(result.parser, "")
        if result.status != "ok" or result.non_space_chars <= 0:
            votes.append({
                "parser": result.parser,
                "status": result.status,
                "eligible": False,
                "vote_score": 0.0,
                "reason": result.error or "no usable output",
            })
            continue

        peers = [peer for peer in ok_results if peer.parser != result.parser]
        if peers:
            similarities = [
                similarity_score(text, normalized_texts.get(peer.parser, ""), similarity_chars)
                for peer in peers
            ]
            consensus = round(sum(similarities) / len(similarities), 4)
        else:
            consensus = 0.0

        coverage = round(result.non_space_chars / max(1, max_chars), 4)
        structure = structure_score(text)
        repetition = repetition_metrics(text)
        repetition_score = float(repetition.get("score") or 0.0)
        duplicate_penalty = float(repetition.get("penalty") or 0.0)
        invoice_fields: dict[str, object] | None = None
        invoice_metrics: dict[str, object] | None = None
        invoice_score = 0.0
        structured_fields: dict[str, object] | None = None
        business_metrics: dict[str, object] | None = None
        business_score = 0.0
        if invoice_enabled:
            try:
                invoice_fields = extract_invoice_fields_from_text(text, file_path, result.parser, result, strict=True)
                invoice_metrics = invoice_vote_metrics(invoice_fields)
                invoice_score = float(invoice_metrics.get("score") or 0.0)
            except Exception as exc:
                invoice_metrics = {
                    "score": 0.0,
                    "validation_status": "failed_extraction",
                    "error": repr(exc),
                }
        if business_enabled:
            try:
                structured_fields = extract_business_fields_from_text(resolved_profile, text, file_path, result.parser, result)
                business_metrics = business_vote_metrics(structured_fields)
                business_score = float(business_metrics.get("score") or 0.0)
            except Exception as exc:
                business_metrics = {
                    "score": 0.0,
                    "status": "failed_extraction",
                    "error": repr(exc),
                }
        penalty = 0.0
        if result.quality_label in QUALITY_BAD_LABELS:
            penalty += 0.35
        if min_quality > 0 and result.quality_score < min_quality:
            penalty += 0.15
        if result.cid_markers:
            penalty += min(0.25, result.cid_markers / max(1, result.non_space_chars) * 8)
        if result.control_chars:
            penalty += min(0.2, result.control_chars / max(1, result.non_space_chars) * 6)
        penalty += duplicate_penalty

        raw_score = (
            result.quality_score * weights["quality_score"]
            + consensus * weights["consensus_score"]
            + coverage * weights["coverage_score"]
            + structure * weights["structure_score"]
            + repetition_score * weights["repetition_score"]
        )
        if invoice_enabled:
            raw_score += invoice_score * weights["invoice_score"]
        if business_enabled:
            raw_score += business_score * weights["business_score"]
        vote_score = round(
            max(
                0.0,
                min(
                    1.0,
                    raw_score - penalty,
                ),
            ),
            4,
        )
        vote_item: dict[str, object] = {
            "parser": result.parser,
            "status": result.status,
            "eligible": True,
            "vote_score": vote_score,
            "raw_vote_score": round(raw_score, 4),
            "quality_score": result.quality_score,
            "quality_label": result.quality_label,
            "consensus_score": consensus,
            "coverage_score": coverage,
            "structure_score": structure,
            "repetition_score": repetition_score,
            "repetition_metrics": repetition,
            "duplicate_penalty": round(duplicate_penalty, 4),
            "penalty": round(penalty, 4),
            "chars": result.chars,
            "non_space_chars": result.non_space_chars,
            "seconds": result.seconds,
            "output_file": result.output_file,
            "cache_status": result.cache_status,
        }
        if invoice_enabled:
            vote_item["invoice_score"] = round(invoice_score, 4)
            vote_item["invoice_metrics"] = invoice_metrics
            if invoice_fields is not None:
                vote_item["invoice_fields"] = invoice_fields
        if business_enabled:
            vote_item["business_score"] = round(business_score, 4)
            vote_item["business_metrics"] = business_metrics
            if structured_fields is not None:
                vote_item["structured_fields"] = structured_fields
        votes.append(vote_item)

    eligible_votes = [vote for vote in votes if vote.get("eligible")]
    eligible_votes.sort(
        key=lambda vote: (
            float(vote.get("vote_score", 0.0)),
            float(vote.get("invoice_score", 0.0)),
            float(vote.get("business_score", 0.0)),
            float(vote.get("repetition_score", 0.0)),
            float(vote.get("quality_score", 0.0)),
            int(vote.get("non_space_chars", 0)),
        ),
        reverse=True,
    )
    winner = eligible_votes[0] if eligible_votes else None
    recommended = next((r for r in results if winner and r.parser == winner["parser"]), None)
    gate = quality_gate_status(recommended, min_quality, fail_on_bad)
    diagnostics = build_diagnostics(results)
    if len(ok_results) < 2:
        diagnostics.append("Only one parser produced usable output, so the vote is mostly quality-based rather than consensus-based.")
    if not winner:
        diagnostics.append("No parser produced a customer-deliverable output.")
    if invoice_enabled:
        diagnostics.append("Invoice profile is enabled: structured field completeness, validation status, and duplicate line-item checks are included in the vote.")
    if business_enabled:
        diagnostics.append(f"{resolved_profile} profile is enabled: business field coverage and keyword extraction are included in the vote.")
    if any(float(vote.get("duplicate_penalty") or 0.0) > 0 for vote in votes if isinstance(vote, dict)):
        diagnostics.append("One or more parser outputs were penalized for repeated lines or repeated text spans.")
    if preflight_probe:
        probe_summary = preflight_probe.get("summary") if isinstance(preflight_probe.get("summary"), dict) else {}
        diagnostics.append(
            "Preflight probe enabled: only ready parsers entered the final vote "
            f"({probe_summary.get('ready_count', 0)}/{probe_summary.get('parser_count', 0)} ready)."
        )

    return {
        "version": __version__,
        "type": "vote",
        "source_file": str(file_path),
        "actual_file": str(actual_path),
        "format": fmt,
        "generated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        "similarity_chars": similarity_chars,
        "profile": profile_info,
        "preflight_probe": preflight_probe,
        "weights": weights,
        "winner": winner,
        "recommendation": asdict(recommended) if recommended else None,
        "quality_gate": gate,
        "diagnostics": diagnostics,
        "votes": votes,
        "results": [asdict(r) for r in results],
    }


def vote_to_markdown(payload: dict[str, object]) -> str:
    """Render a human-readable vote report."""
    winner = payload.get("winner") if isinstance(payload.get("winner"), dict) else None
    gate = payload.get("quality_gate") if isinstance(payload.get("quality_gate"), dict) else {}
    profile = payload.get("profile") if isinstance(payload.get("profile"), dict) else {}
    lines = [
        "# Parser Vote Report",
        "",
        f"- Source: `{payload.get('source_file')}`",
        f"- Format: `{payload.get('format')}`",
        f"- Profile: `{profile.get('resolved', 'none')}` (`{profile.get('reason', '')}`)",
        f"- Winner: `{winner.get('parser') if winner else None}`",
        f"- Vote score: `{winner.get('vote_score') if winner else 0}`",
        f"- Quality gate: `{gate.get('passed')}` (`{gate.get('reason')}`)",
        "",
        "## Decision",
        "",
    ]
    if winner:
        lines.extend([
            f"- Selected parser: `{winner.get('parser')}`",
            f"- Quality: `{winner.get('quality_score')}` (`{winner.get('quality_label')}`)",
            f"- Consensus: `{winner.get('consensus_score')}`",
            f"- Coverage: `{winner.get('coverage_score')}`",
            f"- Structure: `{winner.get('structure_score')}`",
            f"- Repetition: `{winner.get('repetition_score')}` (penalty `{winner.get('duplicate_penalty')}`)",
            f"- Output: `{winner.get('output_file')}`",
            "",
        ])
        if profile.get("resolved") == "invoice":
            metrics = winner.get("invoice_metrics") if isinstance(winner.get("invoice_metrics"), dict) else {}
            lines.extend([
                f"- Invoice score: `{winner.get('invoice_score')}`",
                f"- Invoice validation: `{metrics.get('validation_status')}`",
                f"- Invoice items: `{metrics.get('item_count')}`",
                "",
            ])
    else:
        lines.extend(["No parser produced a usable output.", ""])

    preflight_probe = payload.get("preflight_probe") if isinstance(payload.get("preflight_probe"), dict) else None
    if preflight_probe:
        summary = preflight_probe.get("summary") if isinstance(preflight_probe.get("summary"), dict) else {}
        ready = summary.get("ready_parsers") if isinstance(summary.get("ready_parsers"), list) else []
        lines.extend([
            "## Preflight Probe",
            "",
            f"- Ready parsers: `{summary.get('ready_count', 0)}` / `{summary.get('parser_count', 0)}`",
            f"- Dependency missing: `{summary.get('dependency_missing_count', 0)}`",
            f"- Runtime failed: `{summary.get('runtime_failed_count', 0)}`",
            f"- Quality failed: `{summary.get('quality_failed_count', 0)}`",
            f"- Final vote parsers: `{', '.join(str(name) for name in ready)}`",
            "",
        ])

    lines.extend([
        "## Votes",
        "",
        "| Parser | Eligible | Vote | Quality | Consensus | Coverage | Structure | Repetition | Invoice | Business | Penalty | Chars | Time(s) | Cache |",
        "|---|---:|---:|---:|---:|---:|---:|---:|---:|---:|---:|---:|---:|---|",
    ])
    votes = payload.get("votes") if isinstance(payload.get("votes"), list) else []
    for item in votes:
        if not isinstance(item, dict):
            continue
        lines.append(
            f"| {item.get('parser')} | {item.get('eligible')} | {item.get('vote_score', 0)} | "
            f"{item.get('quality_score', '')} | {item.get('consensus_score', '')} | "
            f"{item.get('coverage_score', '')} | {item.get('structure_score', '')} | "
            f"{item.get('repetition_score', '')} | {item.get('invoice_score', '')} | "
            f"{item.get('business_score', '')} | {item.get('penalty', '')} | {item.get('chars', '')} | "
            f"{item.get('seconds', '')} | {item.get('cache_status', '')} |"
        )

    diagnostics = payload.get("diagnostics") if isinstance(payload.get("diagnostics"), list) else []
    lines.extend(["", "## Diagnostics", ""])
    for item in diagnostics:
        lines.append(f"- {item}")
    return "\n".join(lines) + "\n"


def write_vote_outputs(
    payload: dict[str, object],
    normalized_texts: dict[str, str],
    out_dir: Path,
    output_format: str,
) -> list[Path]:
    """Write vote reports and best raw text output files."""
    out_dir.mkdir(parents=True, exist_ok=True)
    written: list[Path] = []
    report_json = write_json_file(out_dir / "vote_report.json", payload)
    report_md = out_dir / "vote_report.md"
    report_md.write_text(vote_to_markdown(payload), encoding="utf-8")
    written.extend([report_json, report_md])

    winner = payload.get("winner") if isinstance(payload.get("winner"), dict) else None
    winner_parser = str(winner.get("parser")) if winner else ""
    best_text = normalized_texts.get(winner_parser, "")
    if not best_text:
        return written

    suffixes = output_suffixes(output_format)
    for suffix in suffixes:
        path = out_dir / f"best.{suffix}"
        if suffix in ("md", "txt"):
            path.write_text(best_text + "\n", encoding="utf-8")
        elif suffix == "json":
            payload_json = {
                "version": __version__,
                "schema_version": CUSTOMER_BEST_SCHEMA_VERSION,
                "type": "customer_best_text",
                "parser": winner_parser,
                "source_file": payload.get("source_file"),
                "vote_score": winner.get("vote_score") if winner else None,
                "quality_score": winner.get("quality_score") if winner else None,
                "text": best_text,
            }
            write_json_file(path, payload_json)
        written.append(path)
    return written


def customer_value(value: object) -> str:
    """Format structured values for customer-facing text."""
    if value is None:
        return ""
    if isinstance(value, float):
        return f"{value:.2f}"
    return str(value).replace("|", "\\|").strip()


def customer_bbox(value: object) -> str:
    """Format a bbox for customer-facing traceability tables."""
    if not isinstance(value, list) or len(value) < 4:
        return ""
    try:
        return "[" + ", ".join(f"{float(item):.2f}" for item in value[:4]) + "]"
    except (TypeError, ValueError):
        return ""


INVOICE_CONFIDENCE_FIELDS = [
    "invoice_type",
    "invoice_number",
    "invoice_date",
    "buyer_name",
    "buyer_tax_id",
    "seller_name",
    "seller_tax_id",
    "total_amount",
    "total_tax",
    "total_with_tax",
    "total_with_tax_cn",
    "drawer",
]


def normalize_field_value(value: object) -> str:
    """Normalize structured field values for cross-parser comparison."""
    if value is None:
        return ""
    if isinstance(value, float):
        return f"{value:.2f}"
    return re.sub(r"\s+", "", str(value))


def estimate_field_page(text: str, value: object) -> int | None:
    """Estimate the page where a field value appears from page markers."""
    value_text = normalize_field_value(value)
    if not value_text:
        return None
    compact = compact_text(text)
    pos = compact.find(value_text)
    if pos < 0:
        return None
    prefix = text[: min(len(text), pos + 200)]
    markers = re.findall(r"<!--\s*page\s+(\d+)", prefix, flags=re.IGNORECASE)
    if markers:
        try:
            return int(markers[-1])
        except ValueError:
            return None
    return 1


def coerce_bbox(value: object) -> list[float] | None:
    """Return a normalized bbox list when a layout object has usable coordinates."""
    if not isinstance(value, (list, tuple)) or len(value) < 4:
        return None
    try:
        return [round(float(item), 3) for item in value[:4]]
    except (TypeError, ValueError):
        return None


def bbox_union(bboxes: list[list[float]]) -> list[float] | None:
    """Return the union of multiple PDF layout bboxes."""
    usable = [bbox for bbox in bboxes if bbox and len(bbox) >= 4]
    if not usable:
        return None
    return [
        round(min(bbox[0] for bbox in usable), 3),
        round(min(bbox[1] for bbox in usable), 3),
        round(max(bbox[2] for bbox in usable), 3),
        round(max(bbox[3] for bbox in usable), 3),
    ]


def layout_field_candidates(layout: dict[str, object] | None) -> list[dict[str, object]]:
    """Flatten layout JSON into searchable line/span candidates."""
    if not layout:
        return []
    pages = layout.get("pages") if isinstance(layout.get("pages"), list) else []
    candidates: list[dict[str, object]] = []
    for page in pages:
        if not isinstance(page, dict):
            continue
        page_no = page.get("page")
        blocks = page.get("blocks") if isinstance(page.get("blocks"), list) else []
        for block in blocks:
            if not isinstance(block, dict):
                continue
            block_id = block.get("block_id")
            lines = block.get("lines") if isinstance(block.get("lines"), list) else []
            for line in lines:
                if not isinstance(line, dict):
                    continue
                line_text = str(line.get("text") or "")
                spans = line.get("spans") if isinstance(line.get("spans"), list) else []
                span_bboxes: list[list[float]] = []
                for span in spans:
                    if not isinstance(span, dict):
                        continue
                    span_text = str(span.get("text") or "")
                    span_bbox = coerce_bbox(span.get("bbox"))
                    if span_bbox:
                        span_bboxes.append(span_bbox)
                    if span_text.strip() and span_bbox:
                        candidates.append({
                            "kind": "span",
                            "page": page_no,
                            "block_id": block_id,
                            "line_id": line.get("line_id"),
                            "span_id": span.get("span_id"),
                            "text": span_text,
                            "context_text": line_text,
                            "bbox": span_bbox,
                        })
                line_bbox = coerce_bbox(line.get("bbox")) or bbox_union(span_bboxes)
                if line_text.strip() and line_bbox:
                    candidates.append({
                        "kind": "line",
                        "page": page_no,
                        "block_id": block_id,
                        "line_id": line.get("line_id"),
                        "span_id": None,
                        "text": line_text,
                        "context_text": line_text,
                        "bbox": line_bbox,
                    })
    return candidates


FIELD_CONTEXT_KEYWORDS = {
    "invoice_type": ["发票"],
    "invoice_number": ["发票号码", "号码"],
    "invoice_date": ["开票日期", "日期"],
    "buyer_name": ["购买方", "购买方信息", "名称"],
    "buyer_tax_id": ["购买方", "纳税人识别号", "统一社会信用代码"],
    "seller_name": ["销售方", "销售方信息", "名称"],
    "seller_tax_id": ["销售方", "纳税人识别号", "统一社会信用代码"],
    "total_amount": ["合计", "金额"],
    "total_tax": ["合计", "税额"],
    "total_with_tax": ["价税合计", "小写"],
    "total_with_tax_cn": ["价税合计", "大写"],
    "drawer": ["开票人"],
}


def field_value_variants(value: object) -> list[str]:
    """Build normalized variants for matching structured values back to layout text."""
    base = normalize_field_value(value)
    variants = {base} if base else set()
    if base:
        variants.add(base.replace(",", ""))
        variants.add(base.replace("¥", "").replace("￥", "").replace(",", ""))
    number = decimal_from_text(str(value)) if value is not None else None
    if number is not None:
        formatted = f"{number:.2f}"
        variants.update({formatted, f"¥{formatted}", f"￥{formatted}"})
    return sorted((item for item in variants if item), key=len, reverse=True)


def find_field_in_layout(
    layout: dict[str, object] | None,
    value: object,
    field: str | None = None,
    candidates: list[dict[str, object]] | None = None,
) -> dict[str, object] | None:
    """Find a conservative field value match in PDF layout coordinates."""
    variants = field_value_variants(value)
    if not variants:
        return None
    candidates = candidates if candidates is not None else layout_field_candidates(layout)
    keywords = FIELD_CONTEXT_KEYWORDS.get(field or "", [])
    best: dict[str, object] | None = None
    best_score = -1
    for candidate in candidates:
        text = str(candidate.get("text") or "")
        context = str(candidate.get("context_text") or text)
        compact_candidate = normalize_field_value(text)
        if not compact_candidate:
            continue
        compact_context = normalize_field_value(context)
        context_bonus = 40 if keywords and any(keyword in context for keyword in keywords) else 0
        for variant in variants:
            exact = compact_candidate == variant
            contextual_short = bool(field and context_bonus and len(variant) >= 2)
            contains = variant in compact_candidate and (len(variant) >= 4 or contextual_short)
            if not exact and not contains:
                continue
            score = (100 if exact else 70) + context_bonus
            if candidate.get("kind") == "line":
                score -= 5
            if score <= best_score:
                continue
            bbox = coerce_bbox(candidate.get("bbox"))
            if not bbox:
                continue
            best_score = score
            best = {
                "page": candidate.get("page"),
                "bbox": bbox,
                "match_type": f"{candidate.get('kind')}_{'exact' if exact else 'contains'}",
                "text": text[:180],
                "block_id": candidate.get("block_id"),
                "line_id": candidate.get("line_id"),
                "span_id": candidate.get("span_id"),
            }
    return best


def build_field_confidence(
    payload: dict[str, object],
    normalized_texts: dict[str, str],
    layout: dict[str, object] | None = None,
    final_fields: dict[str, object] | None = None,
    final_field_decisions: dict[str, object] | None = None,
) -> dict[str, object]:
    """Build field-level confidence from invoice vote candidates."""
    winner = payload.get("winner") if isinstance(payload.get("winner"), dict) else {}
    winner_parser = str(winner.get("parser") or "")
    votes = payload.get("votes") if isinstance(payload.get("votes"), list) else []
    parser_fields: dict[str, dict[str, object]] = {}
    for vote in votes:
        if not isinstance(vote, dict) or not vote.get("eligible"):
            continue
        fields = vote.get("invoice_fields") if isinstance(vote.get("invoice_fields"), dict) else None
        if fields:
            parser_fields[str(vote.get("parser"))] = fields
    parser_count = max(1, len(parser_fields))
    winner_fields = parser_fields.get(winner_parser, {})
    candidates = layout_field_candidates(layout)
    confidence: dict[str, object] = {}
    for field in INVOICE_CONFIDENCE_FIELDS:
        decision = final_field_decisions.get(field) if isinstance(final_field_decisions, dict) and isinstance(final_field_decisions.get(field), dict) else {}
        winning_value = final_fields.get(field) if isinstance(final_fields, dict) and field in final_fields else winner_fields.get(field)
        source_parser = str(decision.get("source_parser") or winner_parser or "")
        normalized_winning = normalize_field_value(winning_value)
        winner_location = find_field_in_layout(layout, winning_value, field, candidates)
        sources: list[dict[str, object]] = []
        support = 0
        for parser_name, fields in parser_fields.items():
            value = fields.get(field)
            normalized_value = normalize_field_value(value)
            if not normalized_value:
                continue
            same = normalized_winning and normalized_value == normalized_winning
            if same:
                support += 1
            text = normalized_texts.get(parser_name, "")
            source_location = find_field_in_layout(layout, value, field, candidates)
            sources.append({
                "parser": parser_name,
                "value": value,
                "matches_final": bool(same),
                "matches_winner": bool(same),
                "page": source_location.get("page") if source_location else estimate_field_page(text, value),
                "bbox": source_location.get("bbox") if source_location else None,
                "location": source_location,
                "confidence": round(1.0 if same else 0.55, 4),
            })
        base_confidence = support / parser_count if normalized_winning else 0.0
        if normalized_winning and source_parser in parser_fields:
            base_confidence = min(1.0, base_confidence + 0.12)
        if decision.get("confidence") is not None:
            try:
                base_confidence = max(base_confidence, float(decision.get("confidence") or 0.0))
            except (TypeError, ValueError):
                pass
        confidence[field] = {
            "value": winning_value,
            "source_parser": source_parser if normalized_winning else None,
            "support_count": support,
            "parser_count": parser_count,
            "confidence": round(base_confidence, 4),
            "page": winner_location.get("page") if winner_location else estimate_field_page(normalized_texts.get(source_parser, ""), winning_value),
            "bbox": winner_location.get("bbox") if winner_location else None,
            "location": winner_location,
            "sources": sources,
        }
    return confidence


def invoice_validation_rank(status: object) -> int:
    """Rank invoice validation statuses for conservative fusion decisions."""
    return {
        "ok": 3,
        "missing_fields": 2,
        "failed_checks": 1,
        "unknown": 0,
    }.get(str(status or "unknown"), 0)


def invoice_value_quality(field: str, value: object) -> float:
    """Score whether an individual invoice value looks plausible."""
    normalized = normalize_field_value(value)
    if not normalized:
        return 0.0
    if field == "invoice_number":
        return 1.0 if re.fullmatch(r"[0-9]{8,30}", normalized) else 0.35
    if field in {"buyer_tax_id", "seller_tax_id"}:
        return 1.0 if re.fullmatch(r"[0-9A-Z]{15,25}", normalized) else 0.45
    if field == "invoice_date":
        return 1.0 if re.fullmatch(r"[0-9]{4}年[0-9]{1,2}月[0-9]{1,2}日|[0-9]{4}[-/][0-9]{1,2}[-/][0-9]{1,2}", str(value)) else 0.55
    if field in {"total_amount", "total_tax", "total_with_tax"}:
        return 1.0 if decimal_from_text(str(value)) is not None else 0.2
    if field in {"buyer_name", "seller_name"}:
        return 1.0 if len(normalized) >= 4 else 0.45
    if field == "total_with_tax_cn":
        return 1.0 if any(ch in normalized for ch in "壹贰叁肆伍陆柒捌玖拾佰仟万圆元整") else 0.55
    return 0.85


def parser_vote_lookup(payload: dict[str, object]) -> dict[str, dict[str, object]]:
    """Return eligible vote rows keyed by parser name."""
    votes = payload.get("votes") if isinstance(payload.get("votes"), list) else []
    return {
        str(vote.get("parser")): vote
        for vote in votes
        if isinstance(vote, dict) and vote.get("eligible")
    }


def choose_fused_invoice_field(
    field: str,
    parser_fields: dict[str, dict[str, object]],
    vote_rows: dict[str, dict[str, object]],
    winner_parser: str,
) -> dict[str, object]:
    """Choose the best value for one invoice field from all parser candidates."""
    groups: dict[str, dict[str, object]] = {}
    for parser_name, fields in parser_fields.items():
        value = fields.get(field)
        normalized = normalize_field_value(value)
        if not normalized:
            continue
        vote = vote_rows.get(parser_name, {})
        group = groups.setdefault(
            normalized,
            {
                "value": value,
                "normalized": normalized,
                "parsers": [],
                "support_count": 0,
                "vote_score_sum": 0.0,
                "invoice_score_sum": 0.0,
                "quality_score_sum": 0.0,
                "winner_supported": False,
                "value_quality": invoice_value_quality(field, value),
            },
        )
        group["parsers"].append(parser_name)
        group["support_count"] = int(group.get("support_count") or 0) + 1
        group["vote_score_sum"] = float(group.get("vote_score_sum") or 0.0) + float(vote.get("vote_score") or 0.0)
        group["invoice_score_sum"] = float(group.get("invoice_score_sum") or 0.0) + float(vote.get("invoice_score") or 0.0)
        group["quality_score_sum"] = float(group.get("quality_score_sum") or 0.0) + float(vote.get("quality_score") or 0.0)
        if parser_name == winner_parser:
            group["winner_supported"] = True

    if not groups:
        return {
            "field": field,
            "value": None,
            "source_parser": None,
            "support_count": 0,
            "parser_count": len(parser_fields),
            "confidence": 0.0,
            "fusion_score": 0.0,
            "selection_reason": "missing_from_all_parsers",
            "supporting_parsers": [],
        }

    parser_count = max(1, len(parser_fields))
    for group in groups.values():
        support_ratio = int(group.get("support_count") or 0) / parser_count
        avg_vote = float(group.get("vote_score_sum") or 0.0) / max(1, int(group.get("support_count") or 0))
        avg_invoice = float(group.get("invoice_score_sum") or 0.0) / max(1, int(group.get("support_count") or 0))
        avg_quality = float(group.get("quality_score_sum") or 0.0) / max(1, int(group.get("support_count") or 0))
        winner_bonus = 0.04 if group.get("winner_supported") else 0.0
        value_quality = float(group.get("value_quality") or 0.0)
        fusion_score = (
            support_ratio * 0.42
            + avg_vote * 0.22
            + avg_invoice * 0.16
            + avg_quality * 0.10
            + value_quality * 0.10
            + winner_bonus
        )
        group["support_ratio"] = round(support_ratio, 4)
        group["avg_vote_score"] = round(avg_vote, 4)
        group["avg_invoice_score"] = round(avg_invoice, 4)
        group["avg_quality_score"] = round(avg_quality, 4)
        group["fusion_score"] = round(min(1.0, fusion_score), 4)
        group["confidence"] = round(min(1.0, fusion_score), 4)

    ranked = sorted(
        groups.values(),
        key=lambda group: (
            float(group.get("fusion_score") or 0.0),
            int(group.get("support_count") or 0),
            bool(group.get("winner_supported")),
            float(group.get("avg_invoice_score") or 0.0),
            float(group.get("avg_vote_score") or 0.0),
        ),
        reverse=True,
    )
    selected = ranked[0]
    support = [str(parser) for parser in selected.get("parsers", []) if parser]
    source_parser = winner_parser if selected.get("winner_supported") else support[0] if support else None
    return {
        "field": field,
        "value": selected.get("value"),
        "source_parser": source_parser,
        "support_count": selected.get("support_count"),
        "parser_count": parser_count,
        "confidence": selected.get("confidence"),
        "fusion_score": selected.get("fusion_score"),
        "selection_reason": "field_consensus_vote",
        "supporting_parsers": support,
        "winner_supported": bool(selected.get("winner_supported")),
        "candidates": [
            {
                "value": group.get("value"),
                "support_count": group.get("support_count"),
                "supporting_parsers": group.get("parsers"),
                "winner_supported": bool(group.get("winner_supported")),
                "fusion_score": group.get("fusion_score"),
                "value_quality": group.get("value_quality"),
            }
            for group in ranked
        ],
    }


def choose_fused_invoice_items(
    parser_fields: dict[str, dict[str, object]],
    vote_rows: dict[str, dict[str, object]],
    winner_parser: str,
) -> dict[str, object]:
    """Choose line items from the strongest parser candidate."""
    best: dict[str, object] | None = None
    for parser_name, fields in parser_fields.items():
        items = fields.get("items") if isinstance(fields.get("items"), list) else []
        if not items:
            continue
        vote = vote_rows.get(parser_name, {})
        repetition = invoice_item_repetition_metrics(items)
        score = (
            min(1.0, len(items) / 3) * 0.25
            + float(vote.get("invoice_score") or 0.0) * 0.35
            + float(vote.get("vote_score") or 0.0) * 0.25
            + float(repetition.get("score") or 0.0) * 0.15
            + (0.08 if parser_name == winner_parser else 0.0)
        )
        candidate = {
            "parser": parser_name,
            "items": items,
            "item_count": len(items),
            "score": round(min(1.0, score), 4),
            "repetition": repetition,
        }
        if best is None or (
            float(candidate["score"]),
            int(candidate["item_count"]),
            parser_name == winner_parser,
        ) > (
            float(best.get("score") or 0.0),
            int(best.get("item_count") or 0),
            str(best.get("parser")) == winner_parser,
        ):
            best = candidate
    if best:
        return best
    return {
        "parser": None,
        "items": [],
        "item_count": 0,
        "score": 0.0,
        "repetition": invoice_item_repetition_metrics([]),
    }


def invoice_fields_snapshot(fields: dict[str, object]) -> dict[str, object]:
    """Return only customer-facing invoice fields from a larger parser payload."""
    return {
        "invoice_type": fields.get("invoice_type"),
        "invoice_number": fields.get("invoice_number"),
        "invoice_date": fields.get("invoice_date"),
        "buyer_name": fields.get("buyer_name"),
        "buyer_tax_id": fields.get("buyer_tax_id"),
        "seller_name": fields.get("seller_name"),
        "seller_tax_id": fields.get("seller_tax_id"),
        "total_amount": fields.get("total_amount"),
        "total_tax": fields.get("total_tax"),
        "total_with_tax": fields.get("total_with_tax"),
        "total_with_tax_cn": fields.get("total_with_tax_cn"),
        "drawer": fields.get("drawer"),
        "items": fields.get("items") if isinstance(fields.get("items"), list) else [],
    }


def build_fused_invoice_fields(payload: dict[str, object]) -> dict[str, object] | None:
    """Fuse invoice fields across eligible parser outputs and re-run validation."""
    winner = payload.get("winner") if isinstance(payload.get("winner"), dict) else {}
    winner_parser = str(winner.get("parser") or "")
    votes = payload.get("votes") if isinstance(payload.get("votes"), list) else []
    parser_fields: dict[str, dict[str, object]] = {}
    for vote in votes:
        if not isinstance(vote, dict) or not vote.get("eligible"):
            continue
        fields = vote.get("invoice_fields") if isinstance(vote.get("invoice_fields"), dict) else None
        if fields:
            parser_fields[str(vote.get("parser"))] = fields
    if not parser_fields:
        return None

    vote_rows = parser_vote_lookup(payload)
    winner_fields = parser_fields.get(winner_parser, {})
    fused: dict[str, object] = {
        "version": __version__,
        "profile": "invoice",
        "source_file": payload.get("source_file"),
        "parser": "field-fusion",
        "base_parser": winner_parser,
        "quality_score": winner.get("quality_score"),
        "quality_label": winner.get("quality_label"),
        "raw_text_chars": winner_fields.get("raw_text_chars"),
    }
    field_decisions: dict[str, object] = {}
    changed_fields: list[str] = []
    for field in INVOICE_CONFIDENCE_FIELDS:
        decision = choose_fused_invoice_field(field, parser_fields, vote_rows, winner_parser)
        fused[field] = decision.get("value")
        field_decisions[field] = decision
        if normalize_field_value(decision.get("value")) != normalize_field_value(winner_fields.get(field)):
            changed_fields.append(field)

    item_decision = choose_fused_invoice_items(parser_fields, vote_rows, winner_parser)
    fused["items"] = item_decision.get("items") if isinstance(item_decision.get("items"), list) else []
    if item_decision.get("parser") and item_decision.get("parser") != winner_parser:
        changed_fields.append("items")

    fused_validation = validate_invoice_fields(fused, strict=True)
    fused["validation"] = fused_validation
    winner_validation = winner_fields.get("validation") if isinstance(winner_fields.get("validation"), dict) else {}
    use_fused = True
    if invoice_validation_rank(fused_validation.get("status")) < invoice_validation_rank(winner_validation.get("status")):
        use_fused = False

    return {
        "enabled": True,
        "used": use_fused,
        "strategy": "field_consensus_vote",
        "base_parser": winner_parser,
        "parser_count": len(parser_fields),
        "changed_fields": sorted(set(changed_fields)),
        "field_decisions": field_decisions,
        "item_decision": item_decision,
        "candidate_fields": fused,
        "candidate_validation": fused_validation,
        "winner_validation": winner_validation,
        "reason": "fused_fields_passed_validation_rank" if use_fused else "fused_validation_rank_worse_than_winner",
    }


def build_customer_invoice_delivery(
    payload: dict[str, object],
    normalized_texts: dict[str, str] | None = None,
    layout: dict[str, object] | None = None,
    field_fusion: bool = True,
) -> dict[str, object] | None:
    """Build a clean structured invoice payload from the winning vote."""
    winner = payload.get("winner") if isinstance(payload.get("winner"), dict) else None
    if not winner:
        return None
    fields = winner.get("invoice_fields") if isinstance(winner.get("invoice_fields"), dict) else None
    if not fields:
        return None
    normalized_texts = normalized_texts or {}
    fusion = build_fused_invoice_fields(payload) if field_fusion else None
    if fusion and fusion.get("used") and isinstance(fusion.get("candidate_fields"), dict):
        fields = fusion["candidate_fields"]
        parser_name = "field-fusion"
    else:
        parser_name = winner.get("parser")
    fusion_decisions = fusion.get("field_decisions") if isinstance(fusion, dict) and isinstance(fusion.get("field_decisions"), dict) else None
    return {
        "version": __version__,
        "schema_version": CUSTOMER_BEST_SCHEMA_VERSION,
        "type": "customer_invoice_delivery",
        "generated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        "source_file": payload.get("source_file"),
        "parser": parser_name,
        "base_parser": winner.get("parser"),
        "vote_score": winner.get("vote_score"),
        "invoice_score": winner.get("invoice_score"),
        "quality_score": winner.get("quality_score"),
        "validation": fields.get("validation"),
        "field_confidence": build_field_confidence(payload, normalized_texts, layout, fields, fusion_decisions),
        "invoice": {
            "invoice_type": fields.get("invoice_type"),
            "invoice_number": fields.get("invoice_number"),
            "invoice_date": fields.get("invoice_date"),
            "buyer_name": fields.get("buyer_name"),
            "buyer_tax_id": fields.get("buyer_tax_id"),
            "seller_name": fields.get("seller_name"),
            "seller_tax_id": fields.get("seller_tax_id"),
            "total_amount": fields.get("total_amount"),
            "total_tax": fields.get("total_tax"),
            "total_with_tax": fields.get("total_with_tax"),
            "total_with_tax_cn": fields.get("total_with_tax_cn"),
            "drawer": fields.get("drawer"),
            "items": fields.get("items") if isinstance(fields.get("items"), list) else [],
        },
        "audit": {
            "profile": payload.get("profile"),
            "weights": payload.get("weights"),
            "field_fusion": fusion or {"enabled": False, "used": False},
            "winner": {
                key: winner.get(key)
                for key in [
                    "parser",
                    "vote_score",
                    "quality_score",
                    "consensus_score",
                    "coverage_score",
                    "structure_score",
                    "repetition_score",
                    "invoice_score",
                    "penalty",
                ]
            },
        },
    }


def build_structured_field_confidence(
    payload: dict[str, object],
    normalized_texts: dict[str, str],
    layout: dict[str, object] | None,
    fields: dict[str, object],
) -> dict[str, object]:
    """Build field confidence for lightweight business structured fields."""
    winner = payload.get("winner") if isinstance(payload.get("winner"), dict) else {}
    winner_parser = str(winner.get("parser") or "")
    votes = payload.get("votes") if isinstance(payload.get("votes"), list) else []
    parser_fields: dict[str, dict[str, object]] = {}
    for vote in votes:
        if not isinstance(vote, dict) or not vote.get("eligible"):
            continue
        row_fields = vote.get("structured_fields") if isinstance(vote.get("structured_fields"), dict) else None
        if row_fields:
            parser_fields[str(vote.get("parser"))] = row_fields
    parser_count = max(1, len(parser_fields))
    candidates = layout_field_candidates(layout)
    ignored = {"version", "schema_version", "profile", "source_file", "parser", "quality_score", "quality_label", "raw_text_chars", "validation"}
    confidence: dict[str, object] = {}
    for field, value in fields.items():
        if field in ignored:
            continue
        if isinstance(value, (dict, list)):
            if not value:
                continue
            value_for_match = value[0] if isinstance(value, list) and value else None
        else:
            value_for_match = value
        normalized_value = normalize_field_value(value_for_match)
        if not normalized_value:
            continue
        support = 0
        sources: list[dict[str, object]] = []
        for parser_name, parser_data in parser_fields.items():
            candidate_value = parser_data.get(field)
            if isinstance(candidate_value, list):
                candidate_match_value = candidate_value[0] if candidate_value else None
            else:
                candidate_match_value = candidate_value
            same = normalize_field_value(candidate_match_value) == normalized_value
            if same:
                support += 1
            location = find_field_in_layout(layout, candidate_match_value, field, candidates)
            sources.append({
                "parser": parser_name,
                "value": candidate_value,
                "matches_final": bool(same),
                "page": location.get("page") if location else estimate_field_page(normalized_texts.get(parser_name, ""), candidate_match_value),
                "bbox": location.get("bbox") if location else None,
                "location": location,
            })
        location = find_field_in_layout(layout, value_for_match, field, candidates)
        confidence[field] = {
            "value": value,
            "source_parser": winner_parser,
            "support_count": support,
            "parser_count": parser_count,
            "confidence": round(min(1.0, support / parser_count + 0.12), 4),
            "page": location.get("page") if location else estimate_field_page(normalized_texts.get(winner_parser, ""), value_for_match),
            "bbox": location.get("bbox") if location else None,
            "location": location,
            "sources": sources,
        }
    return confidence


def build_customer_structured_delivery(
    payload: dict[str, object],
    normalized_texts: dict[str, str],
    layout: dict[str, object] | None = None,
) -> dict[str, object] | None:
    """Build customer delivery for lightweight non-invoice business profiles."""
    profile = payload.get("profile") if isinstance(payload.get("profile"), dict) else {}
    resolved = str(profile.get("resolved") or "")
    if resolved not in BUSINESS_STRUCTURED_PROFILES:
        return None
    winner = payload.get("winner") if isinstance(payload.get("winner"), dict) else None
    if not winner:
        return None
    fields = winner.get("structured_fields") if isinstance(winner.get("structured_fields"), dict) else None
    if not fields:
        return None
    return {
        "version": __version__,
        "schema_version": CUSTOMER_BEST_SCHEMA_VERSION,
        "type": "customer_structured_delivery",
        "generated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        "source_file": payload.get("source_file"),
        "profile": resolved,
        "parser": winner.get("parser"),
        "vote_score": winner.get("vote_score"),
        "business_score": winner.get("business_score"),
        "quality_score": winner.get("quality_score"),
        "validation": fields.get("validation"),
        "fields": {
            key: value
            for key, value in fields.items()
            if key not in {"version", "schema_version", "profile", "source_file", "parser", "quality_score", "quality_label", "raw_text_chars", "validation"}
        },
        "field_confidence": build_structured_field_confidence(payload, normalized_texts, layout, fields),
        "audit": {
            "profile": profile,
            "weights": payload.get("weights"),
            "winner": {
                key: winner.get(key)
                for key in [
                    "parser",
                    "vote_score",
                    "quality_score",
                    "consensus_score",
                    "coverage_score",
                    "structure_score",
                    "repetition_score",
                    "business_score",
                    "penalty",
                ]
            },
        },
    }


def customer_structured_to_markdown(delivery: dict[str, object]) -> str:
    """Render lightweight business structured delivery as Markdown."""
    fields = delivery.get("fields") if isinstance(delivery.get("fields"), dict) else {}
    validation = delivery.get("validation") if isinstance(delivery.get("validation"), dict) else {}
    lines = [
        f"# {delivery.get('profile')} 结构化解析结果",
        "",
        f"- 来源文件：`{delivery.get('source_file')}`",
        f"- 推荐解析器：`{delivery.get('parser')}`",
        f"- 投票得分：`{delivery.get('vote_score')}`",
        f"- 业务字段得分：`{delivery.get('business_score')}`",
        f"- 字段状态：`{validation.get('status')}`",
        "",
        "## 字段",
        "",
        "| 字段 | 值 |",
        "|---|---|",
    ]
    for key, value in fields.items():
        if isinstance(value, list):
            rendered = "<br>".join(customer_value(item) for item in value)
        else:
            rendered = customer_value(value)
        lines.append(f"| `{key}` | {rendered} |")
    confidence = delivery.get("field_confidence") if isinstance(delivery.get("field_confidence"), dict) else {}
    if confidence:
        lines.extend(["", "## 字段置信度", "", "| 字段 | 来源解析器 | 置信度 | 支持数 | 页码 | 坐标 |", "|---|---|---:|---:|---:|---|"])
        for key, info in confidence.items():
            if not isinstance(info, dict):
                continue
            lines.append(
                f"| `{key}` | {customer_value(info.get('source_parser'))} | {customer_value(info.get('confidence'))} | "
                f"{customer_value(info.get('support_count'))}/{customer_value(info.get('parser_count'))} | "
                f"{customer_value(info.get('page'))} | {customer_bbox(info.get('bbox'))} |"
            )
    return "\n".join(lines) + "\n"


def customer_structured_to_txt(delivery: dict[str, object]) -> str:
    """Render lightweight business structured delivery as plain text."""
    fields = delivery.get("fields") if isinstance(delivery.get("fields"), dict) else {}
    validation = delivery.get("validation") if isinstance(delivery.get("validation"), dict) else {}
    lines = [
        f"{delivery.get('profile')} 结构化解析结果",
        f"来源文件: {delivery.get('source_file')}",
        f"推荐解析器: {delivery.get('parser')}",
        f"投票得分: {delivery.get('vote_score')}",
        f"业务字段得分: {delivery.get('business_score')}",
        f"字段状态: {validation.get('status')}",
        "",
        "字段",
    ]
    for key, value in fields.items():
        if isinstance(value, list):
            rendered = "; ".join(customer_value(item) for item in value)
        else:
            rendered = customer_value(value)
        lines.append(f"{key}: {rendered}")
    return "\n".join(lines) + "\n"


def customer_invoice_to_markdown(delivery: dict[str, object]) -> str:
    """Render customer invoice delivery as Markdown."""
    invoice = delivery.get("invoice") if isinstance(delivery.get("invoice"), dict) else {}
    validation = delivery.get("validation") if isinstance(delivery.get("validation"), dict) else {}
    audit = delivery.get("audit") if isinstance(delivery.get("audit"), dict) else {}
    fusion = audit.get("field_fusion") if isinstance(audit.get("field_fusion"), dict) else {}
    lines = [
        "# 发票解析结果",
        "",
        f"- 来源文件：`{delivery.get('source_file')}`",
        f"- 推荐解析器：`{delivery.get('parser')}`",
        f"- 基础胜出解析器：`{delivery.get('base_parser')}`",
        f"- 投票得分：`{delivery.get('vote_score')}`",
        f"- 发票校验：`{validation.get('status')}`",
        f"- 字段级融合：`{fusion.get('used')}`",
        "",
        "## 基本信息",
        "",
        "| 字段 | 值 |",
        "|---|---|",
    ]
    labels = [
        ("invoice_type", "发票类型"),
        ("invoice_number", "发票号码"),
        ("invoice_date", "开票日期"),
        ("buyer_name", "购买方名称"),
        ("buyer_tax_id", "购买方税号"),
        ("seller_name", "销售方名称"),
        ("seller_tax_id", "销售方税号"),
        ("total_amount", "合计金额"),
        ("total_tax", "合计税额"),
        ("total_with_tax", "价税合计"),
        ("total_with_tax_cn", "价税合计大写"),
        ("drawer", "开票人"),
    ]
    for key, label in labels:
        lines.append(f"| {label} | {customer_value(invoice.get(key))} |")

    field_confidence = delivery.get("field_confidence") if isinstance(delivery.get("field_confidence"), dict) else {}
    if field_confidence:
        lines.extend(["", "## 字段置信度", "", "| 字段 | 来源解析器 | 置信度 | 支持数 | 页码 | 坐标 |", "|---|---|---:|---:|---:|---|"])
        for key, label in labels:
            info = field_confidence.get(key) if isinstance(field_confidence.get(key), dict) else {}
            fusion_decisions = fusion.get("field_decisions") if isinstance(fusion.get("field_decisions"), dict) else {}
            decision = fusion_decisions.get(key) if isinstance(fusion_decisions.get(key), dict) else {}
            source_parser = decision.get("source_parser") if fusion.get("used") and decision else info.get("source_parser")
            confidence = decision.get("confidence") if fusion.get("used") and decision else info.get("confidence")
            support = decision.get("support_count") if fusion.get("used") and decision else info.get("support_count")
            parser_count = decision.get("parser_count") if fusion.get("used") and decision else info.get("parser_count")
            lines.append(
                f"| {label} | {customer_value(source_parser)} | "
                f"{customer_value(confidence)} | {customer_value(support)}/{customer_value(parser_count)} | "
                f"{customer_value(info.get('page'))} | {customer_bbox(info.get('bbox'))} |"
            )

    if fusion.get("used"):
        changed = fusion.get("changed_fields") if isinstance(fusion.get("changed_fields"), list) else []
        lines.extend(["", "## 字段级融合", "", f"- 策略：`{fusion.get('strategy')}`", f"- 基础解析器：`{fusion.get('base_parser')}`", f"- 变更字段：`{', '.join(str(item) for item in changed) if changed else '无'}`"])

    lines.extend(["", "## 明细", ""])
    items = invoice.get("items") if isinstance(invoice.get("items"), list) else []
    if items:
        lines.extend([
            "| 项目名称 | 单位 | 数量 | 单价 | 金额 | 税率 | 税额 |",
            "|---|---|---:|---:|---:|---:|---:|",
        ])
        for item in items:
            if not isinstance(item, dict):
                continue
            lines.append(
                f"| {customer_value(item.get('name'))} | {customer_value(item.get('unit'))} | "
                f"{customer_value(item.get('quantity'))} | {customer_value(item.get('unit_price'))} | "
                f"{customer_value(item.get('amount'))} | {customer_value(item.get('tax_rate'))} | "
                f"{customer_value(item.get('tax_amount'))} |"
            )
    else:
        lines.append("未抽取到明细。")

    checks = validation.get("checks") if isinstance(validation.get("checks"), list) else []
    if checks:
        lines.extend(["", "## 校验", "", "| 检查项 | 结果 | 期望 | 实际 |", "|---|---:|---|---|"])
        for check in checks:
            if not isinstance(check, dict):
                continue
            lines.append(
                f"| {customer_value(check.get('name'))} | {customer_value(check.get('passed'))} | "
                f"{customer_value(check.get('expected'))} | {customer_value(check.get('actual'))} |"
            )
    return "\n".join(lines) + "\n"


def customer_invoice_to_txt(delivery: dict[str, object]) -> str:
    """Render customer invoice delivery as plain text."""
    invoice = delivery.get("invoice") if isinstance(delivery.get("invoice"), dict) else {}
    validation = delivery.get("validation") if isinstance(delivery.get("validation"), dict) else {}
    audit = delivery.get("audit") if isinstance(delivery.get("audit"), dict) else {}
    fusion = audit.get("field_fusion") if isinstance(audit.get("field_fusion"), dict) else {}
    labels = [
        ("invoice_type", "发票类型"),
        ("invoice_number", "发票号码"),
        ("invoice_date", "开票日期"),
        ("buyer_name", "购买方名称"),
        ("buyer_tax_id", "购买方税号"),
        ("seller_name", "销售方名称"),
        ("seller_tax_id", "销售方税号"),
        ("total_amount", "合计金额"),
        ("total_tax", "合计税额"),
        ("total_with_tax", "价税合计"),
        ("total_with_tax_cn", "价税合计大写"),
        ("drawer", "开票人"),
    ]
    lines = [
        "发票解析结果",
        f"来源文件: {delivery.get('source_file')}",
        f"推荐解析器: {delivery.get('parser')}",
        f"基础胜出解析器: {delivery.get('base_parser')}",
        f"投票得分: {delivery.get('vote_score')}",
        f"发票校验: {validation.get('status')}",
        f"字段级融合: {fusion.get('used')}",
        "",
        "基本信息",
    ]
    for key, label in labels:
        lines.append(f"{label}: {customer_value(invoice.get(key))}")
    field_confidence = delivery.get("field_confidence") if isinstance(delivery.get("field_confidence"), dict) else {}
    if field_confidence:
        lines.extend(["", "字段置信度"])
        for key, label in labels:
            info = field_confidence.get(key) if isinstance(field_confidence.get(key), dict) else {}
            fusion_decisions = fusion.get("field_decisions") if isinstance(fusion.get("field_decisions"), dict) else {}
            decision = fusion_decisions.get(key) if isinstance(fusion_decisions.get(key), dict) else {}
            source_parser = decision.get("source_parser") if fusion.get("used") and decision else info.get("source_parser")
            confidence = decision.get("confidence") if fusion.get("used") and decision else info.get("confidence")
            support = decision.get("support_count") if fusion.get("used") and decision else info.get("support_count")
            parser_count = decision.get("parser_count") if fusion.get("used") and decision else info.get("parser_count")
            lines.append(
                f"{label}: 置信度 {customer_value(confidence)}, "
                f"来源 {customer_value(source_parser)}, "
                f"支持 {customer_value(support)}/{customer_value(parser_count)}, "
                f"页码 {customer_value(info.get('page'))}, "
                f"坐标 {customer_bbox(info.get('bbox'))}"
            )
    if fusion.get("used"):
        changed = fusion.get("changed_fields") if isinstance(fusion.get("changed_fields"), list) else []
        lines.extend(["", "字段级融合", f"策略: {fusion.get('strategy')}", f"基础解析器: {fusion.get('base_parser')}", f"变更字段: {', '.join(str(item) for item in changed) if changed else '无'}"])
    lines.extend(["", "明细"])
    items = invoice.get("items") if isinstance(invoice.get("items"), list) else []
    if items:
        for index, item in enumerate(items, 1):
            if not isinstance(item, dict):
                continue
            lines.append(
                f"{index}. {customer_value(item.get('name'))} | 单位 {customer_value(item.get('unit'))} | "
                f"数量 {customer_value(item.get('quantity'))} | 单价 {customer_value(item.get('unit_price'))} | "
                f"金额 {customer_value(item.get('amount'))} | 税率 {customer_value(item.get('tax_rate'))} | "
                f"税额 {customer_value(item.get('tax_amount'))}"
            )
    else:
        lines.append("未抽取到明细。")
    return "\n".join(lines) + "\n"


def build_customer_text_delivery(
    payload: dict[str, object],
    normalized_texts: dict[str, str],
) -> dict[str, object] | None:
    """Build a generic customer payload when no structured profile is active."""
    winner = payload.get("winner") if isinstance(payload.get("winner"), dict) else None
    if not winner:
        return None
    winner_parser = str(winner.get("parser") or "")
    text = normalized_texts.get(winner_parser, "")
    if not text:
        return None
    return {
        "version": __version__,
        "schema_version": CUSTOMER_BEST_SCHEMA_VERSION,
        "type": "customer_best_text",
        "generated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        "source_file": payload.get("source_file"),
        "parser": winner_parser,
        "vote_score": winner.get("vote_score"),
        "quality_score": winner.get("quality_score"),
        "text": text,
    }


def build_customer_delivery(
    payload: dict[str, object],
    normalized_texts: dict[str, str],
    layout: dict[str, object] | None = None,
    field_fusion: bool = True,
) -> dict[str, object] | None:
    """Build the best customer delivery payload for the resolved profile."""
    profile = payload.get("profile") if isinstance(payload.get("profile"), dict) else {}
    delivery: dict[str, object] | None = None
    if profile.get("resolved") == "invoice":
        delivery = build_customer_invoice_delivery(payload, normalized_texts, layout, field_fusion)
    if delivery is None and profile.get("resolved") in BUSINESS_STRUCTURED_PROFILES:
        delivery = build_customer_structured_delivery(payload, normalized_texts, layout)
    if delivery is None:
        delivery = build_customer_text_delivery(payload, normalized_texts)
    return delivery


def customer_best_json_schema() -> dict[str, object]:
    """Return the customer_best.json delivery contract."""
    return {
        "$schema": "https://json-schema.org/draft/2020-12/schema",
        "$id": "customer_best.schema.json",
        "title": "PDF Skill Customer Best Delivery",
        "type": "object",
        "required": ["version", "schema_version", "type", "generated_at", "source_file", "parser"],
        "properties": {
            "version": {"type": "string"},
            "schema_version": {"const": CUSTOMER_BEST_SCHEMA_VERSION},
            "type": {"enum": sorted(CUSTOMER_TEXT_TYPES)},
            "generated_at": {"type": "string"},
            "source_file": {"type": ["string", "null"]},
            "parser": {"type": ["string", "null"]},
            "base_parser": {"type": ["string", "null"]},
            "profile": {"type": ["string", "null"]},
            "vote_score": {"type": ["number", "null"]},
            "quality_score": {"type": ["number", "null"]},
            "invoice_score": {"type": ["number", "null"]},
            "business_score": {"type": ["number", "null"]},
            "validation": {"type": ["object", "null"]},
            "invoice": {"type": "object"},
            "fields": {"type": "object"},
            "field_confidence": {
                "type": "object",
                "additionalProperties": {
                    "type": "object",
                    "properties": {
                        "value": {},
                        "source_parser": {"type": ["string", "null"]},
                        "support_count": {"type": "integer"},
                        "parser_count": {"type": "integer"},
                        "confidence": {"type": "number"},
                        "page": {"type": ["integer", "null"]},
                        "bbox": {
                            "anyOf": [
                                {"type": "null"},
                                {"type": "array", "items": {"type": "number"}, "minItems": 4, "maxItems": 4},
                            ]
                        },
                        "location": {"type": ["object", "null"]},
                        "evidence_image": {"type": ["string", "null"]},
                        "sources": {"type": "array"},
                    },
                },
            },
            "audit": {"type": "object"},
            "text": {"type": "string"},
        },
        "additionalProperties": True,
    }


def write_customer_schema(out_dir: Path) -> Path:
    """Write the stable customer_best.json schema contract."""
    return write_json_file(out_dir / "customer_best.schema.json", customer_best_json_schema())


def write_customer_outputs(
    payload: dict[str, object],
    normalized_texts: dict[str, str],
    out_dir: Path,
    output_format: str,
    layout: dict[str, object] | None = None,
    field_fusion: bool = True,
) -> list[Path]:
    """Write customer_best.* outputs from structured fields when available."""
    delivery = build_customer_delivery(payload, normalized_texts, layout, field_fusion)
    if delivery is None:
        return []
    return write_customer_delivery_outputs(delivery, out_dir, output_format)


def write_customer_delivery_outputs(
    delivery: dict[str, object],
    out_dir: Path,
    output_format: str,
) -> list[Path]:
    """Write an already-built customer delivery payload."""
    out_dir.mkdir(parents=True, exist_ok=True)
    written: list[Path] = []
    written.append(write_customer_schema(out_dir))
    for suffix in output_suffixes(output_format):
        path = out_dir / f"customer_best.{suffix}"
        if suffix == "json":
            write_json_file(path, delivery)
        elif suffix == "md":
            if delivery.get("type") == "customer_invoice_delivery":
                path.write_text(customer_invoice_to_markdown(delivery), encoding="utf-8")
            elif delivery.get("type") == "customer_structured_delivery":
                path.write_text(customer_structured_to_markdown(delivery), encoding="utf-8")
            else:
                path.write_text(str(delivery.get("text") or "") + "\n", encoding="utf-8")
        elif suffix == "txt":
            if delivery.get("type") == "customer_invoice_delivery":
                path.write_text(customer_invoice_to_txt(delivery), encoding="utf-8")
            elif delivery.get("type") == "customer_structured_delivery":
                path.write_text(customer_structured_to_txt(delivery), encoding="utf-8")
            else:
                path.write_text(str(delivery.get("text") or "") + "\n", encoding="utf-8")
        written.append(path)
    return written


def safe_filename_part(value: object, fallback: str = "item") -> str:
    """Return a filesystem-friendly short name part."""
    text = str(value or fallback)
    text = re.sub(r"[^A-Za-z0-9._\-\u4e00-\u9fff]+", "_", text).strip("._")
    return text[:80] or fallback


def iter_customer_field_confidence(delivery: dict[str, object]) -> list[tuple[str, dict[str, object]]]:
    """Return field confidence rows from a customer delivery."""
    confidence = delivery.get("field_confidence") if isinstance(delivery.get("field_confidence"), dict) else {}
    rows: list[tuple[str, dict[str, object]]] = []
    for field, info in confidence.items():
        if isinstance(info, dict):
            rows.append((str(field), info))
    return rows


def crop_pdf_bbox_image(
    pdf_path: Path,
    page_number: int,
    bbox: list[float],
    out_path: Path,
    dpi: int = 220,
    padding_points: float = 6.0,
) -> dict[str, object]:
    """Crop one PDF bbox to an evidence PNG."""
    if not module_exists("fitz"):
        raise ImportError("字段截图需要 PyMuPDF/fitz，运行: pip install pymupdf")
    import fitz  # type: ignore

    doc = fitz.open(pdf_path)
    try:
        page_index = max(0, min(doc.page_count - 1, page_number - 1))
        page = doc.load_page(page_index)
        rect = fitz.Rect(*bbox)
        rect.x0 -= padding_points
        rect.y0 -= padding_points
        rect.x1 += padding_points
        rect.y1 += padding_points
        rect = rect & page.rect
        zoom = dpi / 72.0
        pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), clip=rect, alpha=False)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        pix.save(out_path)
        return {
            "page": page_number,
            "bbox": [round(float(item), 3) for item in bbox[:4]],
            "file": str(out_path),
            "width": pix.width,
            "height": pix.height,
            "dpi": dpi,
            "padding_points": padding_points,
        }
    finally:
        doc.close()


def write_field_evidence_pack(
    pdf_path: Path,
    delivery: dict[str, object],
    out_dir: Path,
    dpi: int = 220,
) -> dict[str, object]:
    """Create field evidence crops from field_confidence bbox values."""
    evidence_dir = out_dir / "field_evidence"
    rows: list[dict[str, object]] = []
    diagnostics: list[str] = []
    for field, info in iter_customer_field_confidence(delivery):
        page = info.get("page")
        bbox = coerce_bbox(info.get("bbox"))
        if not page or not bbox:
            rows.append({
                "field": field,
                "value": info.get("value"),
                "source_parser": info.get("source_parser"),
                "page": page,
                "bbox": bbox,
                "status": "missing_bbox",
                "file": None,
            })
            continue
        try:
            path = evidence_dir / f"{safe_filename_part(field)}_p{int(page):03d}.png"
            crop = crop_pdf_bbox_image(pdf_path, int(page), bbox, path, dpi=dpi)
            info["evidence_image"] = str(path)
            rows.append({
                "field": field,
                "value": info.get("value"),
                "source_parser": info.get("source_parser"),
                "confidence": info.get("confidence"),
                "status": "ok",
                **crop,
            })
        except Exception as exc:
            diagnostics.append(f"{field}: {exc!r}")
            rows.append({
                "field": field,
                "value": info.get("value"),
                "source_parser": info.get("source_parser"),
                "page": page,
                "bbox": bbox,
                "status": "failed",
                "file": None,
                "error": repr(exc),
            })
    payload = {
        "version": __version__,
        "type": "field-evidence",
        "source_file": str(pdf_path),
        "generated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        "dpi": dpi,
        "evidence_dir": str(evidence_dir),
        "field_count": len(rows),
        "image_count": sum(1 for row in rows if row.get("status") == "ok"),
        "diagnostics": diagnostics,
        "fields": rows,
    }
    evidence_dir.mkdir(parents=True, exist_ok=True)
    md_lines = [
        "# Field Evidence",
        "",
        f"- Source: `{pdf_path}`",
        f"- Image count: `{payload['image_count']}` / `{payload['field_count']}`",
        "",
        "| Field | Value | Parser | Confidence | Page | BBox | Image | Status |",
        "|---|---|---|---:|---:|---|---|---|",
    ]
    for row in rows:
        image = Path(str(row.get("file"))).name if row.get("file") else ""
        md_lines.append(
            f"| `{row.get('field')}` | {customer_value(row.get('value'))} | {customer_value(row.get('source_parser'))} | "
            f"{customer_value(row.get('confidence'))} | {customer_value(row.get('page'))} | {customer_bbox(row.get('bbox'))} | "
            f"{image} | {row.get('status')} |"
        )
    md_path = evidence_dir / "field_evidence.md"
    md_path.write_text("\n".join(md_lines) + "\n", encoding="utf-8")
    json_path = evidence_dir / "field_evidence.json"
    payload["outputs"] = {"json": str(json_path), "md": str(md_path)}
    write_json_file(json_path, payload)
    return payload


def html_rel_path(target: object, base_dir: Path) -> str:
    """Return a safe relative path for local HTML links."""
    if not target:
        return ""
    try:
        return Path(str(target)).resolve().relative_to(base_dir.resolve()).as_posix()
    except Exception:
        try:
            return Path(str(target)).resolve().as_uri()
        except Exception:
            return html.escape(str(target))


def render_review_pages(pdf_path: Path, out_dir: Path, pages: list[int], dpi: int = 120) -> list[dict[str, object]]:
    """Render PDF pages used by the review HTML."""
    page_count = get_pdf_page_count(pdf_path)
    page_indices = sorted({page - 1 for page in pages if isinstance(page, int) and 1 <= page <= page_count})
    if not page_indices:
        page_indices = [0] if page_count else []
    rendered = render_pdf_pages(pdf_path, page_indices, out_dir / "review_pages", dpi=dpi)
    zoom = dpi / 72.0
    for item in rendered:
        try:
            item["point_width"] = round(float(item.get("width") or 0) / zoom, 3)
            item["point_height"] = round(float(item.get("height") or 0) / zoom, 3)
        except (TypeError, ValueError, ZeroDivisionError):
            pass
    return rendered


def write_review_html(
    pdf_path: Path,
    delivery: dict[str, object],
    out_dir: Path,
    table_payload: dict[str, object] | None = None,
    dpi: int = 120,
) -> Path:
    """Write a local HTML review page with page images and field bbox overlays."""
    field_rows = iter_customer_field_confidence(delivery)
    pages = [int(info.get("page")) for _, info in field_rows if info.get("page")]
    rendered_pages = render_review_pages(pdf_path, out_dir, pages, dpi=dpi)
    page_assets = {int(item.get("page") or 0): item for item in rendered_pages}
    fields_for_js: list[dict[str, object]] = []
    for field, info in field_rows:
        bbox = coerce_bbox(info.get("bbox"))
        page_no = int(info.get("page") or 0) if info.get("page") else None
        page_asset = page_assets.get(page_no or 0)
        fields_for_js.append({
            "field": field,
            "value": info.get("value"),
            "source_parser": info.get("source_parser"),
            "confidence": info.get("confidence"),
            "page": page_no,
            "bbox": bbox,
            "image": html_rel_path(page_asset.get("file"), out_dir) if isinstance(page_asset, dict) else "",
            "page_width": page_asset.get("point_width") if isinstance(page_asset, dict) else None,
            "page_height": page_asset.get("point_height") if isinstance(page_asset, dict) else None,
            "evidence_image": html_rel_path(info.get("evidence_image"), out_dir),
        })
    best_tables = table_payload.get("best_tables") if isinstance(table_payload, dict) and isinstance(table_payload.get("best_tables"), list) else []
    data_json = json.dumps({
        "source_file": str(pdf_path),
        "delivery_type": delivery.get("type"),
        "profile": delivery.get("profile") or (delivery.get("audit", {}) if isinstance(delivery.get("audit"), dict) else {}).get("profile"),
        "fields": fields_for_js,
        "tables": best_tables[:20],
    }, ensure_ascii=False)
    safe_data_json = data_json.replace("</", "<\\/")
    html_text = f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>PDF Review</title>
  <style>
    :root {{ color-scheme: light; --line:#d7dee8; --ink:#16202a; --muted:#667085; --accent:#0f766e; --soft:#eaf7f5; }}
    * {{ box-sizing: border-box; }}
    body {{ margin: 0; font-family: Arial, "Microsoft YaHei", sans-serif; color: var(--ink); background: #f6f8fb; }}
    header {{ padding: 14px 18px; border-bottom: 1px solid var(--line); background: #fff; position: sticky; top: 0; z-index: 10; }}
    header h1 {{ margin: 0; font-size: 18px; }}
    .shell {{ display: grid; grid-template-columns: minmax(360px, 54vw) 1fr; min-height: calc(100vh - 52px); }}
    .pdf {{ padding: 16px; overflow: auto; border-right: 1px solid var(--line); }}
    .panel {{ padding: 16px; overflow: auto; background: #fff; }}
    .page-wrap {{ position: relative; display: inline-block; background: #fff; box-shadow: 0 2px 10px rgba(22,32,42,.12); }}
    .page-wrap img {{ display: block; max-width: 100%; height: auto; }}
    .box {{ position: absolute; border: 3px solid var(--accent); background: rgba(15,118,110,.12); pointer-events: none; }}
    .field {{ width: 100%; text-align: left; padding: 10px 12px; border: 1px solid var(--line); border-radius: 6px; background: #fff; margin-bottom: 8px; cursor: pointer; }}
    .field.active {{ border-color: var(--accent); background: var(--soft); }}
    .field strong {{ display: block; font-size: 13px; }}
    .field span {{ display: block; color: var(--muted); font-size: 12px; margin-top: 4px; word-break: break-all; }}
    table {{ border-collapse: collapse; width: 100%; font-size: 12px; margin: 10px 0 18px; }}
    th, td {{ border: 1px solid var(--line); padding: 6px; vertical-align: top; }}
    h2 {{ font-size: 15px; margin: 18px 0 10px; }}
    @media (max-width: 900px) {{ .shell {{ grid-template-columns: 1fr; }} .pdf {{ border-right: 0; border-bottom: 1px solid var(--line); }} }}
  </style>
</head>
<body>
  <header><h1>PDF Review</h1></header>
  <main class="shell">
    <section class="pdf"><div id="page" class="page-wrap"></div></section>
    <section class="panel">
      <h2>Fields</h2>
      <div id="fields"></div>
      <h2>Tables</h2>
      <div id="tables"></div>
    </section>
  </main>
  <script id="review-data" type="application/json">{safe_data_json}</script>
  <script>
    const data = JSON.parse(document.getElementById('review-data').textContent);
    const pageEl = document.getElementById('page');
    const fieldsEl = document.getElementById('fields');
    const tablesEl = document.getElementById('tables');
    function showField(index) {{
      const row = data.fields[index];
      document.querySelectorAll('.field').forEach((el, i) => el.classList.toggle('active', i === index));
      pageEl.innerHTML = '';
      if (!row || !row.image) {{
        pageEl.textContent = 'No page image for this field.';
        return;
      }}
      const img = document.createElement('img');
      img.src = row.image;
      pageEl.appendChild(img);
      if (row.bbox && row.page_width && row.page_height) {{
        img.addEventListener('load', () => {{
          const scaleX = img.clientWidth / row.page_width;
          const scaleY = img.clientHeight / row.page_height;
          const box = document.createElement('div');
          box.className = 'box';
          box.style.left = (row.bbox[0] * scaleX) + 'px';
          box.style.top = (row.bbox[1] * scaleY) + 'px';
          box.style.width = Math.max(8, (row.bbox[2] - row.bbox[0]) * scaleX) + 'px';
          box.style.height = Math.max(8, (row.bbox[3] - row.bbox[1]) * scaleY) + 'px';
          pageEl.appendChild(box);
        }});
      }}
    }}
    data.fields.forEach((row, index) => {{
      const btn = document.createElement('button');
      btn.className = 'field';
      btn.innerHTML = `<strong>${{row.field}}</strong><span>${{row.value ?? ''}}</span><span>parser: ${{row.source_parser ?? ''}} | confidence: ${{row.confidence ?? ''}} | page: ${{row.page ?? ''}}</span>`;
      btn.addEventListener('click', () => showField(index));
      fieldsEl.appendChild(btn);
    }});
    if (data.tables.length) {{
      data.tables.forEach((table) => {{
        const title = document.createElement('div');
        title.textContent = `Page ${{table.page}} Table ${{table.table_index}}`;
        tablesEl.appendChild(title);
        const tbl = document.createElement('table');
        const head = document.createElement('tr');
        (table.headers || []).forEach(h => {{ const th = document.createElement('th'); th.textContent = h; head.appendChild(th); }});
        tbl.appendChild(head);
        (table.rows || []).slice(0, 30).forEach(r => {{
          const tr = document.createElement('tr');
          r.forEach(c => {{ const td = document.createElement('td'); td.textContent = c; tr.appendChild(td); }});
          tbl.appendChild(tr);
        }});
        tablesEl.appendChild(tbl);
      }});
    }} else {{
      tablesEl.textContent = 'No voted tables in this pack.';
    }}
    if (data.fields.length) showField(0);
  </script>
</body>
</html>
"""
    path = out_dir / "review.html"
    path.write_text(html_text, encoding="utf-8")
    return path


def write_customer_pack_readme(manifest: dict[str, object], out_dir: Path) -> Path:
    """Write a concise README for a customer delivery pack."""
    outputs = manifest.get("outputs") if isinstance(manifest.get("outputs"), dict) else {}
    vote = manifest.get("vote") if isinstance(manifest.get("vote"), dict) else {}
    tables = manifest.get("tables") if isinstance(manifest.get("tables"), dict) else {}
    layout = manifest.get("layout") if isinstance(manifest.get("layout"), dict) else {}
    metadata = manifest.get("metadata") if isinstance(manifest.get("metadata"), dict) else {}
    winner = vote.get("winner") if isinstance(vote.get("winner"), dict) else {}
    lines = [
        "# Customer PDF Pack",
        "",
        f"- Source: `{manifest.get('source_file')}`",
        f"- Generated: `{manifest.get('generated_at')}`",
        f"- Best parser: `{winner.get('parser')}`",
        f"- Vote score: `{winner.get('vote_score')}`",
        f"- Table count: `{tables.get('table_count', 0)}`",
        f"- Layout pages: `{layout.get('page_count', 0)}`",
        "",
        "## Outputs",
        "",
    ]
    for key, value in outputs.items():
        if isinstance(value, list):
            joined = ", ".join(f"`{item}`" for item in value)
            lines.append(f"- `{key}`: {joined}")
        elif value:
            lines.append(f"- `{key}`: `{value}`")
    diagnostics = manifest.get("diagnostics") if isinstance(manifest.get("diagnostics"), list) else []
    if diagnostics:
        lines.extend(["", "## Diagnostics", ""])
        for item in diagnostics:
            lines.append(f"- {item}")
    if metadata:
        lines.extend([
            "",
            "## Metadata Summary",
            "",
            f"- Format: `{metadata.get('format')}`",
            f"- Size bytes: `{metadata.get('size_bytes')}`",
        ])
        pdf = metadata.get("pdf") if isinstance(metadata.get("pdf"), dict) else {}
        if pdf:
            lines.extend([
                f"- PDF pages: `{pdf.get('page_count')}`",
                f"- Has sample text layer: `{pdf.get('has_sample_text_layer')}`",
            ])
    readme_path = out_dir / "README.md"
    readme_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return readme_path


def write_customer_pdf_pack(
    file_path: Path,
    out_dir: Path,
    parser_names: list[str],
    start_page: int,
    max_pages: int | None,
    probe_max_pages: int | None,
    table_pages: str,
    layout_max_pages: int | None,
    profile: str,
    min_quality: float,
    fail_on_bad: bool,
    similarity_chars: int,
    include_tables: bool = True,
    include_layout: bool = True,
    field_fusion: bool = True,
    parallel: bool = False,
    timeout_seconds: float | None = None,
    health_cache_path: Path | None = None,
    use_health_cache: bool = False,
    result_cache_dir: Path | None = None,
    use_result_cache: bool = False,
    include_field_evidence: bool = True,
    include_review_html: bool = True,
) -> dict[str, object]:
    """Build a customer-facing PDF pack with text, vote, tables, layout, and metadata."""
    out_dir.mkdir(parents=True, exist_ok=True)
    diagnostics: list[str] = []
    written: list[Path] = []

    metadata = file_metadata(file_path)
    metadata_path = write_json_file(out_dir / "metadata.json", metadata)
    written.append(metadata_path)

    probe_dir = out_dir / "preflight_probe"
    preflight_probe, probe_written = run_parser_probe(
        file_path,
        "pdf",
        parser_names,
        probe_dir,
        start_page,
        probe_max_pages,
        min_quality,
        fail_on_bad,
        parallel,
        keep_outputs=False,
        progress=True,
        timeout_seconds=timeout_seconds,
        health_cache_path=health_cache_path,
        use_health_cache=use_health_cache,
        result_cache_dir=result_cache_dir,
        use_result_cache=use_result_cache,
    )
    written.extend(probe_written)
    ready_parser_names = ready_parsers_from_probe(preflight_probe)
    if not ready_parser_names:
        raise RuntimeError("预检没有解析器通过运行时体检")

    results, normalized_texts, actual_path = collect_parser_outputs(
        file_path,
        "pdf",
        ready_parser_names,
        out_dir,
        "md",
        start_page,
        max_pages,
        parallel,
        progress=True,
        timeout_seconds=timeout_seconds,
        health_cache_path=health_cache_path,
        use_health_cache=use_health_cache,
        result_cache_dir=result_cache_dir,
        use_result_cache=use_result_cache,
    )
    vote_payload = build_vote_payload(
        file_path,
        actual_path,
        "pdf",
        results,
        normalized_texts,
        similarity_chars,
        min_quality,
        fail_on_bad,
        profile,
        preflight_probe,
    )

    layout: dict[str, object] | None = None
    layout_payload: dict[str, object] = {"enabled": include_layout, "page_count": 0, "output": None, "error": None}
    if include_layout:
        try:
            layout = pdf_layout_json(file_path, start_page, layout_max_pages)
            layout_path = write_json_file(out_dir / "layout.json", layout)
            page_map_path = write_json_file(out_dir / "page_map.json", layout_to_page_map(layout))
            written.extend([layout_path, page_map_path])
            pages = layout.get("pages") if isinstance(layout.get("pages"), list) else []
            layout_payload = {
                "enabled": True,
                "page_count": len(pages),
                "output": str(layout_path),
                "page_map": str(page_map_path),
                "error": None,
            }
        except Exception as exc:
            diagnostics.append(f"Layout extraction failed: {exc!r}")
            layout_payload = {"enabled": True, "page_count": 0, "output": None, "error": repr(exc)}

    vote_written = write_vote_outputs(vote_payload, normalized_texts, out_dir, "all")
    written.extend(vote_written)

    tables_payload: dict[str, object] = {"enabled": include_tables, "table_count": 0, "outputs": {}, "error": None}
    table_vote_payload: dict[str, object] | None = None
    if include_tables:
        try:
            page_count = get_pdf_page_count(file_path)
            page_indices = resolve_page_indices(table_pages, page_count, default_all=True)
            tables_dir = out_dir / "tables"
            table_vote_payload = build_table_vote_payload(
                file_path,
                page_indices,
                ["pdfplumber", "pdfplumber-text", "pymupdf-text"],
            )
            table_written = write_table_vote_outputs(table_vote_payload, tables_dir, "all")
            written.extend(table_written)
            tables = table_vote_payload.get("best_tables") if isinstance(table_vote_payload.get("best_tables"), list) else []
            tables_payload = {
                "enabled": True,
                "page_count": len(page_indices),
                "table_count": len(tables),
                "winner": table_vote_payload.get("winner"),
                "outputs": {
                    "vote_report": str(tables_dir / "table_vote_report.md"),
                    "vote_report_json": str(tables_dir / "table_vote_report.json"),
                    "md": str(tables_dir / "best_tables.md"),
                    "csv": str(tables_dir / "best_tables.csv"),
                    "json": str(tables_dir / "best_tables.json"),
                },
                "error": None,
            }
        except Exception as exc:
            diagnostics.append(f"Table extraction failed: {exc!r}")
            tables_payload = {"enabled": True, "table_count": 0, "outputs": {}, "error": repr(exc)}

    delivery = build_customer_delivery(vote_payload, normalized_texts, layout, field_fusion)
    evidence_payload: dict[str, object] = {"enabled": include_field_evidence, "image_count": 0, "outputs": {}, "error": None}
    if delivery and include_field_evidence:
        try:
            evidence_payload = write_field_evidence_pack(file_path, delivery, out_dir)
            evidence_outputs = evidence_payload.get("outputs") if isinstance(evidence_payload.get("outputs"), dict) else {}
            for path_value in evidence_outputs.values():
                if path_value:
                    written.append(Path(str(path_value)))
        except Exception as exc:
            diagnostics.append(f"Field evidence failed: {exc!r}")
            evidence_payload = {"enabled": True, "image_count": 0, "outputs": {}, "error": repr(exc)}

    review_payload: dict[str, object] = {"enabled": include_review_html, "output": None, "error": None}
    if delivery and include_review_html:
        try:
            review_path = write_review_html(file_path, delivery, out_dir, table_vote_payload)
            written.append(review_path)
            review_payload = {"enabled": True, "output": str(review_path), "error": None}
        except Exception as exc:
            diagnostics.append(f"Review HTML failed: {exc!r}")
            review_payload = {"enabled": True, "output": None, "error": repr(exc)}

    customer_written = write_customer_delivery_outputs(delivery, out_dir, "all") if delivery else []
    written.extend(customer_written)

    winner = vote_payload.get("winner") if isinstance(vote_payload.get("winner"), dict) else {}
    manifest = {
        "version": __version__,
        "type": "customer-pdf-pack",
        "generated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        "source_file": str(file_path),
        "actual_file": str(actual_path),
        "profile": profile,
        "schema_version": CUSTOMER_BEST_SCHEMA_VERSION,
        "start_page": start_page,
        "max_pages": max_pages,
        "options": {
            "probe_max_pages": probe_max_pages,
            "table_pages": table_pages,
            "layout_max_pages": layout_max_pages,
            "include_tables": include_tables,
            "include_layout": include_layout,
            "include_field_evidence": include_field_evidence,
            "include_review_html": include_review_html,
            "field_fusion": field_fusion,
            "timeout_seconds": timeout_seconds,
            "parser_health_cache": use_health_cache,
            "health_cache_path": str(health_cache_path) if use_health_cache and health_cache_path else None,
            "result_cache": use_result_cache,
            "result_cache_dir": str(result_cache_dir) if use_result_cache and result_cache_dir else None,
        },
        "metadata": metadata,
        "vote": {
            "winner": winner,
            "quality_gate": vote_payload.get("quality_gate"),
            "ready_parsers": ready_parser_names,
        },
        "tables": tables_payload,
        "layout": layout_payload,
        "field_evidence": evidence_payload,
        "review_html": review_payload,
        "diagnostics": diagnostics,
        "outputs": {
            "metadata": str(metadata_path),
            "manifest": str(out_dir / "manifest.json"),
            "readme": str(out_dir / "README.md"),
            "vote_report": str(out_dir / "vote_report.md"),
            "vote_report_json": str(out_dir / "vote_report.json"),
            "best": [str(out_dir / f"best.{suffix}") for suffix in ["md", "txt", "json"]],
            "customer_best": [str(out_dir / f"customer_best.{suffix}") for suffix in ["md", "txt", "json"] if (out_dir / f"customer_best.{suffix}").exists()],
            "customer_schema": str(out_dir / "customer_best.schema.json") if (out_dir / "customer_best.schema.json").exists() else None,
            "preflight_probe": str(probe_dir / "probe_report.md"),
            "field_evidence": evidence_payload.get("outputs") if isinstance(evidence_payload, dict) else None,
            "review_html": review_payload.get("output") if isinstance(review_payload, dict) else None,
            "tables": tables_payload.get("outputs"),
            "layout": layout_payload.get("output"),
            "page_map": layout_payload.get("page_map"),
        },
    }
    manifest_path = write_json_file(out_dir / "manifest.json", manifest)
    readme_path = write_customer_pack_readme(manifest, out_dir)
    manifest["outputs"]["manifest"] = str(manifest_path)
    manifest["outputs"]["readme"] = str(readme_path)
    write_json_file(manifest_path, manifest)
    return manifest


def safe_report_name(value: object, fallback: str) -> str:
    """Return a filesystem-safe report directory name."""
    raw = str(value or fallback).strip() or fallback
    safe = re.sub(r"[^0-9A-Za-z._\-\u4e00-\u9fff]+", "_", raw)
    safe = safe.strip("._")
    return safe[:80] or fallback


def golden_case_id(file_path: Path, root: Path | None = None) -> str:
    """Build a stable readable case id for a PDF path."""
    try:
        rel = file_path.relative_to(root) if root else file_path.name
    except ValueError:
        rel = file_path.name
    raw = str(rel).replace("\\", "_").replace("/", "_")
    return safe_report_name(Path(raw).with_suffix("").as_posix(), safe_filename_part(file_path.stem, "case"))


def infer_exam_golden_metadata(file_path: Path, roots: list[Path] | None = None) -> dict[str, object]:
    """Infer lightweight exam-PDF metadata from directory and filename."""
    parts = [file_path.stem, *[part for part in file_path.parts if part]]
    joined = " ".join(parts)
    year_match = re.search(r"(20[0-9]{2})", joined)
    year = year_match.group(1) if year_match else None
    season = "上半年" if "上半年" in joined else "下半年" if "下半年" in joined else None
    section = "上午" if "上午" in joined else "案例" if "案例" in joined else "论文" if "论文" in joined else None
    if section is None and "模拟" in joined:
        section = "模拟题"
    chapter_match = re.search(r"第([一二三四五六七八九十0-9]+)章", joined)
    round_name = "第2套" if "第2套" in joined or "第二轮" in joined else "第1套" if "第1" in joined or "第1 套" in joined else None
    category = "mock" if "模拟" in joined or "题库" in joined else "past_exam"
    source_root = None
    if roots:
        for root in roots:
            try:
                file_path.relative_to(root)
                source_root = str(root)
                break
            except ValueError:
                continue
    return {
        "category": category,
        "year": year,
        "season": season,
        "section": section,
        "chapter": chapter_match.group(0) if chapter_match else None,
        "round": round_name,
        "source_root": source_root,
    }


def int_to_chinese(value: int) -> str:
    """Convert a small positive integer to common Chinese numerals."""
    digits = "零一二三四五六七八九"
    if value <= 0:
        return str(value)
    if value < 10:
        return digits[value]
    if value == 10:
        return "十"
    if value < 20:
        return "十" + digits[value % 10]
    if value < 100:
        tens, ones = divmod(value, 10)
        return digits[tens] + "十" + (digits[ones] if ones else "")
    return str(value)


def golden_keywords_for_pdf(file_path: Path, metadata: dict[str, object]) -> list[str]:
    """Build conservative keyword expectations for exam or mock PDFs."""
    name = file_path.stem
    keywords: list[str] = []
    for token in ["系统规划与管理师", "系统规划", "真题", "模拟题", "答案", "解析", "案例", "论文", "上午"]:
        if token in name and token not in keywords:
            keywords.append(token)
    for key in ["year", "season", "section", "chapter"]:
        value = metadata.get(key)
        if value and str(value) not in keywords:
            keywords.append(str(value))
        if key == "chapter" and value:
            digit_match = re.search(r"第([0-9]+)章", str(value))
            if digit_match:
                chinese_number = int_to_chinese(int(digit_match.group(1)))
                alt = f"第{chinese_number}章"
                if alt not in keywords:
                    keywords.append(alt)
    if not any(token in keywords for token in ["系统规划与管理师", "系统规划"]):
        keywords.append("系统")
    return keywords[:8]


def build_golden_case_payload(
    file_path: Path,
    root: Path | None,
    parsers: str,
    max_pages: int,
    min_quality: float,
    min_chars: int,
    required_keywords: list[str] | None = None,
    name_prefix: str | None = None,
    roots: list[Path] | None = None,
) -> dict[str, object]:
    """Build one generic PDF golden case payload."""
    metadata = infer_exam_golden_metadata(file_path, roots)
    keywords = required_keywords if required_keywords is not None else golden_keywords_for_pdf(file_path, metadata)
    case_id = golden_case_id(file_path, root)
    name = f"{name_prefix}-{case_id}" if name_prefix else case_id
    return {
        "name": name,
        "file": str(file_path.resolve()),
        "command": "vote",
        "profile": "none",
        "parsers": parsers,
        "max_pages": max_pages,
        "expected": {
            "min_vote_score": min_quality,
            "quality_gate_passed": True,
            "min_non_space_chars": min_chars,
            "quality_label_not": ["empty", "bad"],
            "text_contains_any": keywords,
        },
        "metadata": metadata,
    }


def collect_golden_pdf_files(roots: list[Path], recursive: bool = True) -> list[Path]:
    """Collect PDF files from one or more roots."""
    files: list[Path] = []
    seen: set[Path] = set()
    for raw_root in roots:
        root = raw_root.resolve()
        candidates: list[Path]
        if root.is_file() and root.suffix.lower() == ".pdf":
            candidates = [root]
        elif root.is_dir():
            pattern = "**/*.pdf" if recursive else "*.pdf"
            candidates = sorted(path for path in root.glob(pattern) if path.is_file())
        else:
            candidates = []
        for candidate in candidates:
            resolved = candidate.resolve()
            if resolved not in seen:
                seen.add(resolved)
                files.append(resolved)
    return files


def write_golden_library_readme(manifest: dict[str, object], out_dir: Path) -> Path:
    """Write a human-readable README for an initialized golden library."""
    rows = manifest.get("cases") if isinstance(manifest.get("cases"), list) else []
    lines = [
        "# Golden PDF Case Library",
        "",
        f"- Version: `{manifest.get('version')}`",
        f"- Generated at: `{manifest.get('generated_at')}`",
        f"- Case count: `{manifest.get('case_count')}`",
        f"- PDF count: `{manifest.get('pdf_count')}`",
        f"- Parsers: `{manifest.get('parsers')}`",
        f"- Max pages: `{manifest.get('max_pages')}`",
        "",
        "## Run",
        "",
        "```powershell",
        f"python scripts\\parse_document_compare.py eval-golden \"{out_dir}\" --recursive --timeout 120 --parser-health-cache --result-cache",
        "```",
        "",
        "## Cases",
        "",
        "| Case | Category | Year | Section | PDF |",
        "|---|---|---|---|---|",
    ]
    for row in rows:
        if not isinstance(row, dict):
            continue
        metadata = row.get("metadata") if isinstance(row.get("metadata"), dict) else {}
        lines.append(
            f"| `{row.get('name')}` | {metadata.get('category') or ''} | {metadata.get('year') or ''} | "
            f"{metadata.get('section') or ''} | `{row.get('file')}` |"
        )
    path = out_dir / "README.md"
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return path


def write_golden_library_index(manifest: dict[str, object], out_dir: Path) -> tuple[Path, Path]:
    """Write index JSON and Markdown for a generated golden library."""
    index_json = write_json_file(out_dir / "index.json", manifest)
    readme = write_golden_library_readme(manifest, out_dir)
    return index_json, readme


def load_golden_case_files(path: Path, recursive: bool = False) -> list[Path]:
    """Return candidate golden case JSON files."""
    if path.is_file():
        return [path]
    if not path.exists() or not path.is_dir():
        return []
    pattern = "**/*.json" if recursive else "*.json"
    case_files: list[Path] = []
    for candidate in sorted(path.glob(pattern)):
        if not candidate.is_file() or candidate.name in GOLDEN_REPORT_NAMES:
            continue
        try:
            data = json.loads(candidate.read_text(encoding="utf-8"))
        except Exception:
            continue
        if isinstance(data, dict) and ("file" in data or "path" in data) and "expected" in data:
            case_files.append(candidate)
    return case_files


def compare_expected_value(actual: object, expected: object) -> tuple[bool, str]:
    """Compare an actual value to an expected golden value."""
    if isinstance(expected, dict):
        if "present" in expected:
            present = actual is not None and normalize_field_value(actual) != ""
            return present is bool(expected.get("present")), str(expected.get("present"))
        if "equals" in expected:
            return compare_expected_value(actual, expected.get("equals"))
        if "contains" in expected:
            actual_text = normalize_field_value(actual)
            expected_text = normalize_field_value(expected.get("contains"))
            return bool(expected_text and expected_text in actual_text), f"contains {expected.get('contains')}"
        if "regex" in expected:
            pattern = str(expected.get("regex") or "")
            return bool(pattern and re.search(pattern, str(actual or ""))), f"regex {pattern}"
    actual_norm = normalize_field_value(actual)
    expected_norm = normalize_field_value(expected)
    return actual_norm == expected_norm, str(expected)


def evaluate_vote_expectations(payload: dict[str, object], expected: dict[str, object] | None) -> dict[str, object]:
    """Evaluate a vote payload against golden expectations."""
    expected = expected or {}
    winner = payload.get("winner") if isinstance(payload.get("winner"), dict) else {}
    winner_text = ""
    if winner.get("parser"):
        texts = payload.get("_normalized_texts")
        if isinstance(texts, dict):
            winner_text = str(texts.get(str(winner.get("parser"))) or "")
    checks: list[dict[str, object]] = []

    def add_check(name: str, actual: object, expected_value: object, passed: bool) -> None:
        checks.append({
            "name": name,
            "passed": bool(passed),
            "expected": expected_value,
            "actual": actual,
        })

    if "winner_parser" in expected:
        expected_winner = expected.get("winner_parser")
        allowed = expected_winner if isinstance(expected_winner, list) else [expected_winner]
        actual = winner.get("parser")
        add_check("winner_parser", actual, expected_winner, actual in allowed)

    if "profile_resolved" in expected:
        profile = payload.get("profile") if isinstance(payload.get("profile"), dict) else {}
        actual = profile.get("resolved")
        add_check("profile_resolved", actual, expected.get("profile_resolved"), actual == expected.get("profile_resolved"))

    if "min_vote_score" in expected:
        actual_score = float(winner.get("vote_score") or 0.0)
        expected_score = float(expected.get("min_vote_score") or 0.0)
        add_check("min_vote_score", actual_score, expected_score, actual_score >= expected_score)

    if "quality_gate_passed" in expected:
        gate = payload.get("quality_gate") if isinstance(payload.get("quality_gate"), dict) else {}
        actual = bool(gate.get("passed"))
        add_check("quality_gate_passed", actual, expected.get("quality_gate_passed"), actual is bool(expected.get("quality_gate_passed")))

    if "min_non_space_chars" in expected:
        actual_chars = int(winner.get("non_space_chars") or 0)
        expected_chars = int(expected.get("min_non_space_chars") or 0)
        add_check("min_non_space_chars", actual_chars, expected_chars, actual_chars >= expected_chars)

    if "min_chars" in expected:
        actual_chars = int(winner.get("chars") or 0)
        expected_chars = int(expected.get("min_chars") or 0)
        add_check("min_chars", actual_chars, expected_chars, actual_chars >= expected_chars)

    if "quality_label" in expected:
        actual = winner.get("quality_label")
        expected_label = expected.get("quality_label")
        allowed = expected_label if isinstance(expected_label, list) else [expected_label]
        add_check("quality_label", actual, expected_label, actual in allowed)

    if "quality_label_not" in expected:
        actual = winner.get("quality_label")
        forbidden = expected.get("quality_label_not")
        forbidden_values = forbidden if isinstance(forbidden, list) else [forbidden]
        add_check("quality_label_not", actual, forbidden, actual not in forbidden_values)

    if "max_duplicate_line_ratio" in expected:
        repetition = winner.get("repetition_metrics") if isinstance(winner.get("repetition_metrics"), dict) else {}
        actual = float(repetition.get("duplicate_line_ratio") or 0.0)
        expected_ratio = float(expected.get("max_duplicate_line_ratio") or 0.0)
        add_check("max_duplicate_line_ratio", actual, expected_ratio, actual <= expected_ratio)

    if "max_repeated_ngram_ratio" in expected:
        repetition = winner.get("repetition_metrics") if isinstance(winner.get("repetition_metrics"), dict) else {}
        actual = float(repetition.get("repeated_ngram_ratio") or 0.0)
        expected_ratio = float(expected.get("max_repeated_ngram_ratio") or 0.0)
        add_check("max_repeated_ngram_ratio", actual, expected_ratio, actual <= expected_ratio)

    if "text_contains" in expected:
        expected_text = str(expected.get("text_contains") or "")
        add_check("text_contains", expected_text if expected_text in winner_text else None, expected_text, bool(expected_text and expected_text in winner_text))

    if "text_contains_all" in expected:
        expected_items = expected.get("text_contains_all")
        items = [str(item) for item in expected_items] if isinstance(expected_items, list) else [str(expected_items)]
        missing = [item for item in items if item and item not in winner_text]
        add_check("text_contains_all", {"missing": missing}, items, not missing)

    if "text_contains_any" in expected:
        expected_items = expected.get("text_contains_any")
        items = [str(item) for item in expected_items] if isinstance(expected_items, list) else [str(expected_items)]
        matched = [item for item in items if item and item in winner_text]
        add_check("text_contains_any", {"matched": matched}, items, bool(matched))

    if "text_regex" in expected:
        pattern = str(expected.get("text_regex") or "")
        add_check("text_regex", pattern if pattern and re.search(pattern, winner_text) else None, pattern, bool(pattern and re.search(pattern, winner_text)))

    invoice_fields = winner.get("invoice_fields") if isinstance(winner.get("invoice_fields"), dict) else {}
    validation = invoice_fields.get("validation") if isinstance(invoice_fields.get("validation"), dict) else {}
    if "validation_status" in expected:
        actual = validation.get("status")
        add_check("validation_status", actual, expected.get("validation_status"), actual == expected.get("validation_status"))

    expected_fields = expected.get("fields") if isinstance(expected.get("fields"), dict) else {}
    for field, expected_value in expected_fields.items():
        actual = invoice_fields.get(field)
        passed, expected_label = compare_expected_value(actual, expected_value)
        add_check(f"field:{field}", actual, expected_label, passed)

    return {
        "passed": all(bool(check.get("passed")) for check in checks) if checks else True,
        "check_count": len(checks),
        "checks": checks,
    }


def golden_eval_to_markdown(payload: dict[str, object]) -> str:
    """Render golden evaluation report as Markdown."""
    lines = [
        "# Golden PDF Evaluation",
        "",
        f"- Cases: `{payload.get('case_count')}`",
        f"- Passed: `{payload.get('passed_count')}`",
        f"- Failed: `{payload.get('failed_count')}`",
        "",
        "## Cases",
        "",
        "| Case | Passed | Winner | Checks | File |",
        "|---|---:|---|---:|---|",
    ]
    rows = payload.get("rows") if isinstance(payload.get("rows"), list) else []
    for row in rows:
        if not isinstance(row, dict):
            continue
        lines.append(
            f"| {row.get('name')} | {row.get('passed')} | {row.get('winner_parser')} | "
            f"{row.get('passed_checks')}/{row.get('check_count')} | `{row.get('file')}` |"
        )
    for row in rows:
        if not isinstance(row, dict) or row.get("passed"):
            continue
        lines.extend(["", f"## Failed: {row.get('name')}", ""])
        for check in row.get("checks", []):
            if not isinstance(check, dict) or check.get("passed"):
                continue
            lines.append(
                f"- `{check.get('name')}` expected `{check.get('expected')}`, got `{check.get('actual')}`"
            )
    return "\n".join(lines) + "\n"


def cmd_eval_golden(args: argparse.Namespace) -> int:
    """Run golden PDF cases and report parser strategy regressions."""
    root = args.path.resolve()
    case_files = load_golden_case_files(root, args.recursive)
    if not case_files:
        print(f"错误：没有找到 golden case JSON：{root}", file=sys.stderr)
        return 1
    if args.out_dir:
        out_dir = args.out_dir.resolve()
    elif root.is_dir():
        out_dir = root / "golden_eval_output"
    else:
        out_dir = root.parent / "golden_eval_output"
    out_dir.mkdir(parents=True, exist_ok=True)
    health_cache_path = Path(args.health_cache).resolve() if args.health_cache else default_parser_health_cache_path(out_dir)
    result_cache_dir = Path(args.result_cache_dir).resolve() if args.result_cache_dir else default_parser_result_cache_dir(out_dir)

    rows: list[dict[str, object]] = []
    print("Golden PDF 评测")
    print(f"用例数：{len(case_files)}")
    for index, case_file in enumerate(case_files, 1):
        case = json.loads(case_file.read_text(encoding="utf-8"))
        name = str(case.get("name") or case_file.stem)
        print(f"[{index}/{len(case_files)}] {name}...", end=" ", flush=True)
        case_dir = out_dir / safe_report_name(name, f"case_{index}")
        case_dir.mkdir(parents=True, exist_ok=True)
        try:
            if str(case.get("command", "vote")) != "vote":
                raise ValueError("eval-golden 当前只支持 command=vote")
            raw_file = Path(str(case.get("file") or case.get("path") or ""))
            file_path = raw_file if raw_file.is_absolute() else case_file.parent / raw_file
            if not file_path.exists():
                raise FileNotFoundError(str(file_path))
            if detect_format(file_path) != "pdf":
                raise ValueError("eval-golden 当前只评测 PDF vote 用例")
            parser_selection = args.parsers if args.parsers else case.get("parsers")
            if isinstance(parser_selection, list):
                parser_arg = ",".join(str(item) for item in parser_selection)
            else:
                parser_arg = str(parser_selection) if parser_selection else None
            parser_names, unknown = validate_parser_selection("pdf", parser_arg)
            if unknown:
                print(f"警告：跳过不支持解析器 {', '.join(unknown)}", file=sys.stderr)
            if not parser_names:
                raise ValueError("没有可用的 PDF 解析器")
            raw_max_pages = args.max_pages if args.max_pages >= 0 else int(case.get("max_pages", 0))
            max_pages = None if raw_max_pages == 0 else raw_max_pages
            start_page = args.start_page if args.start_page is not None else int(case.get("start_page", 1))
            profile = args.profile if args.profile else str(case.get("profile", "auto"))
            results, normalized_texts, actual_path = collect_parser_outputs(
                file_path,
                "pdf",
                parser_names,
                case_dir,
                "md",
                start_page,
                max_pages,
                args.parallel,
                progress=False,
                timeout_seconds=args.timeout,
                health_cache_path=health_cache_path,
                use_health_cache=bool(args.parser_health_cache),
                result_cache_dir=result_cache_dir,
                use_result_cache=bool(args.result_cache),
            )
            vote_payload = build_vote_payload(
                file_path,
                actual_path,
                "pdf",
                results,
                normalized_texts,
                args.similarity_chars,
                args.min_quality,
                args.fail_on_bad,
                profile,
            )
            written = write_vote_outputs(vote_payload, normalized_texts, case_dir, "all")
            vote_payload["_normalized_texts"] = normalized_texts
            evaluation = evaluate_vote_expectations(
                vote_payload,
                case.get("expected") if isinstance(case.get("expected"), dict) else {},
            )
            vote_payload.pop("_normalized_texts", None)
            winner = vote_payload.get("winner") if isinstance(vote_payload.get("winner"), dict) else {}
            checks = evaluation.get("checks") if isinstance(evaluation.get("checks"), list) else []
            row = {
                "name": name,
                "case_file": str(case_file),
                "file": str(file_path),
                "passed": bool(evaluation.get("passed")),
                "check_count": evaluation.get("check_count"),
                "passed_checks": sum(1 for check in checks if isinstance(check, dict) and check.get("passed")),
                "winner_parser": winner.get("parser"),
                "vote_score": winner.get("vote_score"),
                "checks": checks,
                "outputs": [str(path) for path in written],
                "error": None,
            }
            rows.append(row)
            print("ok" if row["passed"] else "FAIL")
        except Exception as exc:
            rows.append({
                "name": name,
                "case_file": str(case_file),
                "file": str(case.get("file") or case.get("path") or ""),
                "passed": False,
                "check_count": 1,
                "passed_checks": 0,
                "winner_parser": None,
                "vote_score": 0.0,
                "checks": [{"name": "case_runtime", "passed": False, "expected": "ok", "actual": repr(exc)}],
                "outputs": [],
                "error": repr(exc),
            })
            print(f"FAIL {exc}")

    report = {
        "version": __version__,
        "type": "eval-golden",
        "generated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        "root": str(root),
        "out_dir": str(out_dir),
        "case_count": len(rows),
        "passed_count": sum(1 for row in rows if row.get("passed")),
        "failed_count": sum(1 for row in rows if not row.get("passed")),
        "rows": rows,
    }
    json_path = write_json_file(out_dir / "golden_report.json", report)
    md_path = out_dir / "golden_report.md"
    md_path.write_text(golden_eval_to_markdown(report), encoding="utf-8")
    print(f"报告：{md_path}")
    print(f"机器可读报告：{json_path}")
    return 0 if report["failed_count"] == 0 else 2


def cmd_init_golden(args: argparse.Namespace) -> int:
    """Initialize generic PDF golden case JSON files from directories."""
    roots = [path.resolve() for path in args.paths]
    out_dir = args.out_dir.resolve()
    pdf_files = collect_golden_pdf_files(roots, recursive=bool(args.recursive))
    if not pdf_files:
        print("错误：没有找到 PDF 文件", file=sys.stderr)
        return 1
    out_dir.mkdir(parents=True, exist_ok=True)
    cases_dir = out_dir / "cases"
    cases_dir.mkdir(parents=True, exist_ok=True)

    max_pages = args.max_pages
    parser_names = [part.strip() for part in (args.parsers or GOLDEN_DEFAULT_PARSERS).split(",") if part.strip()]
    if args.include_ocr and "ocr-tesseract" not in parser_names:
        parser_names.append("ocr-tesseract")
    parsers = ",".join(parser_names)
    min_chars = int(args.min_non_space_chars)
    min_quality = float(args.min_quality)
    manifest_cases: list[dict[str, object]] = []
    written: list[Path] = []
    for index, pdf_path in enumerate(pdf_files, 1):
        root_for_case = next((root for root in roots if root.is_dir() and pdf_path.is_relative_to(root)), None)
        payload = build_golden_case_payload(
            pdf_path,
            root_for_case,
            parsers,
            max_pages,
            min_quality,
            min_chars,
            name_prefix=args.name_prefix,
            roots=roots,
        )
        case_name = safe_report_name(payload.get("name"), f"case_{index:03d}")
        case_path = cases_dir / f"{index:03d}_{case_name}.json"
        if case_path.exists() and not args.force:
            raise FileExistsError(f"case 已存在，使用 --force 覆盖：{case_path}")
        write_json_file(case_path, payload)
        written.append(case_path)
        manifest_cases.append({
            "name": payload.get("name"),
            "case_file": str(case_path),
            "file": str(pdf_path),
            "size_bytes": pdf_path.stat().st_size,
            "metadata": payload.get("metadata"),
            "expected": payload.get("expected"),
        })

    categories = Counter(
        str((row.get("metadata") or {}).get("category"))
        for row in manifest_cases
        if isinstance(row.get("metadata"), dict)
    )
    sections = Counter(
        str((row.get("metadata") or {}).get("section"))
        for row in manifest_cases
        if isinstance(row.get("metadata"), dict) and (row.get("metadata") or {}).get("section")
    )
    years = Counter(
        str((row.get("metadata") or {}).get("year"))
        for row in manifest_cases
        if isinstance(row.get("metadata"), dict) and (row.get("metadata") or {}).get("year")
    )
    manifest = {
        "version": __version__,
        "type": "golden-case-library",
        "generated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        "roots": [str(root) for root in roots],
        "out_dir": str(out_dir),
        "case_count": len(written),
        "pdf_count": len(pdf_files),
        "parsers": parsers,
        "max_pages": max_pages,
        "min_quality": min_quality,
        "min_non_space_chars": min_chars,
        "summary": {
            "categories": dict(categories),
            "sections": dict(sections),
            "years": dict(years),
        },
        "cases": manifest_cases,
    }
    index_json, readme = write_golden_library_index(manifest, out_dir)
    print("Golden 样本库已生成")
    print(f"PDF 数：{len(pdf_files)}")
    print(f"Case 数：{len(written)}")
    print(f"目录：{out_dir}")
    print(f"索引：{index_json}")
    print(f"说明：{readme}")
    print(f"评测命令：python scripts\\parse_document_compare.py eval-golden \"{out_dir}\" --recursive --timeout 120 --parser-health-cache --result-cache")
    return 0


def probe_status_for_result(
    result: ParseResult,
    dependency_available: bool,
    min_quality: float,
    fail_on_bad: bool,
) -> dict[str, object]:
    """Classify a parser after a real extraction probe."""
    if not dependency_available:
        return {
            "status": "dependency_missing",
            "ready": False,
            "gate": quality_gate_status(None, min_quality, fail_on_bad),
            "reason": "dependency_missing",
        }
    if result.status != "ok":
        return {
            "status": "runtime_failed",
            "ready": False,
            "gate": quality_gate_status(None, min_quality, fail_on_bad),
            "reason": result.error or result.status,
        }
    gate = quality_gate_status(result, min_quality, fail_on_bad)
    if not gate["passed"]:
        return {
            "status": "quality_failed",
            "ready": False,
            "gate": gate,
            "reason": gate["reason"],
        }
    return {
        "status": "ready",
        "ready": True,
        "gate": gate,
        "reason": "ok",
    }


def build_probe_payload(
    file_path: Path,
    actual_path: Path,
    fmt: str,
    parser_names: list[str],
    results: list[ParseResult],
    normalized_texts: dict[str, str],
    min_quality: float,
    fail_on_bad: bool,
    max_pages: int | None,
    start_page: int,
) -> dict[str, object]:
    """Build a runtime parser probe report from actual extraction attempts."""
    dependency_map = {
        str(row.get("parser")): row
        for row in parser_dependency_rows(fmt)
    }
    result_map = {result.parser: result for result in results}
    rows: list[dict[str, object]] = []
    ready_parsers: list[str] = []
    diagnostics: list[str] = []

    for parser_name in parser_names:
        dependency = dependency_map.get(parser_name, {})
        dependency_available = bool(dependency.get("available"))
        result = result_map.get(parser_name)
        if result is None:
            result = ParseResult(
                parser=parser_name,
                status="skipped",
                seconds=0.0,
                error="parser did not run",
            )
        status = probe_status_for_result(result, dependency_available, min_quality, fail_on_bad)
        text = normalized_texts.get(parser_name, "")
        repetition = repetition_metrics(text) if text else None
        row = {
            "parser": parser_name,
            "kind": dependency.get("kind"),
            "dependency_available": dependency_available,
            "commands": dependency.get("commands", []),
            "modules": dependency.get("modules", []),
            "runtime_status": result.status,
            "probe_status": status["status"],
            "ready": status["ready"],
            "reason": status["reason"],
            "seconds": result.seconds,
            "quality_score": result.quality_score,
            "quality_label": result.quality_label,
            "chars": result.chars,
            "non_space_chars": result.non_space_chars,
            "chinese_chars": result.chinese_chars,
            "cid_markers": result.cid_markers,
            "control_chars": result.control_chars,
            "repetition": repetition,
            "output_file": result.output_file,
            "error": result.error,
            "quality_gate": status["gate"],
            "preview": result.preview,
        }
        if row["ready"]:
            ready_parsers.append(parser_name)
        elif row["probe_status"] == "runtime_failed" and dependency_available:
            diagnostics.append(f"{parser_name} dependency is available but the real extraction failed: {result.error}")
        elif row["probe_status"] == "dependency_missing":
            diagnostics.append(f"{parser_name} is registered but its dependency is missing.")
        elif row["probe_status"] == "quality_failed":
            diagnostics.append(f"{parser_name} ran but did not pass the quality gate ({status['reason']}).")
        rows.append(row)

    if not diagnostics:
        diagnostics.append("All probed parsers passed the runtime quality gate.")

    summary = {
        "parser_count": len(parser_names),
        "ready_count": len(ready_parsers),
        "dependency_missing_count": sum(1 for row in rows if row["probe_status"] == "dependency_missing"),
        "runtime_failed_count": sum(1 for row in rows if row["probe_status"] == "runtime_failed"),
        "quality_failed_count": sum(1 for row in rows if row["probe_status"] == "quality_failed"),
        "ready_parsers": ready_parsers,
    }
    return {
        "version": __version__,
        "type": "probe",
        "source_file": str(file_path),
        "actual_file": str(actual_path),
        "format": fmt,
        "generated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        "start_page": start_page,
        "max_pages": max_pages,
        "min_quality": min_quality,
        "fail_on_bad": fail_on_bad,
        "summary": summary,
        "diagnostics": diagnostics,
        "rows": rows,
        "results": [asdict(result) for result in results],
    }


def probe_to_markdown(payload: dict[str, object]) -> str:
    """Render a runtime parser probe report."""
    summary = payload.get("summary") if isinstance(payload.get("summary"), dict) else {}
    rows = payload.get("rows") if isinstance(payload.get("rows"), list) else []
    lines = [
        "# Parser Runtime Probe",
        "",
        f"- Source: `{payload.get('source_file')}`",
        f"- Format: `{payload.get('format')}`",
        f"- Ready parsers: `{summary.get('ready_count', 0)}` / `{summary.get('parser_count', 0)}`",
        f"- Dependency missing: `{summary.get('dependency_missing_count', 0)}`",
        f"- Runtime failed: `{summary.get('runtime_failed_count', 0)}`",
        f"- Quality failed: `{summary.get('quality_failed_count', 0)}`",
        "",
        "## Parsers",
        "",
        "| Parser | Probe | Dep | Runtime | Quality | Score | Chars | Time(s) | Reason |",
        "|---|---|---:|---|---|---:|---:|---:|---|",
    ]
    for row in rows:
        if not isinstance(row, dict):
            continue
        reason = str(row.get("reason") or row.get("error") or "").replace("|", "\\|")
        lines.append(
            f"| {row.get('parser')} | {row.get('probe_status')} | {row.get('dependency_available')} | "
            f"{row.get('runtime_status')} | {row.get('quality_label')} | {row.get('quality_score')} | "
            f"{row.get('chars')} | {row.get('seconds')} | {reason[:240]} |"
        )

    diagnostics = payload.get("diagnostics") if isinstance(payload.get("diagnostics"), list) else []
    lines.extend(["", "## Diagnostics", ""])
    for item in diagnostics:
        lines.append(f"- {item}")
    return "\n".join(lines) + "\n"


def write_probe_outputs(payload: dict[str, object], out_dir: Path) -> tuple[Path, Path]:
    """Write JSON and Markdown probe reports."""
    out_dir.mkdir(parents=True, exist_ok=True)
    json_path = write_json_file(out_dir / "probe_report.json", payload)
    md_path = out_dir / "probe_report.md"
    md_path.write_text(probe_to_markdown(payload), encoding="utf-8")
    return json_path, md_path


def run_parser_probe(
    file_path: Path,
    fmt: str,
    parser_names: list[str],
    out_dir: Path,
    start_page: int,
    max_pages: int | None,
    min_quality: float,
    fail_on_bad: bool,
    parallel: bool = False,
    keep_outputs: bool = False,
    progress: bool = False,
    timeout_seconds: float | None = None,
    health_cache_path: Path | None = None,
    use_health_cache: bool = False,
    result_cache_dir: Path | None = None,
    use_result_cache: bool = False,
) -> tuple[dict[str, object], list[Path]]:
    """Run parser runtime probe and write probe reports."""
    results, normalized_texts, actual_path = collect_parser_outputs(
        file_path,
        fmt,
        parser_names,
        out_dir,
        "md" if keep_outputs else "txt",
        start_page,
        max_pages,
        parallel,
        progress=progress,
        timeout_seconds=timeout_seconds,
        health_cache_path=health_cache_path,
        use_health_cache=use_health_cache,
        result_cache_dir=result_cache_dir,
        use_result_cache=use_result_cache,
    )
    payload = build_probe_payload(
        file_path,
        actual_path,
        fmt,
        parser_names,
        results,
        normalized_texts,
        min_quality,
        fail_on_bad,
        max_pages,
        start_page,
    )
    json_path, md_path = write_probe_outputs(payload, out_dir)
    return payload, [json_path, md_path]


def ready_parsers_from_probe(payload: dict[str, object]) -> list[str]:
    """Return parsers that passed a probe payload."""
    summary = payload.get("summary") if isinstance(payload.get("summary"), dict) else {}
    ready = summary.get("ready_parsers") if isinstance(summary.get("ready_parsers"), list) else []
    return [str(name) for name in ready if name]


def build_diagnostics(results: list[ParseResult]) -> list[str]:
    """Build concise diagnostics for extraction quality issues."""
    ok_results = [r for r in results if r.status == "ok"]
    diagnostics: list[str] = []
    if not ok_results:
        return ["No parser produced usable output."]

    bad_or_empty = [r for r in ok_results if r.quality_label in ("empty", "bad")]
    if len(bad_or_empty) == len(ok_results):
        diagnostics.append("All successful parser outputs look empty or poor; OCR or a cleaner source file is likely needed.")

    if any(r.cid_markers > 0 for r in ok_results):
        diagnostics.append("One or more outputs contain `(cid:...)` markers, which usually indicates embedded font mapping problems.")

    if any(r.control_chars > max(20, r.non_space_chars * 0.05) for r in ok_results):
        diagnostics.append("One or more outputs contain many control characters, so visual searchability may not equal extractable text quality.")

    if not diagnostics:
        diagnostics.append("No obvious extraction-quality warning was detected from text-only heuristics.")

    return diagnostics


def ocr_fallback_notice(enabled: bool) -> str | None:
    """Return a clear OCR fallback notice."""
    if not enabled:
        return None
    return "OCR fallback was requested. Use the `ocr` command when ordinary extraction is bad or empty; OCR requires PyMuPDF, Pillow, pytesseract, and the system Tesseract binary."


def render_pdf_pages(pdf_path: Path, page_indices: list[int], out_dir: Path, dpi: int = 200) -> list[dict[str, object]]:
    """Render selected PDF pages to PNG files using PyMuPDF."""
    if not module_exists("fitz"):
        raise ImportError("页面截图需要 PyMuPDF/fitz，运行: pip install pymupdf")
    import fitz  # type: ignore

    out_dir.mkdir(parents=True, exist_ok=True)
    rendered: list[dict[str, object]] = []
    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)
    doc = fitz.open(pdf_path)
    try:
        for page_idx in page_indices:
            if page_idx < 0 or page_idx >= doc.page_count:
                continue
            page = doc.load_page(page_idx)
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            out_path = out_dir / f"{pdf_path.stem}_page_{page_idx + 1:03d}.png"
            pix.save(out_path)
            rendered.append(
                {
                    "page": page_idx + 1,
                    "file": str(out_path),
                    "width": pix.width,
                    "height": pix.height,
                    "dpi": dpi,
                }
            )
    finally:
        doc.close()
    return rendered


def ensure_tesseract_available() -> None:
    """Validate OCR dependencies and raise a readable setup error."""
    missing_python = []
    if not module_exists("pytesseract"):
        missing_python.append("pytesseract")
    if not module_exists("PIL"):
        missing_python.append("pillow")
    if missing_python:
        raise ImportError(
            "OCR 需要 Python 依赖："
            + ", ".join(missing_python)
            + "。安装示例：pip install pytesseract pillow"
        )
    if shutil.which("tesseract") is None:
        raise RuntimeError("OCR 需要系统 Tesseract 可执行文件，请安装 tesseract-ocr 并加入 PATH")


def ocr_install_hint(missing_python: list[str], missing_system: list[str]) -> str | None:
    """Build a concise OCR install hint without empty commands."""
    hints: list[str] = []
    if missing_python:
        hints.append(f"pip install {' '.join(missing_python)}")
    if missing_system:
        hints.append("install system Tesseract OCR, add tesseract to PATH, and install requested language data such as chi_sim")
    return "; ".join(hints) if hints else None


def opendataloader_install_hint(missing_python: list[str], missing_system: list[str]) -> str | None:
    """Build a concise opendataloader-pdf install hint without empty commands."""
    hints: list[str] = []
    if missing_python:
        hints.append(f"pip install {' '.join(missing_python)}")
    if missing_system:
        hints.append("install Java JDK/JRE and ensure java is on PATH")
    return "; ".join(hints) if hints else None


def command_install_hint(missing_commands: list[str]) -> str | None:
    """Build install hints for external command parsers."""
    if not missing_commands:
        return None
    hints: list[str] = []
    if any("@pspdfkit/pdf-to-markdown" in item for item in missing_commands):
        hints.append(
            "install Nutrient/PSPDFKit PDF to Markdown CLI: npm install -g @pspdfkit/pdf-to-markdown; "
            "ensure pdf-to-markdown is on PATH. Official binaries currently target macOS/Linux."
        )
    other = [item for item in missing_commands if "@pspdfkit/pdf-to-markdown" not in item]
    hints.extend(f"install external command for {item}" for item in other)
    return "; ".join(hints)


def ocr_pdf_pages(pdf_path: Path, page_indices: list[int], lang: str, dpi: int) -> list[dict[str, object]]:
    """Run Tesseract OCR on selected PDF pages."""
    ensure_tesseract_available()
    if not module_exists("fitz"):
        raise ImportError("OCR 渲染 PDF 页面需要 PyMuPDF/fitz，运行: pip install pymupdf")

    import fitz  # type: ignore
    import pytesseract  # type: ignore
    from PIL import Image  # type: ignore

    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)
    results: list[dict[str, object]] = []
    doc = fitz.open(pdf_path)
    try:
        for page_idx in page_indices:
            page = doc.load_page(page_idx)
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            image = Image.open(io.BytesIO(pix.tobytes("png")))
            start = time.perf_counter()
            text = pytesseract.image_to_string(image, lang=lang)
            normalized = normalize_text(text)
            quality_score, quality_label = quality_from_text(normalized)
            results.append(
                {
                    "page": page_idx + 1,
                    "seconds": round(time.perf_counter() - start, 3),
                    "chars": len(normalized),
                    "quality_score": quality_score,
                    "quality_label": quality_label,
                    "text": normalized,
                }
            )
    finally:
        doc.close()
    return results


def ocr_pages_to_text(pages: list[dict[str, object]]) -> str:
    """Join OCR page results into a page-marked text stream."""
    parts: list[str] = []
    for page in pages:
        parts.append(f"\n\n<!-- page {page.get('page')} ocr -->\n\n")
        parts.append(str(page.get("text") or ""))
    return normalize_text("\n".join(parts))


def maybe_ocr_fallback(
    file_path: Path,
    current_text: str,
    current_result: ParseResult,
    enabled: bool,
    min_quality: float,
    pages: str,
    lang: str,
    dpi: int,
) -> tuple[str, ParseResult, dict[str, object] | None]:
    """Run OCR when enabled and current quality is below the threshold."""
    if not enabled:
        return current_text, current_result, None
    if file_path.suffix.lower() != ".pdf":
        return current_text, current_result, {"attempted": False, "reason": "not_pdf"}
    if current_result.quality_score >= min_quality and current_result.quality_label not in QUALITY_BAD_LABELS:
        return current_text, current_result, {"attempted": False, "reason": "text_quality_ok"}

    try:
        page_count = get_pdf_page_count(file_path)
        page_indices = resolve_page_indices(pages, page_count, default_all=False)
        start = time.perf_counter()
        ocr_pages = ocr_pdf_pages(file_path, page_indices, lang, dpi)
        ocr_text = ocr_pages_to_text(ocr_pages)
        ocr_result = result_from_text("ocr-tesseract", ocr_text, time.perf_counter() - start)
        used = ocr_result.quality_score > current_result.quality_score or current_result.quality_label in QUALITY_BAD_LABELS
        payload = {
            "attempted": True,
            "used": used,
            "backend": "tesseract",
            "pages": [page["page"] for page in ocr_pages],
            "lang": lang,
            "dpi": dpi,
            "quality_score": ocr_result.quality_score,
            "quality_label": ocr_result.quality_label,
            "error": None,
        }
        if used:
            return ocr_text, ocr_result, payload
        return current_text, current_result, payload
    except Exception as exc:
        return current_text, current_result, {
            "attempted": True,
            "used": False,
            "backend": "tesseract",
            "pages": pages,
            "lang": lang,
            "dpi": dpi,
            "quality_score": 0.0,
            "quality_label": "failed",
            "error": repr(exc),
        }


def write_ocr_output(
    out_path: Path,
    source_file: Path,
    pages: list[dict[str, object]],
    output_format: str,
    lang: str,
    dpi: int,
) -> Path:
    """Write OCR output as txt, md, or json."""
    out_path.parent.mkdir(parents=True, exist_ok=True)
    if output_format == "json":
        payload = {
            "version": __version__,
            "source_file": str(source_file),
            "backend": "tesseract",
            "lang": lang,
            "dpi": dpi,
            "page_count": len(pages),
            "pages": pages,
        }
        out_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    elif output_format == "md":
        parts = [
            "# OCR Output",
            "",
            f"- Source: `{source_file}`",
            "- Backend: `tesseract`",
            f"- Language: `{lang}`",
            f"- DPI: `{dpi}`",
            "",
        ]
        for page in pages:
            parts.extend([f"## Page {page['page']}", "", str(page["text"]), ""])
        out_path.write_text("\n".join(parts), encoding="utf-8")
    else:
        parts = []
        for page in pages:
            parts.append(f"--- page {page['page']} ---")
            parts.append(str(page["text"]))
        out_path.write_text("\n\n".join(parts) + ("\n" if parts else ""), encoding="utf-8")
    return out_path


def pdf_layout_json(pdf_path: Path, start_page: int = 1, max_pages: int | None = None) -> dict[str, object]:
    """Extract PDF page/block/line/span layout with coordinates using PyMuPDF."""
    if not module_exists("fitz"):
        raise ImportError("layout-json 需要 PyMuPDF/fitz，运行: pip install pymupdf")
    import fitz  # type: ignore

    doc = fitz.open(pdf_path)
    pages: list[dict[str, object]] = []
    try:
        start_index, end_index = compute_page_range(doc.page_count, start_page, max_pages)
        for page_index in range(start_index, end_index + 1):
            page = doc.load_page(page_index)
            data = page.get_text("dict")
            page_payload: dict[str, object] = {
                "page": page_index + 1,
                "width": page.rect.width,
                "height": page.rect.height,
                "rotation": page.rotation,
                "blocks": [],
            }
            for block_index, block in enumerate(data.get("blocks", []), start=1):
                block_payload: dict[str, object] = {
                    "block_id": block_index,
                    "type": block.get("type"),
                    "bbox": block.get("bbox"),
                    "lines": [],
                }
                for line_index, line in enumerate(block.get("lines", []), start=1):
                    spans_payload = []
                    line_text_parts = []
                    for span_index, span in enumerate(line.get("spans", []), start=1):
                        span_text = span.get("text", "")
                        line_text_parts.append(span_text)
                        spans_payload.append(
                            {
                                "span_id": span_index,
                                "text": span_text,
                                "bbox": span.get("bbox"),
                                "font": span.get("font"),
                                "size": span.get("size"),
                                "flags": span.get("flags"),
                                "color": span.get("color"),
                            }
                        )
                    block_payload["lines"].append(
                        {
                            "line_id": line_index,
                            "bbox": line.get("bbox"),
                            "text": "".join(line_text_parts),
                            "spans": spans_payload,
                        }
                    )
                page_payload["blocks"].append(block_payload)
            pages.append(page_payload)
    finally:
        doc.close()

    return {
        "version": __version__,
        "source_file": str(pdf_path),
        "format": "pdf-layout",
        "start_page": start_page,
        "max_pages": max_pages,
        "page_count": len(pages),
        "pages": pages,
    }


def layout_to_page_map(layout: dict[str, object]) -> list[dict[str, object]]:
    """Build a lightweight page map from layout JSON."""
    page_map: list[dict[str, object]] = []
    for page in layout.get("pages", []):  # type: ignore[union-attr]
        if not isinstance(page, dict):
            continue
        texts: list[str] = []
        for block in page.get("blocks", []):
            if not isinstance(block, dict):
                continue
            for line in block.get("lines", []):
                if isinstance(line, dict) and line.get("text"):
                    texts.append(str(line["text"]))
        joined = normalize_text("\n".join(texts))
        page_map.append(
            {
                "page": page.get("page"),
                "width": page.get("width"),
                "height": page.get("height"),
                "chars": len(joined),
                "preview": joined[:500],
            }
        )
    return page_map


INVOICE_SUMMARY_COLUMNS = [
    "source_file",
    "parser",
    "quality_label",
    "quality_score",
    "invoice_type",
    "invoice_number",
    "invoice_date",
    "buyer_name",
    "buyer_tax_id",
    "seller_name",
    "seller_tax_id",
    "total_amount",
    "total_tax",
    "total_with_tax",
    "total_with_tax_cn",
    "drawer",
    "validation_status",
    "missing_fields",
]


def clean_invoice_scan_text(text: str) -> str:
    """Remove layout noise that commonly appears in invoice parser output."""
    return re.sub(r"[\s|`]+", "", text or "")


def regex_first(pattern: str, text: str, flags: int = 0) -> str | None:
    match = re.search(pattern, text, flags)
    if not match:
        return None
    value = match.group(1).strip()
    return value or None


def extract_invoice_items(text: str) -> list[dict[str, object]]:
    """Extract invoice line items from text-oriented parser output."""
    items: list[dict[str, object]] = []
    item_pattern = re.compile(
        r"(?P<name>\*[^*\n]+\*[^0-9\n|]+?)\s+"
        r"(?P<unit>[\u4e00-\u9fffA-Za-z]{1,8})\s+"
        r"(?P<quantity>[0-9]+(?:\.[0-9]+)?)\s+"
        r"(?P<unit_price>[0-9]+(?:\.[0-9]+)?)\s+"
        r"(?P<amount>[0-9]+(?:\.[0-9]{2})?)\s+"
        r"(?P<tax_rate>[0-9]+%)\s+"
        r"(?P<tax_amount>[0-9]+(?:\.[0-9]{2})?)"
    )
    for raw_line in text.splitlines():
        line = raw_line.replace("|", " ")
        line = re.sub(r"\s+", " ", line).strip()
        if not line or "项目名称" in line or "---" in line:
            continue
        match = item_pattern.search(line)
        if not match:
            continue
        data = match.groupdict()
        items.append(
            {
                "name": data["name"].strip(),
                "unit": data["unit"],
                "quantity": decimal_from_text(data["quantity"]),
                "unit_price": decimal_from_text(data["unit_price"]),
                "amount": decimal_from_text(data["amount"]),
                "tax_rate": data["tax_rate"],
                "tax_amount": decimal_from_text(data["tax_amount"]),
            }
        )
    return items


def invoice_item_repetition_metrics(items: object) -> dict[str, object]:
    """Measure duplicate invoice line items after structured extraction."""
    if not isinstance(items, list) or not items:
        return {
            "duplicate_item_ratio": 0.0,
            "duplicate_item_count": 0,
            "unique_item_count": 0,
            "max_item_repeat": 0,
            "penalty": 0.0,
            "score": 1.0,
        }

    signatures: list[tuple[object, ...]] = []
    for item in items:
        if not isinstance(item, dict):
            continue
        signatures.append((
            re.sub(r"\s+", "", str(item.get("name") or "")),
            item.get("unit"),
            item.get("quantity"),
            item.get("unit_price"),
            item.get("amount"),
            item.get("tax_rate"),
            item.get("tax_amount"),
        ))

    if not signatures:
        return {
            "duplicate_item_ratio": 0.0,
            "duplicate_item_count": 0,
            "unique_item_count": 0,
            "max_item_repeat": 0,
            "penalty": 0.0,
            "score": 1.0,
        }

    counts = Counter(signatures)
    duplicate_count = sum(count - 1 for count in counts.values() if count > 1)
    duplicate_ratio = duplicate_count / max(1, len(signatures))
    max_repeat = max(counts.values())
    penalty = min(0.2, duplicate_ratio * 0.35 + max(0, max_repeat - 2) * 0.03)
    return {
        "duplicate_item_ratio": round(duplicate_ratio, 4),
        "duplicate_item_count": duplicate_count,
        "unique_item_count": len(counts),
        "max_item_repeat": max_repeat,
        "penalty": round(penalty, 4),
        "score": round(max(0.0, 1.0 - penalty / 0.2), 4),
    }


def validate_invoice_fields(fields: dict[str, object], strict: bool = False) -> dict[str, object]:
    """Run invoice consistency checks."""
    required = [
        "invoice_number",
        "invoice_date",
        "buyer_name",
        "buyer_tax_id",
        "seller_name",
        "seller_tax_id",
        "total_with_tax",
    ]
    missing = [key for key in required if not fields.get(key)]
    checks: list[dict[str, object]] = []

    total_amount = fields.get("total_amount")
    total_tax = fields.get("total_tax")
    total_with_tax = fields.get("total_with_tax")
    if isinstance(total_amount, float) and isinstance(total_tax, float) and isinstance(total_with_tax, float):
        expected = round(total_amount + total_tax, 2)
        passed = abs(expected - total_with_tax) <= 0.01
        checks.append(
            {
                "name": "amount_plus_tax_equals_total",
                "passed": passed,
                "expected": expected,
                "actual": total_with_tax,
            }
        )

    items = fields.get("items")
    if isinstance(items, list) and items:
        repetition = invoice_item_repetition_metrics(items)
        if int(repetition.get("duplicate_item_count") or 0) > 0:
            checks.append(
                {
                    "name": "line_items_not_repeated",
                    "passed": False,
                    "expected": "unique line items",
                    "actual": repetition,
                }
            )
        item_amount = round(sum(float(item.get("amount") or 0.0) for item in items if isinstance(item, dict)), 2)
        item_tax = round(sum(float(item.get("tax_amount") or 0.0) for item in items if isinstance(item, dict)), 2)
        if isinstance(total_amount, float):
            checks.append(
                {
                    "name": "item_amount_sum_equals_total_amount",
                    "passed": abs(item_amount - total_amount) <= 0.01,
                    "expected": item_amount,
                    "actual": total_amount,
                }
            )
        if isinstance(total_tax, float):
            checks.append(
                {
                    "name": "item_tax_sum_equals_total_tax",
                    "passed": abs(item_tax - total_tax) <= 0.01,
                    "expected": item_tax,
                    "actual": total_tax,
                }
            )

    if strict:
        if not isinstance(items, list) or not items:
            checks.append({"name": "has_line_items", "passed": False, "expected": ">=1", "actual": 0})
        if fields.get("invoice_number") and not re.fullmatch(r"[0-9]{8,30}", str(fields["invoice_number"])):
            checks.append({"name": "invoice_number_format", "passed": False, "expected": "8-30 digits", "actual": fields["invoice_number"]})
        for key in ("buyer_tax_id", "seller_tax_id"):
            value = fields.get(key)
            if value and not re.fullmatch(r"[0-9A-Z]{15,25}", str(value)):
                checks.append({"name": f"{key}_format", "passed": False, "expected": "15-25 alnum", "actual": value})

    status = "ok"
    if missing:
        status = "missing_fields"
    if any(not check["passed"] for check in checks):
        status = "failed_checks"
    return {"status": status, "missing_fields": missing, "checks": checks}


def extract_invoice_fields_from_text(
    text: str,
    source_file: Path | None = None,
    parser_name: str | None = None,
    result: ParseResult | None = None,
    strict: bool = False,
) -> dict[str, object]:
    """Extract structured fields from a Chinese electronic invoice."""
    normalized = normalize_text(text)
    scan = clean_invoice_scan_text(normalized)

    tax_ids = re.findall(r"统一社会信用代码/纳税人识别号[:：]?([0-9A-Z]{15,25})", scan)
    amount_pairs = re.findall(r"合计[¥￥]?([0-9]+(?:\.[0-9]{2}))[¥￥]?([0-9]+(?:\.[0-9]{2}))", scan)
    total_amount = decimal_from_text(amount_pairs[-1][0]) if amount_pairs else None
    total_tax = decimal_from_text(amount_pairs[-1][1]) if amount_pairs else None
    drawer = regex_first(r"开票人[:：]?\s*([^\s\n|]+)", normalized)

    fields: dict[str, object] = {
        "version": __version__,
        "profile": "invoice",
        "source_file": str(source_file) if source_file else None,
        "parser": parser_name,
        "quality_score": result.quality_score if result else None,
        "quality_label": result.quality_label if result else None,
        "invoice_type": regex_first(r"(电子发票[（(][^）)]+[）)]|增值税电子普通发票|增值税专用发票|增值税普通发票)", scan),
        "invoice_number": regex_first(r"发票号码[:：]?([0-9]{8,30})", scan),
        "invoice_date": regex_first(r"开票日期[:：]?([0-9]{4}年[0-9]{1,2}月[0-9]{1,2}日|[0-9]{4}[-/][0-9]{1,2}[-/][0-9]{1,2})", scan),
        "buyer_name": regex_first(r"购.*?名称[:：]?(.+?)(?:销|买售|统一社会信用代码)", scan),
        "buyer_tax_id": tax_ids[0] if len(tax_ids) >= 1 else None,
        "seller_name": regex_first(r"销.*?名称[:：]?(.+?)(?:买售|售方|方信|统一社会信用代码|项目名称|息)", scan),
        "seller_tax_id": tax_ids[1] if len(tax_ids) >= 2 else None,
        "total_amount": total_amount,
        "total_tax": total_tax,
        "total_with_tax": decimal_from_text(regex_first(r"[（(]小写[）)][¥￥]?([0-9]+(?:\.[0-9]{2}))", scan)),
        "total_with_tax_cn": regex_first(r"价税合计[（(]大写[）)](.+?)[（(]小写[）)]", scan),
        "drawer": drawer or regex_first(r"开票人[:：]?([\u4e00-\u9fffA-Za-z0-9·]{1,30})", scan),
        "items": extract_invoice_items(normalized),
        "raw_text_chars": len(normalized),
    }
    fields["validation"] = validate_invoice_fields(fields, strict=strict)
    return fields


def score_invoice_fields(fields: dict[str, object]) -> int:
    """Score invoice extraction completeness for choosing among parsers."""
    score = 0
    for key in [
        "invoice_number",
        "invoice_date",
        "buyer_name",
        "buyer_tax_id",
        "seller_name",
        "seller_tax_id",
        "total_with_tax",
        "drawer",
    ]:
        if fields.get(key):
            score += 2
    if fields.get("items"):
        score += 4
    validation = fields.get("validation")
    if isinstance(validation, dict) and validation.get("status") == "ok":
        score += 3
    return score


def extract_invoice_fields_from_file(
    file_path: Path,
    parser_name: str = "auto",
    start_page: int = 1,
    max_pages: int | None = None,
    strict: bool = False,
) -> dict[str, object]:
    """Extract invoice fields from a file, trying multiple parsers in auto mode."""
    fmt = detect_format(file_path)
    if fmt is None:
        raise ValueError(f"不支持的文件格式：{file_path.suffix}")

    candidates = [resolve_parser_name(parser_name)] if parser_name != "auto" else get_parsers_for_format(fmt)
    extracted: list[dict[str, object]] = []
    errors: list[str] = []
    for candidate in candidates:
        start = time.perf_counter()
        try:
            raw_text = extract_text_with_parser(file_path, candidate, start_page, max_pages)
            normalized = normalize_text(raw_text)
            result = text_stats(candidate, "ok", time.perf_counter() - start, normalized, None)
            fields = extract_invoice_fields_from_text(normalized, file_path, candidate, result, strict=strict)
            fields["extraction_seconds"] = result.seconds
            extracted.append(fields)
        except Exception as exc:
            errors.append(f"{candidate}: {exc}")

    if not extracted:
        raise RuntimeError("发票字段抽取失败：" + "; ".join(errors))
    return max(
        extracted,
        key=lambda item: (
            score_invoice_fields(item),
            float(item.get("quality_score") or 0.0),
            int(item.get("raw_text_chars") or 0),
        ),
    )


def classify_document(text: str, metadata: dict[str, object] | None = None) -> dict[str, object]:
    """Classify documents for document-intelligence workflows."""
    scan = clean_invoice_scan_text(text)
    invoice_hits = sum(1 for token in ["发票号码", "开票日期", "购买方", "销售方", "价税合计", "税率"] if token in scan)
    contract_hits = sum(1 for token in ["合同", "甲方", "乙方", "签订", "违约", "协议", "履行期限", "争议解决"] if token in scan)
    bank_statement_hits = sum(1 for token in ["银行流水", "交易日期", "借方", "贷方", "余额", "对方户名", "交易金额", "账户"] if token in scan)
    quotation_hits = sum(1 for token in ["报价单", "报价", "单价", "含税", "有效期", "付款方式", "报价人"] if token in scan)
    purchase_order_hits = sum(1 for token in ["采购订单", "采购单", "订单编号", "供应商", "采购方", "交货日期", "PO"] if token in scan)
    report_hits = sum(1 for token in ["报告", "摘要", "结论", "数据分析", "指标", "图表", "附录"] if token in scan)
    annual_report_hits = sum(1 for token in ["年度报告", "年报", "资产负债表", "利润表", "现金流量表", "董事会", "审计报告"] if token in scan)
    textbook_hits = sum(1 for token in ["目录", "第", "章", "练习", "考试", "知识点"] if token in scan)
    notice_hits = sum(1 for token in ["通知", "公告", "请", "附件", "日期"] if token in scan)

    fmt = metadata.get("format") if isinstance(metadata, dict) else None
    pdf_meta = metadata.get("pdf") if isinstance(metadata, dict) else None
    has_text_layer = None
    page_count = None
    if isinstance(pdf_meta, dict):
        has_text_layer = pdf_meta.get("has_sample_text_layer")
        page_count = pdf_meta.get("page_count")

    scores = {
        "invoice": invoice_hits / 6,
        "contract": contract_hits / 8,
        "bank_statement": bank_statement_hits / 8,
        "quotation": quotation_hits / 7,
        "purchase_order": purchase_order_hits / 7,
        "report": report_hits / 7,
        "annual_report": annual_report_hits / 7,
        "textbook": textbook_hits / 6,
        "notice": notice_hits / 5,
    }
    best_profile, best_score = max(scores.items(), key=lambda item: item[1])
    if invoice_hits >= 3:
        best_profile = "invoice"
        best_score = scores["invoice"]
    elif fmt == "pdf" and has_text_layer is False:
        best_profile = "scanned_pdf"
        best_score = 0.8
    elif fmt == "pdf" and not compact_text(text):
        best_profile = "image_or_empty_pdf"
        best_score = 0.7
    elif best_score < 0.34:
        best_profile = "generic"
        best_score = 0.5

    strategy = {
        "invoice": "extract_fields_and_verify",
        "contract": "extract_best_text_tables_and_key_terms",
        "bank_statement": "extract_tables_and_preserve_transaction_rows",
        "quotation": "extract_tables_and_price_terms",
        "purchase_order": "extract_order_tables_and_supplier_terms",
        "report": "extract_best_text_tables_and_summary_sections",
        "annual_report": "extract_financial_tables_and_section_chunks",
        "textbook": "chunk_by_page_and_build_knowledge_pack",
        "notice": "extract_text_and_metadata",
        "scanned_pdf": "ocr_then_chunk",
        "image_or_empty_pdf": "ocr_required",
        "generic": "extract_best_text_and_chunk",
    }.get(best_profile, "extract_best_text_and_chunk")

    return {
        "profile": best_profile,
        "confidence": round(min(1.0, best_score), 3),
        "strategy": strategy,
        "signals": {
            "invoice_hits": invoice_hits,
            "contract_hits": contract_hits,
            "bank_statement_hits": bank_statement_hits,
            "quotation_hits": quotation_hits,
            "purchase_order_hits": purchase_order_hits,
            "report_hits": report_hits,
            "annual_report_hits": annual_report_hits,
            "textbook_hits": textbook_hits,
            "notice_hits": notice_hits,
            "format": fmt,
            "page_count": page_count,
            "has_sample_text_layer": has_text_layer,
        },
    }


def write_invoice_fields(fields: dict[str, object], out_path: Path, output_format: str) -> Path:
    """Write invoice fields as JSON or Markdown."""
    out_path.parent.mkdir(parents=True, exist_ok=True)
    if output_format == "json":
        out_path.write_text(json.dumps(fields, ensure_ascii=False, indent=2), encoding="utf-8")
    else:
        validation = fields.get("validation") if isinstance(fields.get("validation"), dict) else {}
        lines = [
            "# Invoice Fields",
            "",
            f"- Source: `{fields.get('source_file')}`",
            f"- Parser: `{fields.get('parser')}`",
            f"- Quality: `{fields.get('quality_score')}` (`{fields.get('quality_label')}`)",
            f"- Validation: `{validation.get('status')}`",
            "",
            "| Field | Value |",
            "|---|---|",
        ]
        for key in INVOICE_SUMMARY_COLUMNS[4:16]:
            lines.append(f"| `{key}` | {fields.get(key) or ''} |")
        lines.extend(["", "## Items", ""])
        items = fields.get("items")
        if isinstance(items, list) and items:
            lines.extend(["| Name | Unit | Quantity | Unit Price | Amount | Tax Rate | Tax Amount |", "|---|---|---:|---:|---:|---|---:|"])
            for item in items:
                if isinstance(item, dict):
                    lines.append(
                        f"| {item.get('name') or ''} | {item.get('unit') or ''} | {item.get('quantity') or ''} | "
                        f"{item.get('unit_price') or ''} | {item.get('amount') or ''} | {item.get('tax_rate') or ''} | {item.get('tax_amount') or ''} |"
                    )
        else:
            lines.append("No line items extracted.")
        out_path.write_text("\n".join(lines), encoding="utf-8")
    return out_path


def invoice_summary_row(fields: dict[str, object]) -> dict[str, object]:
    """Flatten invoice fields for tabular exports."""
    validation = fields.get("validation") if isinstance(fields.get("validation"), dict) else {}
    missing = validation.get("missing_fields") if isinstance(validation, dict) else []
    row = {key: fields.get(key) for key in INVOICE_SUMMARY_COLUMNS}
    row["validation_status"] = validation.get("status") if isinstance(validation, dict) else None
    row["missing_fields"] = ",".join(missing) if isinstance(missing, list) else ""
    return row


def collect_input_files(path: Path, extensions: set[str], recursive: bool) -> list[Path]:
    """Collect one input file or files under a directory."""
    if path.is_file():
        return [path]
    if not path.exists() or not path.is_dir():
        raise FileNotFoundError(f"路径不存在：{path}")
    files: list[Path] = []
    for ext in sorted(extensions):
        iterator = path.rglob(f"*{ext}") if recursive else path.glob(f"*{ext}")
        files.extend(item for item in iterator if item.is_file())
    return sorted(set(files))


def write_invoice_xlsx(rows: list[dict[str, object]], out_path: Path) -> Path:
    """Write invoice summary and line items to XLSX."""
    if not module_exists("openpyxl"):
        raise ImportError("导出 XLSX 需要 openpyxl，运行: pip install openpyxl")
    from openpyxl import Workbook  # type: ignore

    out_path.parent.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    ws = wb.active
    ws.title = "invoices"
    ws.append(INVOICE_SUMMARY_COLUMNS)
    for fields in rows:
        flat = invoice_summary_row(fields)
        ws.append([flat.get(column) for column in INVOICE_SUMMARY_COLUMNS])

    item_ws = wb.create_sheet("items")
    item_columns = ["source_file", "invoice_number", "name", "unit", "quantity", "unit_price", "amount", "tax_rate", "tax_amount"]
    item_ws.append(item_columns)
    for fields in rows:
        items = fields.get("items")
        if not isinstance(items, list):
            continue
        for item in items:
            if not isinstance(item, dict):
                continue
            item_ws.append([
                fields.get("source_file"),
                fields.get("invoice_number"),
                item.get("name"),
                item.get("unit"),
                item.get("quantity"),
                item.get("unit_price"),
                item.get("amount"),
                item.get("tax_rate"),
                item.get("tax_amount"),
            ])

    wb.save(out_path)
    return out_path


def write_json_file(path: Path, payload: object) -> Path:
    """Write a JSON file with UTF-8 encoding."""
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return path


def write_knowledge_pack(
    file_path: Path,
    out_dir: Path,
    parser_name: str = "auto",
    start_page: int = 1,
    max_pages: int | None = None,
    chunk_by: str = "page",
    chunk_size: int = 2000,
    overlap: int = 200,
    min_quality: float = 0.5,
) -> dict[str, object]:
    """Generate a traceable RAG/knowledge package for one document."""
    out_dir.mkdir(parents=True, exist_ok=True)
    metadata = file_metadata(file_path)
    text, result = extract_best_text(file_path, parser_name, start_page, max_pages)
    classification = classify_document(text, metadata)
    gate = quality_gate_status(result, min_quality, False)

    layout: dict[str, object] | None = None
    page_map: list[dict[str, object]] = []
    if file_path.suffix.lower() == ".pdf":
        try:
            layout = pdf_layout_json(file_path, start_page, max_pages)
            page_map = layout_to_page_map(layout)
        except Exception as exc:
            page_map = [{"error": repr(exc)}]

    if chunk_by == "page" and file_path.suffix.lower() == ".pdf" and result.parser in PAGE_CHUNK_PARSERS:
        page_extractor = get_pdf_page_extractor(result.parser)
        chunks = split_pages_into_chunks(page_extractor(file_path, start_page, max_pages) if page_extractor else [])
    else:
        chunks = split_text_into_chunks(text, chunk_size, overlap)

    source_md = out_dir / "source.md"
    source_txt = out_dir / "source.txt"
    source_md.write_text(text + "\n", encoding="utf-8")
    source_txt.write_text(text + "\n", encoding="utf-8")

    metadata_path = write_json_file(out_dir / "metadata.json", metadata)
    quality_path = write_json_file(out_dir / "quality_report.json", {"quality": asdict(result), "gate": gate})
    page_map_path = write_json_file(out_dir / "page_map.json", page_map)
    layout_path = write_json_file(out_dir / "layout.json", layout) if layout is not None else None
    chunk_paths = write_chunks(
        chunks,
        out_dir / "chunks",
        "all",
        {
            "version": __version__,
            "source_file": str(file_path),
            "parser": result.parser,
            "format": detect_format(file_path),
            "chunk_by": chunk_by,
        },
    )

    manifest = {
        "version": __version__,
        "type": "knowledge-pack",
        "source_file": str(file_path),
        "parser": result.parser,
        "classification": classification,
        "quality": asdict(result),
        "quality_gate": gate,
        "chunk_count": len(chunks),
        "outputs": {
            "source_md": str(source_md),
            "source_txt": str(source_txt),
            "metadata": str(metadata_path),
            "quality_report": str(quality_path),
            "page_map": str(page_map_path),
            "layout": str(layout_path) if layout_path else None,
            "chunks": [str(path) for path in chunk_paths],
        },
    }
    manifest_path = write_json_file(out_dir / "manifest.json", manifest)
    readme_lines = [
        "# Knowledge Pack",
        "",
        f"- Source: `{file_path}`",
        f"- Parser: `{result.parser}`",
        f"- Classification: `{classification['profile']}`",
        f"- Quality: `{result.quality_score}` (`{result.quality_label}`)",
        f"- Chunk count: `{len(chunks)}`",
        f"- Manifest: `{manifest_path}`",
    ]
    readme_path = out_dir / "README.md"
    readme_path.write_text("\n".join(readme_lines), encoding="utf-8")
    manifest["outputs"]["manifest"] = str(manifest_path)
    manifest["outputs"]["readme"] = str(readme_path)
    write_json_file(manifest_path, manifest)
    return manifest


# ---------------------------------------------------------------------------
# RAG/问答和文档差异辅助
# ---------------------------------------------------------------------------

SEARCH_STOPWORDS = {
    "the", "and", "or", "of", "to", "in", "for", "a", "an", "is", "are", "this", "that",
    "什么", "哪些", "如何", "是否", "可以", "这个", "那个", "里面", "文档", "文件",
    "的", "了", "和", "与", "或", "是", "在", "有", "我", "你", "他", "她", "它",
}


def tokenize_for_search(text: str) -> list[str]:
    """Tokenize mixed Chinese/English text for deterministic local retrieval."""
    tokens = [token.lower() for token in _RE_TOKEN.findall(text or "")]
    return [token for token in tokens if token and token not in SEARCH_STOPWORDS]


def chunk_snippet(text: str, query_tokens: list[str], max_chars: int = RAG_SNIPPET_CHARS) -> str:
    """Return a compact snippet around the first query-token hit."""
    normalized = normalize_text(text)
    if len(normalized) <= max_chars:
        return normalized
    lowered = normalized.lower()
    hit_positions = [lowered.find(token.lower()) for token in query_tokens if token and lowered.find(token.lower()) >= 0]
    center = min(hit_positions) if hit_positions else 0
    start = max(0, center - max_chars // 3)
    end = min(len(normalized), start + max_chars)
    start = max(0, end - max_chars)
    prefix = "..." if start > 0 else ""
    suffix = "..." if end < len(normalized) else ""
    return prefix + normalized[start:end] + suffix


def load_chunks_jsonl(path: Path) -> list[dict[str, object]]:
    chunks: list[dict[str, object]] = []
    for line in path.read_text(encoding="utf-8").splitlines():
        if not line.strip():
            continue
        payload = json.loads(line)
        if isinstance(payload, dict) and payload.get("text"):
            chunks.append(payload)
    return chunks


def load_chunks_json(path: Path) -> list[dict[str, object]]:
    payload = json.loads(path.read_text(encoding="utf-8"))
    if isinstance(payload, dict):
        metadata = payload.get("metadata") if isinstance(payload.get("metadata"), dict) else {}
        chunks = payload.get("chunks")
        if isinstance(chunks, list):
            return [{**metadata, **chunk} for chunk in chunks if isinstance(chunk, dict) and chunk.get("text")]
    if isinstance(payload, list):
        return [chunk for chunk in payload if isinstance(chunk, dict) and chunk.get("text")]
    return []


def load_chunks_from_manifest(manifest_path: Path) -> list[dict[str, object]]:
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    outputs = manifest.get("outputs") if isinstance(manifest, dict) else {}
    if not isinstance(outputs, dict):
        return []
    chunk_paths = outputs.get("chunks")
    if not isinstance(chunk_paths, list):
        return []
    chunks: list[dict[str, object]] = []
    for item in chunk_paths:
        path = Path(str(item))
        if not path.exists():
            continue
        if path.suffix.lower() == ".jsonl":
            chunks.extend(load_chunks_jsonl(path))
        elif path.suffix.lower() == ".json":
            chunks.extend(load_chunks_json(path))
    return chunks


def load_chunks_from_index(index_path: Path) -> list[dict[str, object]]:
    index = json.loads(index_path.read_text(encoding="utf-8"))
    manifests = index.get("manifests") if isinstance(index, dict) else []
    chunks: list[dict[str, object]] = []
    if isinstance(manifests, list):
        for item in manifests:
            manifest_path = Path(str(item))
            if manifest_path.exists():
                chunks.extend(load_chunks_from_manifest(manifest_path))
    return chunks


def collect_chunks_from_pack(path: Path) -> list[dict[str, object]]:
    """Load chunks from a knowledge pack, batch index, chunks file, or pack directory."""
    if path.is_file():
        lower_name = path.name.lower()
        if lower_name == "manifest.json":
            return load_chunks_from_manifest(path)
        if lower_name == "index.json":
            return load_chunks_from_index(path)
        if path.suffix.lower() == ".jsonl":
            return load_chunks_jsonl(path)
        if path.suffix.lower() == ".json":
            return load_chunks_json(path)
        return []

    candidates = [
        path / "chunks" / "chunks.jsonl",
        path / "chunks.jsonl",
        path / "manifest.json",
        path / "index.json",
    ]
    for candidate in candidates:
        if candidate.exists():
            loaded = collect_chunks_from_pack(candidate)
            if loaded:
                return loaded

    chunks: list[dict[str, object]] = []
    for candidate in sorted(path.rglob("chunks.jsonl")):
        chunks.extend(load_chunks_jsonl(candidate))
    return chunks


def collect_qa_chunks(
    source: Path,
    parser_name: str,
    start_page: int,
    max_pages: int | None,
    chunk_by: str,
    chunk_size: int,
    overlap: int,
) -> tuple[list[dict[str, object]], dict[str, object]]:
    """Collect searchable chunks from a knowledge pack or directly from a document."""
    if source.is_dir() or source.suffix.lower() in {".jsonl", ".json"}:
        chunks = collect_chunks_from_pack(source)
        if not chunks:
            raise RuntimeError(f"未在知识包或 chunks 文件中找到可检索文本：{source}")
        return chunks, {"source_type": "knowledge_pack_or_chunks", "source": str(source)}

    fmt = detect_format(source)
    if fmt is None:
        raise ValueError(f"不支持的输入：{source.suffix}")

    text, result = extract_best_text(source, parser_name, start_page, max_pages)
    if fmt == "pdf" and chunk_by == "page" and result.parser in PAGE_CHUNK_PARSERS:
        page_extractor = get_pdf_page_extractor(result.parser)
        chunks = split_pages_into_chunks(page_extractor(source, start_page, max_pages) if page_extractor else [])
    else:
        chunks = split_text_into_chunks(text, chunk_size, overlap)

    for chunk in chunks:
        chunk.setdefault("source_file", str(source))
        chunk.setdefault("parser", result.parser)
        chunk.setdefault("format", fmt)
        chunk.setdefault("chunk_by", chunk_by if fmt == "pdf" else "char")
    return chunks, {"source_type": "document", "source": str(source), "parser": result.parser, "quality": asdict(result)}


def rank_chunks(question: str, chunks: list[dict[str, object]], top_k: int) -> tuple[list[dict[str, object]], list[str]]:
    """Rank chunks with a transparent token-overlap score."""
    query_tokens = tokenize_for_search(question)
    query_set = set(query_tokens)
    scored: list[dict[str, object]] = []
    compact_question = "".join(query_tokens)
    for chunk in chunks:
        text = str(chunk.get("text") or "")
        if not text.strip():
            continue
        text_tokens = tokenize_for_search(text)
        if not text_tokens:
            continue
        token_counts: dict[str, int] = {}
        for token in text_tokens:
            token_counts[token] = token_counts.get(token, 0) + 1
        overlap = sum(min(3, token_counts.get(token, 0)) for token in query_set)
        density = overlap / max(1, len(set(text_tokens)) ** 0.5)
        phrase_bonus = 1.0 if compact_question and compact_question in "".join(text_tokens) else 0.0
        score = round(density + phrase_bonus, 4)
        if score <= 0:
            continue
        scored.append(
            {
                **chunk,
                "retrieval_score": score,
                "snippet": chunk_snippet(text, query_tokens),
            }
        )
    scored.sort(key=lambda item: (float(item.get("retrieval_score") or 0.0), int(item.get("chars") or 0)), reverse=True)
    return scored[:top_k], query_tokens


def extract_relevant_sentences(question: str, matches: list[dict[str, object]], max_sentences: int) -> list[dict[str, object]]:
    """Build an extractive answer from the best-scoring source sentences."""
    query_tokens = set(tokenize_for_search(question))
    candidates: list[dict[str, object]] = []
    for index, match in enumerate(matches, start=1):
        text = normalize_text(str(match.get("text") or match.get("snippet") or ""))
        sentences = [part.strip() for part in re.split(r"(?<=[。！？!?\.])\s+|\n+", text) if part.strip()]
        if not sentences and text:
            sentences = [text[:RAG_SNIPPET_CHARS]]
        for sentence in sentences:
            tokens = set(tokenize_for_search(sentence))
            score = len(query_tokens & tokens)
            if score > 0:
                candidates.append(
                    {
                        "citation": index,
                        "score": score,
                        "text": sentence[:RAG_SNIPPET_CHARS],
                        "source_file": match.get("source_file"),
                        "page": match.get("page"),
                        "chunk_id": match.get("chunk_id"),
                    }
                )
    candidates.sort(key=lambda item: (int(item["score"]), len(str(item["text"]))), reverse=True)
    return candidates[:max_sentences]


def source_citation(match: dict[str, object], index: int) -> dict[str, object]:
    return {
        "citation": index,
        "source_file": match.get("source_file"),
        "page": match.get("page"),
        "chunk_id": match.get("chunk_id"),
        "parser": match.get("parser"),
        "score": match.get("retrieval_score"),
        "snippet": match.get("snippet"),
    }


def build_qa_payload(
    source: Path,
    question: str,
    parser_name: str,
    start_page: int,
    max_pages: int | None,
    chunk_by: str,
    chunk_size: int,
    overlap: int,
    top_k: int,
    answer_sentences: int,
) -> dict[str, object]:
    chunks, source_info = collect_qa_chunks(source, parser_name, start_page, max_pages, chunk_by, chunk_size, overlap)
    matches, query_tokens = rank_chunks(question, chunks, top_k)
    answers = extract_relevant_sentences(question, matches, answer_sentences)
    citations = [source_citation(match, index) for index, match in enumerate(matches, start=1)]
    return {
        "version": __version__,
        "type": "qa",
        "method": "extractive_local_retrieval",
        "question": question,
        "query_tokens": query_tokens,
        "source": source_info,
        "chunk_count": len(chunks),
        "answer": {
            "status": "found" if answers else "no_direct_answer",
            "sentences": answers,
            "note": "This command retrieves and quotes relevant source text; it does not call an LLM.",
        },
        "citations": citations,
    }


def qa_to_markdown(payload: dict[str, object]) -> str:
    answer = payload.get("answer") if isinstance(payload.get("answer"), dict) else {}
    sentences = answer.get("sentences") if isinstance(answer, dict) else []
    citations = payload.get("citations") if isinstance(payload.get("citations"), list) else []
    lines = [
        "# QA Report",
        "",
        f"- Question: {payload.get('question')}",
        f"- Method: `{payload.get('method')}`",
        f"- Status: `{answer.get('status') if isinstance(answer, dict) else None}`",
        f"- Chunk count: `{payload.get('chunk_count')}`",
        "",
        "## Extractive Answer",
        "",
    ]
    if isinstance(sentences, list) and sentences:
        for item in sentences:
            if isinstance(item, dict):
                lines.append(f"- {item.get('text')} [{item.get('citation')}]")
    else:
        lines.append("No direct answer was found from the retrieved chunks.")
    lines.extend(["", "## Citations", "", "| Ref | Page | Chunk | Score | Source | Snippet |", "|---:|---:|---:|---:|---|---|"])
    for item in citations:
        if isinstance(item, dict):
            snippet = str(item.get("snippet") or "").replace("\n", " ")[:220]
            lines.append(
                f"| {item.get('citation')} | {item.get('page') or ''} | {item.get('chunk_id') or ''} | "
                f"{item.get('score') or ''} | `{item.get('source_file') or ''}` | {snippet} |"
            )
    return "\n".join(lines) + "\n"


def write_qa_output(payload: dict[str, object], out_path: Path, output_format: str) -> list[Path]:
    if output_format == "all":
        out_path.mkdir(parents=True, exist_ok=True)
        targets = {"json": out_path / "qa_report.json", "md": out_path / "qa_report.md"}
    else:
        out_path.parent.mkdir(parents=True, exist_ok=True)
        targets = {output_format: out_path}

    written: list[Path] = []
    for suffix, target in targets.items():
        if suffix == "json":
            write_json_file(target, payload)
        elif suffix == "md":
            target.write_text(qa_to_markdown(payload), encoding="utf-8")
        written.append(target)
    return written


def diff_line_stats(left_text: str, right_text: str) -> dict[str, object]:
    left_lines = left_text.splitlines()
    right_lines = right_text.splitlines()
    matcher = difflib.SequenceMatcher(None, left_lines, right_lines, autojunk=True)
    stats = {"equal": 0, "replace": 0, "delete": 0, "insert": 0}
    changed_blocks: list[dict[str, object]] = []
    for tag, i1, i2, j1, j2 in matcher.get_opcodes():
        stats[tag] = stats.get(tag, 0) + max(i2 - i1, j2 - j1)
        if tag != "equal":
            changed_blocks.append(
                {
                    "type": tag,
                    "left_lines": [i1 + 1, i2],
                    "right_lines": [j1 + 1, j2],
                    "left_preview": "\n".join(left_lines[i1:min(i2, i1 + 3)])[:600],
                    "right_preview": "\n".join(right_lines[j1:min(j2, j1 + 3)])[:600],
                }
            )
    return {
        "left_line_count": len(left_lines),
        "right_line_count": len(right_lines),
        "stats": stats,
        "changed_block_count": len(changed_blocks),
        "changed_blocks": changed_blocks[:20],
    }


def limited_unified_diff(left_text: str, right_text: str, left_name: str, right_name: str, max_lines: int) -> list[str]:
    diff_iter = difflib.unified_diff(
        left_text.splitlines(),
        right_text.splitlines(),
        fromfile=left_name,
        tofile=right_name,
        lineterm="",
    )
    lines: list[str] = []
    for index, line in enumerate(diff_iter):
        if index >= max_lines:
            lines.append(f"... diff truncated after {max_lines} lines ...")
            break
        lines.append(line)
    return lines


def scalar_for_diff(value: object) -> str:
    if isinstance(value, float):
        return f"{value:.2f}"
    if value is None:
        return ""
    return str(value)


def diff_invoice_fields(left_fields: dict[str, object], right_fields: dict[str, object]) -> list[dict[str, object]]:
    keys = INVOICE_SUMMARY_COLUMNS[4:16]
    changes: list[dict[str, object]] = []
    for key in keys:
        left_value = scalar_for_diff(left_fields.get(key))
        right_value = scalar_for_diff(right_fields.get(key))
        if left_value != right_value:
            changes.append({"field": key, "left": left_value, "right": right_value})
    left_items = left_fields.get("items") if isinstance(left_fields.get("items"), list) else []
    right_items = right_fields.get("items") if isinstance(right_fields.get("items"), list) else []
    if len(left_items) != len(right_items):
        changes.append({"field": "items_count", "left": len(left_items), "right": len(right_items)})
    return changes


def build_diff_payload(
    left_path: Path,
    right_path: Path,
    parser_name: str,
    start_page: int,
    max_pages: int | None,
    profile: str,
    max_diff_lines: int,
) -> dict[str, object]:
    left_text, left_result = extract_best_text(left_path, parser_name, start_page, max_pages)
    right_text, right_result = extract_best_text(right_path, parser_name, start_page, max_pages)
    left_metadata = file_metadata(left_path)
    right_metadata = file_metadata(right_path)
    left_classification = classify_document(left_text, left_metadata)
    right_classification = classify_document(right_text, right_metadata)

    field_changes: list[dict[str, object]] = []
    should_extract_invoice = profile == "invoice" or (
        profile == "auto"
        and (left_classification.get("profile") == "invoice" or right_classification.get("profile") == "invoice")
    )
    if should_extract_invoice:
        left_fields = extract_invoice_fields_from_text(left_text, left_path, left_result.parser, left_result, strict=True)
        right_fields = extract_invoice_fields_from_text(right_text, right_path, right_result.parser, right_result, strict=True)
        field_changes = diff_invoice_fields(left_fields, right_fields)

    return {
        "version": __version__,
        "type": "diff-docs",
        "left": {"file": str(left_path), "parser": left_result.parser, "quality": asdict(left_result), "classification": left_classification},
        "right": {"file": str(right_path), "parser": right_result.parser, "quality": asdict(right_result), "classification": right_classification},
        "similarity": {
            "full_text_ratio": round(difflib.SequenceMatcher(None, left_text, right_text, autojunk=True).ratio(), 4),
            "first_50000_chars": similarity_score(left_text, right_text, 50000),
        },
        "line_diff": diff_line_stats(left_text, right_text),
        "unified_diff": limited_unified_diff(left_text, right_text, left_path.name, right_path.name, max_diff_lines),
        "field_changes": field_changes,
    }


def diff_to_markdown(payload: dict[str, object]) -> str:
    similarity = payload.get("similarity") if isinstance(payload.get("similarity"), dict) else {}
    line_diff = payload.get("line_diff") if isinstance(payload.get("line_diff"), dict) else {}
    stats = line_diff.get("stats") if isinstance(line_diff.get("stats"), dict) else {}
    field_changes = payload.get("field_changes") if isinstance(payload.get("field_changes"), list) else []
    unified = payload.get("unified_diff") if isinstance(payload.get("unified_diff"), list) else []
    lines = [
        "# Document Diff Report",
        "",
        f"- Left: `{payload.get('left', {}).get('file') if isinstance(payload.get('left'), dict) else ''}`",
        f"- Right: `{payload.get('right', {}).get('file') if isinstance(payload.get('right'), dict) else ''}`",
        f"- Full text similarity: `{similarity.get('full_text_ratio')}`",
        f"- First 50000 chars similarity: `{similarity.get('first_50000_chars')}`",
        f"- Changed blocks: `{line_diff.get('changed_block_count')}`",
        f"- Line stats: `{stats}`",
        "",
        "## Field Changes",
        "",
    ]
    if field_changes:
        lines.extend(["| Field | Left | Right |", "|---|---|---|"])
        for change in field_changes:
            if isinstance(change, dict):
                lines.append(f"| `{change.get('field')}` | {change.get('left') or ''} | {change.get('right') or ''} |")
    else:
        lines.append("No structured field changes detected.")
    lines.extend(["", "## Unified Diff", "", "```diff"])
    lines.extend(str(line) for line in unified)
    lines.extend(["```", ""])
    return "\n".join(lines)


def write_diff_output(payload: dict[str, object], out_path: Path, output_format: str) -> list[Path]:
    if output_format == "all":
        out_path.mkdir(parents=True, exist_ok=True)
        targets = {"json": out_path / "diff_report.json", "md": out_path / "diff_report.md"}
    else:
        out_path.parent.mkdir(parents=True, exist_ok=True)
        targets = {output_format: out_path}

    written: list[Path] = []
    for suffix, target in targets.items():
        if suffix == "json":
            write_json_file(target, payload)
        elif suffix == "md":
            target.write_text(diff_to_markdown(payload), encoding="utf-8")
        written.append(target)
    return written


# ---------------------------------------------------------------------------
# 报告生成
# ---------------------------------------------------------------------------

def make_report(
    file_path: Path,
    actual_path: Path,
    out_dir: Path,
    fmt: str,
    max_pages: int | None,
    start_page: int,
    results: list[ParseResult],
    normalized_texts: dict[str, str],
    similarity_chars: int,
    output_format: str,
    ocr_fallback: bool = False,
    min_quality: float = 0.0,
    fail_on_bad: bool = False,
) -> tuple[Path, Path]:
    """生成 JSON 和 Markdown 对比报告，返回两个报告文件路径。"""
    ok_results = [r for r in results if r.status == "ok"]

    pairwise: list[dict[str, object]] = []
    for i, left in enumerate(ok_results):
        for right in ok_results[i + 1:]:
            left_text = normalized_texts.get(left.parser, "")
            right_text = normalized_texts.get(right.parser, "")
            pairwise.append({
                "left": left.parser,
                "right": right.parser,
                "similarity_first_chars": similarity_chars,
                "similarity": similarity_score(left_text, right_text, similarity_chars),
            })

    report = {
        "version": __version__,
        "format": fmt,
        "source_file": str(file_path),
        "actual_file": str(actual_path),
        "start_page": start_page,
        "max_pages": max_pages,
        "output_format": output_format,
        "generated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        "results": [asdict(r) for r in results],
        "pairwise_similarity": pairwise,
    }
    recommended = choose_recommended_result(results)
    diagnostics = build_diagnostics(results)
    ocr_notice = ocr_fallback_notice(ocr_fallback)
    if ocr_notice:
        diagnostics.append(ocr_notice)
    quality_gate = quality_gate_status(recommended, min_quality, fail_on_bad)
    report["recommendation"] = asdict(recommended) if recommended else None
    report["diagnostics"] = diagnostics
    report["quality_gate"] = quality_gate
    report["ocr_fallback_requested"] = ocr_fallback
    report["ocr_fallback_available"] = False
    json_path = out_dir / "compare_report.json"
    json_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")

    fmt_label = {"pdf": "PDF", "word": "Word", "ppt": "PPT", "excel": "Excel", "html": "HTML"}.get(fmt, fmt)
    md_lines = [
        f"# {fmt_label} 解析对比报告",
        "",
        f"- 版本：`{__version__}`",
        f"- 源文件：`{file_path}`",
        f"- 实际解析文件：`{actual_path}`",
        f"- 文件格式：`{fmt_label}`",
    ]
    if fmt == "pdf":
        md_lines.extend([
            f"- 起始页：`{start_page}`",
            f"- 页数限制：`{max_pages if max_pages is not None else '全量'}`",
        ])
    md_lines.extend([
        f"- 解析输出格式：`{output_format}`",
        f"- 输出目录：`{out_dir}`",
        "",
        "## 推荐",
        "",
    ])
    if recommended:
        md_lines.extend([
            f"- 推荐解析器：`{recommended.parser}`",
            f"- 质量评分：`{recommended.quality_score}` (`{recommended.quality_label}`)",
            f"- 非空白字符数：`{recommended.non_space_chars}`",
            f"- 耗时：`{recommended.seconds}s`",
            "",
        ])
    else:
        md_lines.extend(["- 未找到可推荐的成功解析结果。", ""])

    if quality_gate["enabled"]:
        md_lines.extend([
            "## 质量门禁",
            "",
            f"- 最低质量：`{quality_gate['min_quality']}`",
            f"- 失败标签拦截：`{quality_gate['fail_on_bad']}`",
            f"- 通过：`{quality_gate['passed']}`",
            f"- 原因：`{quality_gate['reason']}`",
            "",
        ])

    md_lines.extend(["## 质量诊断", ""])
    for item in diagnostics:
        md_lines.append(f"- {item}")

    md_lines.extend([
        "",
        "## 解析结果",
        "",
        "| 解析器 | 状态 | 质量 | 评分 | 耗时(s) | 字符数 | 非空白字符 | 行数 | 段落数 | 中文字符 | cid | 控制字符 | 输出 | 错误 |",
        "|---|---:|---|---:|---:|---:|---:|---:|---:|---:|---:|---:|---|---|",
    ])
    for r in results:
        output_name = Path(r.output_file).name if r.output_file else ""
        error = (r.error or "").replace("|", "\\|")
        md_lines.append(
            f"| {r.parser} | {r.status} | {r.quality_label} | {r.quality_score} | "
            f"{r.seconds} | {r.chars} | {r.non_space_chars} | {r.lines} | {r.paragraphs} | "
            f"{r.chinese_chars} | {r.cid_markers} | {r.control_chars} | {output_name} | {error} |"
        )

    md_lines.extend(["", "## 相似度", ""])
    if pairwise:
        md_lines.extend([
            f"> 仅比较每份结果的前 {similarity_chars} 个字符，适合快速判断抽取顺序和文本差异。",
            "",
            "| 左侧解析器 | 右侧解析器 | 相似度 |",
            "|---|---|---:|",
        ])
        for item in pairwise:
            md_lines.append(f"| {item['left']} | {item['right']} | {item['similarity']} |")
    else:
        md_lines.append("没有两个以上成功输出，无法计算相似度。")

    md_lines.extend(["", "## 预览", ""])
    for r in results:
        if r.status != "ok":
            continue
        md_lines.extend([
            f"### {r.parser}", "",
            "```text", r.preview, "```", "",
        ])

    md_path = out_dir / "compare_report.md"
    md_path.write_text("\n".join(md_lines), encoding="utf-8")

    return json_path, md_path


# ---------------------------------------------------------------------------
# 子命令：compare（多解析器对比）
# ---------------------------------------------------------------------------

def cmd_compare(args: argparse.Namespace) -> int:
    file_path = args.file.resolve()
    out_dir = args.out_dir.resolve() if args.out_dir else file_path.parent / f"{file_path.stem}_parse_output"
    max_pages = None if args.max_pages == 0 else args.max_pages
    start_page = args.start_page

    if not file_path.exists():
        print(f"错误：文件不存在：{file_path}", file=sys.stderr)
        return 1

    fmt = detect_format(file_path)
    if fmt is None:
        supported = ", ".join(sorted(EXTENSION_FORMAT_MAP.keys()))
        print(f"错误：不支持的文件格式：{file_path.suffix}。支持：{supported}", file=sys.stderr)
        return 1

    parser_names, unknown = validate_parser_selection(fmt, args.parsers)
    if unknown:
        print(f"警告：{fmt} 格式不支持以下解析器，将被跳过：{', '.join(unknown)}", file=sys.stderr)

    if not parser_names:
        print("错误：没有可用的解析器", file=sys.stderr)
        return 1

    fmt_label = {"pdf": "PDF", "word": "Word", "ppt": "PPT", "excel": "Excel", "html": "HTML"}.get(fmt, fmt)
    print(f"文件格式：{fmt_label}")
    print(f"解析器：{', '.join(parser_names)}")

    out_dir.mkdir(parents=True, exist_ok=True)

    results, normalized_texts, actual_path = collect_parser_outputs(
        file_path, fmt, parser_names, out_dir, args.output_format, start_page, max_pages, args.parallel, progress=True
    )
    json_path, md_path = make_report(
        file_path, actual_path, out_dir, fmt, max_pages, start_page,
        results, normalized_texts, args.similarity_chars, args.output_format, args.ocr_fallback,
        args.min_quality, args.fail_on_bad,
    )
    print(f"\n对比报告：{md_path}")
    print(f"机器可读报告：{json_path}")

    recommended = choose_recommended_result(results)
    gate = quality_gate_status(recommended, args.min_quality, args.fail_on_bad)
    if gate["enabled"] and not gate["passed"]:
        print(f"质量门禁未通过：{gate['reason']}（best={gate['best_quality_score']} {gate['best_quality_label']}）", file=sys.stderr)
        return 2
    return 0


# ---------------------------------------------------------------------------
# 子命令：vote（多解析器投票并输出客户交付文本）
# ---------------------------------------------------------------------------

def cmd_vote(args: argparse.Namespace) -> int:
    file_path = args.file.resolve()
    out_dir = args.out_dir.resolve() if args.out_dir else file_path.parent / f"{file_path.stem}_vote_output"
    max_pages = None if args.max_pages == 0 else args.max_pages
    start_page = args.start_page

    if not file_path.exists():
        print(f"错误：文件不存在：{file_path}", file=sys.stderr)
        return 1

    fmt = detect_format(file_path)
    if fmt is None:
        supported = ", ".join(sorted(EXTENSION_FORMAT_MAP.keys()))
        print(f"错误：不支持的文件格式：{file_path.suffix}。支持：{supported}", file=sys.stderr)
        return 1
    if fmt != "pdf":
        print("错误：vote 当前用于 PDF 多解析器投票。其它格式请使用 auto/convert。", file=sys.stderr)
        return 1

    parser_names, unknown = validate_parser_selection(fmt, args.parsers)
    if unknown:
        print(f"警告：{fmt} 格式不支持以下解析器，将被跳过：{', '.join(unknown)}", file=sys.stderr)
    if not parser_names:
        print("错误：没有可用的解析器", file=sys.stderr)
        return 1

    print("PDF 投票解析")
    print(f"输入：{file_path}")
    print(f"解析器：{', '.join(parser_names)}")
    print(f"Profile：{args.profile}")
    print(f"输出目录：{out_dir}")

    preflight_probe: dict[str, object] | None = None
    preflight_written: list[Path] = []
    health_cache_path = Path(args.health_cache).resolve() if args.health_cache else default_parser_health_cache_path(out_dir)
    use_health_cache = bool(args.parser_health_cache)
    result_cache_dir = Path(args.result_cache_dir).resolve() if args.result_cache_dir else default_parser_result_cache_dir(out_dir)
    use_result_cache = bool(args.result_cache)
    if args.probe_before_vote:
        probe_max_pages = None if args.probe_max_pages == 0 else args.probe_max_pages
        probe_out_dir = out_dir / "preflight_probe"
        print(f"预检体检：前 {args.probe_max_pages if args.probe_max_pages else '全量'} 页 -> {probe_out_dir}")
        preflight_probe, preflight_written = run_parser_probe(
            file_path,
            fmt,
            parser_names,
            probe_out_dir,
            start_page,
            probe_max_pages,
            args.probe_min_quality,
            args.fail_on_bad,
            args.parallel,
            keep_outputs=False,
            progress=True,
            timeout_seconds=args.timeout,
            health_cache_path=health_cache_path,
            use_health_cache=use_health_cache,
            result_cache_dir=result_cache_dir,
            use_result_cache=use_result_cache,
        )
        ready_parser_names = ready_parsers_from_probe(preflight_probe)
        if not ready_parser_names:
            print("错误：预检没有解析器通过运行时体检，终止投票。", file=sys.stderr)
            print(f"预检输出：{', '.join(str(path) for path in preflight_written)}")
            return 2
        skipped = [name for name in parser_names if name not in ready_parser_names]
        if skipped:
            print(f"预检剔除解析器：{', '.join(skipped)}")
        parser_names = ready_parser_names
        print(f"进入正式投票：{', '.join(parser_names)}")

    results, normalized_texts, actual_path = collect_parser_outputs(
        file_path,
        fmt,
        parser_names,
        out_dir,
        "md",
        start_page,
        max_pages,
        args.parallel,
        progress=True,
        timeout_seconds=args.timeout,
        health_cache_path=health_cache_path,
        use_health_cache=use_health_cache,
        result_cache_dir=result_cache_dir,
        use_result_cache=use_result_cache,
    )
    payload = build_vote_payload(
        file_path,
        actual_path,
        fmt,
        results,
        normalized_texts,
        args.similarity_chars,
        args.min_quality,
        args.fail_on_bad,
        args.profile,
        preflight_probe,
    )
    layout: dict[str, object] | None = None
    layout_written: list[Path] = []
    needs_layout = args.customer and (args.field_layout or args.field_evidence or args.review_html)
    if needs_layout:
        try:
            layout_max_pages = None if args.field_layout_max_pages == 0 else args.field_layout_max_pages
            layout = pdf_layout_json(file_path, start_page, layout_max_pages)
            layout_path = write_json_file(out_dir / "field_layout.json", layout)
            layout_written.append(layout_path)
            print(f"字段 layout：{layout_path}")
        except Exception as exc:
            print(f"警告：字段 layout 提取失败，customer 字段坐标将保持为空：{exc}", file=sys.stderr)
    written = write_vote_outputs(payload, normalized_texts, out_dir, args.format)
    written.extend(preflight_written)
    written.extend(layout_written)
    if args.customer:
        delivery = build_customer_delivery(payload, normalized_texts, layout, not args.no_field_fusion)
        if delivery and args.field_evidence:
            evidence_payload = write_field_evidence_pack(file_path, delivery, out_dir)
            evidence_outputs = evidence_payload.get("outputs") if isinstance(evidence_payload.get("outputs"), dict) else {}
            for path_value in evidence_outputs.values():
                if path_value:
                    written.append(Path(str(path_value)))
        if delivery and args.review_html:
            review_path = write_review_html(file_path, delivery, out_dir, None)
            written.append(review_path)
        if delivery:
            written.extend(write_customer_delivery_outputs(delivery, out_dir, args.format))
    winner = payload.get("winner") if isinstance(payload.get("winner"), dict) else None
    if winner:
        print(f"投票胜出：{winner.get('parser')}，得分：{winner.get('vote_score')}")
    else:
        print("没有可交付的解析结果", file=sys.stderr)
    print(f"输出：{', '.join(str(path) for path in written)}")

    gate = payload.get("quality_gate") if isinstance(payload.get("quality_gate"), dict) else {}
    if not winner:
        return 2
    if gate.get("enabled") and not gate.get("passed"):
        print(f"质量门禁未通过：{gate.get('reason')}", file=sys.stderr)
        return 2
    return 0


# ---------------------------------------------------------------------------
# 子命令：customer-pack（复杂 PDF 客户交付包）
# ---------------------------------------------------------------------------

def cmd_customer_pack(args: argparse.Namespace) -> int:
    file_path = args.file.resolve()
    out_dir = args.out_dir.resolve() if args.out_dir else file_path.parent / f"{file_path.stem}_customer_pack"
    max_pages = None if args.max_pages == 0 else args.max_pages
    probe_max_pages = None if args.probe_max_pages == 0 else args.probe_max_pages
    layout_max_pages = None if args.layout_max_pages == 0 else args.layout_max_pages

    if not file_path.exists():
        print(f"错误：文件不存在：{file_path}", file=sys.stderr)
        return 1
    if file_path.suffix.lower() != ".pdf":
        print("错误：customer-pack 当前用于 PDF。其它格式请使用 auto/convert/knowledge-pack。", file=sys.stderr)
        return 1

    parser_names, unknown = validate_parser_selection("pdf", args.parsers)
    if unknown:
        print(f"警告：PDF 不支持以下解析器，将被跳过：{', '.join(unknown)}", file=sys.stderr)
    if not parser_names:
        print("错误：没有可用的解析器", file=sys.stderr)
        return 1

    print("PDF 客户交付包")
    print(f"输入：{file_path}")
    print(f"输出目录：{out_dir}")
    print(f"解析器：{', '.join(parser_names)}")
    print(f"Profile：{args.profile}")

    try:
        health_cache_path = Path(args.health_cache).resolve() if args.health_cache else default_parser_health_cache_path(out_dir)
        result_cache_dir = Path(args.result_cache_dir).resolve() if args.result_cache_dir else default_parser_result_cache_dir(out_dir)
        manifest = write_customer_pdf_pack(
            file_path,
            out_dir,
            parser_names,
            args.start_page,
            max_pages,
            probe_max_pages,
            args.table_pages,
            layout_max_pages,
            args.profile,
            args.min_quality,
            args.fail_on_bad,
            args.similarity_chars,
            include_tables=not args.no_tables,
            include_layout=not args.no_layout,
            field_fusion=not args.no_field_fusion,
            parallel=args.parallel,
            timeout_seconds=args.timeout,
            health_cache_path=health_cache_path,
            use_health_cache=bool(args.parser_health_cache),
            result_cache_dir=result_cache_dir,
            use_result_cache=bool(args.result_cache),
            include_field_evidence=not args.no_field_evidence,
            include_review_html=not args.no_review_html,
        )
    except Exception as exc:
        print(f"错误：customer-pack 失败：{exc}", file=sys.stderr)
        return 1

    outputs = manifest.get("outputs") if isinstance(manifest.get("outputs"), dict) else {}
    vote = manifest.get("vote") if isinstance(manifest.get("vote"), dict) else {}
    winner = vote.get("winner") if isinstance(vote.get("winner"), dict) else {}
    tables = manifest.get("tables") if isinstance(manifest.get("tables"), dict) else {}
    print(f"交付包：{out_dir}")
    print(f"manifest：{outputs.get('manifest')}")
    print(f"README：{outputs.get('readme')}")
    print(f"最佳解析器：{winner.get('parser')}，得分：{winner.get('vote_score')}")
    print(f"表格数：{tables.get('table_count', 0)}")
    return 0


def cmd_batch_customer_pack(args: argparse.Namespace) -> int:
    """Build customer packs for every PDF in a directory and write a batch index."""
    input_dir = Path(args.dir).resolve()
    if not input_dir.exists() or not input_dir.is_dir():
        print(f"错误：目录不存在：{input_dir}", file=sys.stderr)
        return 1
    extensions = parse_extensions(args.ext)
    if ".pdf" not in extensions:
        extensions.add(".pdf")
    try:
        files = collect_input_files(input_dir, extensions, args.recursive)
    except Exception as exc:
        print(f"错误：收集文件失败：{exc}", file=sys.stderr)
        return 1
    files = [path for path in files if path.suffix.lower() == ".pdf"]
    if not files:
        print("错误：目录下没有找到 PDF 文件", file=sys.stderr)
        return 1

    parser_names, unknown = validate_parser_selection("pdf", args.parsers)
    if unknown:
        print(f"警告：PDF 不支持以下解析器，将被跳过：{', '.join(unknown)}", file=sys.stderr)
    if not parser_names:
        print("错误：没有可用的解析器", file=sys.stderr)
        return 1

    out_dir = args.out_dir.resolve() if args.out_dir else input_dir / "customer_packs"
    out_dir.mkdir(parents=True, exist_ok=True)
    max_pages = None if args.max_pages == 0 else args.max_pages
    probe_max_pages = None if args.probe_max_pages == 0 else args.probe_max_pages
    layout_max_pages = None if args.layout_max_pages == 0 else args.layout_max_pages
    health_cache_path = Path(args.health_cache).resolve() if args.health_cache else default_parser_health_cache_path(out_dir)
    result_cache_dir = Path(args.result_cache_dir).resolve() if args.result_cache_dir else default_parser_result_cache_dir(out_dir)

    print("批量 PDF 客户交付包")
    print(f"目录：{input_dir}")
    print(f"文件数：{len(files)}")
    print(f"输出目录：{out_dir}")

    manifests: list[dict[str, object]] = []
    failures: list[dict[str, object]] = []
    for file_path in files:
        print(f"customer-pack: {file_path.name}...", end=" ", flush=True)
        try:
            pack_dir = out_dir / file_path.stem
            manifest = write_customer_pdf_pack(
                file_path,
                pack_dir,
                parser_names,
                args.start_page,
                max_pages,
                probe_max_pages,
                args.table_pages,
                layout_max_pages,
                args.profile,
                args.min_quality,
                args.fail_on_bad,
                args.similarity_chars,
                include_tables=not args.no_tables,
                include_layout=not args.no_layout,
                field_fusion=not args.no_field_fusion,
                parallel=args.parallel,
                timeout_seconds=args.timeout,
                health_cache_path=health_cache_path,
                use_health_cache=bool(args.parser_health_cache),
                result_cache_dir=result_cache_dir,
                use_result_cache=bool(args.result_cache),
                include_field_evidence=not args.no_field_evidence,
                include_review_html=not args.no_review_html,
            )
            manifests.append(manifest)
            print("ok")
        except Exception as exc:
            failures.append({"source_file": str(file_path), "error": repr(exc)})
            print(f"FAIL {exc}")

    index = {
        "version": __version__,
        "type": "batch-customer-pack",
        "generated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        "input_dir": str(input_dir),
        "output_dir": str(out_dir),
        "file_count": len(files),
        "success_count": len(manifests),
        "failure_count": len(failures),
        "profile": args.profile,
        "manifests": [
            {
                "source_file": manifest.get("source_file"),
                "manifest": (manifest.get("outputs") if isinstance(manifest.get("outputs"), dict) else {}).get("manifest"),
                "readme": (manifest.get("outputs") if isinstance(manifest.get("outputs"), dict) else {}).get("readme"),
                "winner": (manifest.get("vote") if isinstance(manifest.get("vote"), dict) else {}).get("winner"),
                "tables": manifest.get("tables"),
            }
            for manifest in manifests
        ],
        "failures": failures,
    }
    index_path = write_json_file(out_dir / "index.json", index)
    md_lines = [
        "# Batch Customer Pack Index",
        "",
        f"- Input directory: `{input_dir}`",
        f"- Output directory: `{out_dir}`",
        f"- Files: `{len(files)}`",
        f"- Success: `{len(manifests)}`",
        f"- Failed: `{len(failures)}`",
        "",
        "| Source | Manifest | Winner | Vote | Tables |",
        "|---|---|---|---:|---:|",
    ]
    for item in index["manifests"]:
        if not isinstance(item, dict):
            continue
        winner = item.get("winner") if isinstance(item.get("winner"), dict) else {}
        tables = item.get("tables") if isinstance(item.get("tables"), dict) else {}
        md_lines.append(
            f"| {Path(str(item.get('source_file'))).name} | `{item.get('manifest')}` | "
            f"{winner.get('parser')} | {winner.get('vote_score')} | {tables.get('table_count', 0)} |"
        )
    if failures:
        md_lines.extend(["", "## Failures", ""])
        for failure in failures:
            md_lines.append(f"- `{failure.get('source_file')}`: {failure.get('error')}")
    index_md = out_dir / "index.md"
    index_md.write_text("\n".join(md_lines) + "\n", encoding="utf-8")
    print(f"索引：{index_path}")
    print(f"索引 Markdown：{index_md}")
    return 0 if not failures else 1


# ---------------------------------------------------------------------------
# 子命令：probe（真实文件解析器运行时体检）
# ---------------------------------------------------------------------------

def cmd_probe(args: argparse.Namespace) -> int:
    file_path = args.file.resolve()
    out_dir = args.out_dir.resolve() if args.out_dir else file_path.parent / f"{file_path.stem}_probe_output"
    max_pages = None if args.max_pages == 0 else args.max_pages
    start_page = args.start_page

    if not file_path.exists():
        print(f"错误：文件不存在：{file_path}", file=sys.stderr)
        return 1

    fmt = detect_format(file_path)
    if fmt is None:
        supported = ", ".join(sorted(EXTENSION_FORMAT_MAP.keys()))
        print(f"错误：不支持的文件格式：{file_path.suffix}。支持：{supported}", file=sys.stderr)
        return 1

    available_parsers = get_parsers_for_format(fmt)
    if args.parsers:
        parser_names = [resolve_parser_name(n.strip().lower()) for n in args.parsers.split(",") if n.strip()]
    else:
        parser_names = available_parsers
    unknown = [name for name in parser_names if name not in available_parsers]
    if unknown:
        print(f"警告：{fmt} 格式不支持以下解析器，将被跳过：{', '.join(unknown)}", file=sys.stderr)
    parser_names = [name for name in parser_names if name in available_parsers]
    if not parser_names:
        print("错误：没有可探测的解析器", file=sys.stderr)
        return 1

    print("解析器运行时体检")
    print(f"输入：{file_path}")
    print(f"格式：{fmt}")
    print(f"解析器：{', '.join(parser_names)}")
    print(f"输出目录：{out_dir}")

    health_cache_path = Path(args.health_cache).resolve() if args.health_cache else default_parser_health_cache_path(out_dir)
    result_cache_dir = Path(args.result_cache_dir).resolve() if args.result_cache_dir else default_parser_result_cache_dir(out_dir)
    payload, _written = run_parser_probe(
        file_path,
        fmt,
        parser_names,
        out_dir,
        start_page,
        max_pages,
        args.min_quality,
        args.fail_on_bad,
        args.parallel,
        keep_outputs=args.keep_outputs,
        progress=True,
        timeout_seconds=args.timeout,
        health_cache_path=health_cache_path,
        use_health_cache=bool(args.parser_health_cache),
        result_cache_dir=result_cache_dir,
        use_result_cache=bool(args.result_cache),
    )
    json_path, md_path = out_dir / "probe_report.json", out_dir / "probe_report.md"
    summary = payload.get("summary") if isinstance(payload.get("summary"), dict) else {}
    print(f"可用解析器：{summary.get('ready_count')}/{summary.get('parser_count')} -> {', '.join(summary.get('ready_parsers', []))}")
    print(f"体检报告：{md_path}")
    print(f"机器可读报告：{json_path}")

    if args.fail_on_bad and int(summary.get("ready_count") or 0) == 0:
        print("质量门禁未通过：没有解析器通过运行时体检", file=sys.stderr)
        return 2
    return 0


# ---------------------------------------------------------------------------
# 子命令：convert（文档转 md/txt/json）
# ---------------------------------------------------------------------------

def cmd_convert(args: argparse.Namespace) -> int:
    file_path = args.file.resolve()
    parser_name = resolve_parser_name(args.parser)

    if not file_path.exists():
        print(f"错误：文件不存在：{file_path}", file=sys.stderr)
        return 1

    fmt = detect_format(file_path)
    if fmt is None:
        supported = ", ".join(sorted(EXTENSION_FORMAT_MAP.keys()))
        print(f"错误：不支持的文件格式：{file_path.suffix}。支持：{supported}", file=sys.stderr)
        return 1

    available = get_parsers_for_format(fmt)
    if parser_name not in available:
        print(f"错误：{fmt} 格式不支持解析器 {parser_name}。可用：{', '.join(available)}", file=sys.stderr)
        return 1

    extract_func = get_extractor(fmt, parser_name)
    if extract_func is None:
        print(f"错误：未找到解析器 {parser_name}", file=sys.stderr)
        return 1

    output_format = args.format

    # 确定输出路径
    if args.output:
        out_path = Path(args.output).resolve()
    else:
        out_path = default_output_path(file_path, output_format)

    # PDF 特殊处理
    max_pages = None if args.max_pages == 0 else args.max_pages
    start_page = args.start_page

    print(f"解析器：{parser_name}")
    print(f"输入：{file_path}")
    print(f"输出格式：{output_format}")
    print(f"输出：{out_path}")
    if args.ocr_fallback:
        print("OCR fallback requested; use the `ocr` subcommand to run Tesseract OCR.", file=sys.stderr)

    start = time.perf_counter()
    temp_dir_obj: tempfile.TemporaryDirectory[str] | None = None
    try:
        target_path = file_path
        if fmt == "pdf" and parser_name in PDF_SUBSET_TARGET_PARSERS:
            temp_dir_obj = tempfile.TemporaryDirectory(prefix="parse_convert_")
            target_path, skip_reason = prepare_pdf_subset_target(
                file_path, start_page, max_pages, Path(temp_dir_obj.name)
            )
            if skip_reason:
                print(f"错误：{skip_reason}", file=sys.stderr)
                return 1

        raw_text = extract_func(target_path, start_page, max_pages)
        written = write_single_conversion_output(out_path, parser_name, file_path, raw_text, output_format)
        elapsed = time.perf_counter() - start
        chars = len(normalize_text(raw_text))
        print(f"完成：{chars} 字符，{elapsed:.3f}s")
        for path in written:
            print(f"输出文件：{path}")
        return 0
    except ImportError as exc:
        print(f"错误：{exc}", file=sys.stderr)
        return 1
    except Exception as exc:
        print(f"错误：{exc}", file=sys.stderr)
        return 1
    finally:
        if temp_dir_obj:
            temp_dir_obj.cleanup()


# ---------------------------------------------------------------------------
# 子命令：batch（批量转换）
# ---------------------------------------------------------------------------

def cmd_batch(args: argparse.Namespace) -> int:
    input_dir = Path(args.dir).resolve()
    parser_name = resolve_parser_name(args.parser)
    out_dir = Path(args.output_dir).resolve() if args.output_dir else input_dir / "batch_output"
    output_format = args.format
    extensions = parse_extensions(args.ext)

    if not input_dir.exists() or not input_dir.is_dir():
        print(f"错误：目录不存在：{input_dir}", file=sys.stderr)
        return 1

    # 收集文件
    files: list[Path] = []
    for ext in sorted(extensions):
        files.extend(input_dir.glob(f"*{ext}"))
    files.sort()

    if not files:
        print(f"错误：目录下没有找到 {', '.join(extensions)} 文件", file=sys.stderr)
        return 1

    print(f"输入目录：{input_dir}")
    print(f"输出目录：{out_dir}")
    print(f"输出格式：{output_format}")
    print(f"解析器：{parser_name}")
    print(f"文件数：{len(files)}")

    out_dir.mkdir(parents=True, exist_ok=True)

    # 处理单个文件
    def process_one(file_path: Path) -> tuple[str, str, bool]:
        fmt = detect_format(file_path)
        if fmt is None:
            return file_path.name, f"不支持的格式: {file_path.suffix}", False

        available = get_parsers_for_format(fmt)
        if parser_name not in available:
            return file_path.name, f"解析器 {parser_name} 不支持 {fmt}", False

        extract_func = get_extractor(fmt, parser_name)
        if extract_func is None:
            return file_path.name, f"未找到解析器", False

        try:
            raw_text = extract_func(file_path, 1, None)
            normalized = normalize_text(raw_text)
            if output_format == "all":
                target = out_dir / file_path.stem
            else:
                target = out_dir / f"{file_path.stem}.{output_format}"
            written = write_single_conversion_output(target, parser_name, file_path, raw_text, output_format)
            return file_path.name, f"{len(normalized)} chars -> {', '.join(p.name for p in written)}", True
        except Exception as exc:
            return file_path.name, repr(exc), False

    # 执行
    success = 0
    fail = 0
    if args.parallel:
        with ThreadPoolExecutor(max_workers=min(len(files), 4)) as pool:
            futures = {pool.submit(process_one, f): f for f in files}
            for future in as_completed(futures):
                name, msg, ok = future.result()
                icon = "OK" if ok else "FAIL"
                print(f"  [{icon}] {name}: {msg}")
                if ok:
                    success += 1
                else:
                    fail += 1
    else:
        for f in files:
            name, msg, ok = process_one(f)
            icon = "OK" if ok else "FAIL"
            print(f"  [{icon}] {name}: {msg}")
            if ok:
                success += 1
            else:
                fail += 1

    print(f"\n完成：{success} 成功，{fail} 失败")
    print(f"输出目录：{out_dir}")
    return 0 if fail == 0 else 1


# ---------------------------------------------------------------------------
# 子命令：scan-dir（批量质量扫描）
# ---------------------------------------------------------------------------

def choose_scan_parser(fmt: str, requested_parser: str) -> str | None:
    """Choose a parser for quality scan; prefer fast native parsers for PDF."""
    available = get_parsers_for_format(fmt)
    if not available:
        return None
    requested = resolve_parser_name(requested_parser.strip().lower())
    if requested != "auto":
        return requested if requested in available else None

    preferences = {
        "pdf": ["pymupdf", "pdfplumber", "pypdf", "pdfminer", "pymupdf4llm", "markitdown", "liteparse"],
        "word": ["markitdown", "python-docx"],
        "ppt": ["markitdown", "python-pptx"],
        "excel": ["markitdown", "openpyxl"],
        "html": ["markitdown", "beautifulsoup4"],
    }
    for parser_name in preferences.get(fmt, available):
        if parser_name in available and parser_available(parser_name):
            return parser_name
    return available[0]


def scan_one_document(
    file_path: Path,
    requested_parser: str,
    start_page: int,
    max_pages: int | None,
    min_quality: float,
    fail_on_bad: bool,
) -> dict[str, object]:
    """Extract a small sample and return a quality row for scan-dir."""
    fmt = detect_format(file_path)
    if fmt is None:
        return {
            "source_file": str(file_path),
            "file_name": file_path.name,
            "format": None,
            "parser": None,
            "status": "skipped",
            "seconds": 0.0,
            "quality_score": 0.0,
            "quality_label": "unknown",
            "chars": 0,
            "non_space_chars": 0,
            "lines": 0,
            "chinese_chars": 0,
            "cid_markers": 0,
            "control_chars": 0,
            "needs_ocr": False,
            "gate_passed": False,
            "gate_reason": "unsupported_format",
            "error": f"不支持的格式：{file_path.suffix}",
        }

    parser_name = choose_scan_parser(fmt, requested_parser)
    if parser_name is None:
        return {
            "source_file": str(file_path),
            "file_name": file_path.name,
            "format": fmt,
            "parser": None,
            "status": "failed",
            "seconds": 0.0,
            "quality_score": 0.0,
            "quality_label": "unknown",
            "chars": 0,
            "non_space_chars": 0,
            "lines": 0,
            "chinese_chars": 0,
            "cid_markers": 0,
            "control_chars": 0,
            "needs_ocr": fmt == "pdf",
            "gate_passed": False,
            "gate_reason": "no_parser",
            "error": f"没有可用于 {fmt} 的解析器",
        }

    extract_func = get_extractor(fmt, parser_name)
    if extract_func is None:
        return {
            "source_file": str(file_path),
            "file_name": file_path.name,
            "format": fmt,
            "parser": parser_name,
            "status": "failed",
            "seconds": 0.0,
            "quality_score": 0.0,
            "quality_label": "unknown",
            "chars": 0,
            "non_space_chars": 0,
            "lines": 0,
            "chinese_chars": 0,
            "cid_markers": 0,
            "control_chars": 0,
            "needs_ocr": fmt == "pdf",
            "gate_passed": False,
            "gate_reason": "missing_extractor",
            "error": f"未找到解析器 {parser_name}",
        }

    temp_dir_obj: tempfile.TemporaryDirectory[str] | None = None
    start = time.perf_counter()
    try:
        target_path = file_path
        if fmt == "pdf" and parser_name in PDF_SUBSET_TARGET_PARSERS and max_pages is not None:
            temp_dir_obj = tempfile.TemporaryDirectory(prefix="parse_scan_")
            target_path, skip_reason = prepare_pdf_subset_target(
                file_path, start_page, max_pages, Path(temp_dir_obj.name)
            )
            if skip_reason:
                raise RuntimeError(skip_reason)
        raw_text = extract_func(target_path, start_page, max_pages)
        normalized = normalize_text(raw_text)
        result = text_stats(parser_name, "ok", time.perf_counter() - start, normalized, None)
        gate = quality_gate_status(result, min_quality, fail_on_bad)
        needs_ocr = (
            fmt == "pdf"
            and (
                result.quality_label in QUALITY_BAD_LABELS
                or (min_quality > 0 and result.quality_score < min_quality)
            )
        )
        return {
            "source_file": str(file_path),
            "file_name": file_path.name,
            "format": fmt,
            "parser": parser_name,
            "status": result.status,
            "seconds": result.seconds,
            "quality_score": result.quality_score,
            "quality_label": result.quality_label,
            "chars": result.chars,
            "non_space_chars": result.non_space_chars,
            "lines": result.lines,
            "chinese_chars": result.chinese_chars,
            "cid_markers": result.cid_markers,
            "control_chars": result.control_chars,
            "needs_ocr": needs_ocr,
            "gate_passed": gate["passed"],
            "gate_reason": gate["reason"],
            "error": None,
        }
    except Exception as exc:
        gate = quality_gate_status(None, min_quality, fail_on_bad)
        return {
            "source_file": str(file_path),
            "file_name": file_path.name,
            "format": fmt,
            "parser": parser_name,
            "status": "failed",
            "seconds": round(time.perf_counter() - start, 3),
            "quality_score": 0.0,
            "quality_label": "failed",
            "chars": 0,
            "non_space_chars": 0,
            "lines": 0,
            "chinese_chars": 0,
            "cid_markers": 0,
            "control_chars": 0,
            "needs_ocr": fmt == "pdf",
            "gate_passed": gate["passed"],
            "gate_reason": gate["reason"],
            "error": repr(exc),
        }
    finally:
        if temp_dir_obj:
            temp_dir_obj.cleanup()


def write_scan_reports(
    rows: list[dict[str, object]],
    out_dir: Path,
    input_dir: Path,
    requested_parser: str,
    extensions: set[str],
    recursive: bool,
    min_quality: float,
    fail_on_bad: bool,
) -> tuple[Path, Path, Path]:
    """Write scan-dir JSON, Markdown, and CSV reports."""
    out_dir.mkdir(parents=True, exist_ok=True)
    summary = {
        "version": __version__,
        "input_dir": str(input_dir),
        "extensions": sorted(extensions),
        "recursive": recursive,
        "requested_parser": requested_parser,
        "min_quality": min_quality,
        "fail_on_bad": fail_on_bad,
        "file_count": len(rows),
        "ok_count": sum(1 for row in rows if row["status"] == "ok"),
        "failed_count": sum(1 for row in rows if row["status"] == "failed"),
        "needs_ocr_count": sum(1 for row in rows if row["needs_ocr"]),
        "gate_failed_count": sum(1 for row in rows if not row["gate_passed"]),
        "generated_at": time.strftime("%Y-%m-%d %H:%M:%S"),
    }

    json_path = out_dir / "scan_report.json"
    json_path.write_text(
        json.dumps({"summary": summary, "files": rows}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    csv_path = out_dir / "scan_report.csv"
    fieldnames = [
        "source_file", "file_name", "format", "parser", "status", "quality_label",
        "quality_score", "chars", "non_space_chars", "lines", "chinese_chars",
        "cid_markers", "control_chars", "needs_ocr", "gate_passed", "gate_reason",
        "seconds", "error",
    ]
    with csv_path.open("w", encoding="utf-8-sig", newline="") as fh:
        writer = csv.DictWriter(fh, fieldnames=fieldnames)
        writer.writeheader()
        for row in rows:
            writer.writerow({key: row.get(key) for key in fieldnames})

    md_path = out_dir / "scan_report.md"
    md_lines = [
        "# Directory Quality Scan",
        "",
        f"- Version: `{__version__}`",
        f"- Input directory: `{input_dir}`",
        f"- Extensions: `{', '.join(sorted(extensions))}`",
        f"- Recursive: `{recursive}`",
        f"- Requested parser: `{requested_parser}`",
        f"- Minimum quality: `{min_quality}`",
        f"- Fail on bad: `{fail_on_bad}`",
        f"- Files: `{summary['file_count']}`",
        f"- OK: `{summary['ok_count']}`",
        f"- Failed: `{summary['failed_count']}`",
        f"- Needs OCR: `{summary['needs_ocr_count']}`",
        f"- Gate failed: `{summary['gate_failed_count']}`",
        "",
        "| File | Format | Parser | Status | Quality | Score | Chars | Needs OCR | Gate | Error |",
        "|---|---|---|---|---|---:|---:|---|---|---|",
    ]
    for row in rows:
        error = str(row.get("error") or "").replace("|", "\\|")
        md_lines.append(
            f"| {row['file_name']} | {row['format']} | {row['parser']} | {row['status']} | "
            f"{row['quality_label']} | {row['quality_score']} | {row['chars']} | "
            f"{row['needs_ocr']} | {row['gate_passed']} | {error} |"
        )
    md_path.write_text("\n".join(md_lines), encoding="utf-8")
    return json_path, md_path, csv_path


def cmd_scan_dir(args: argparse.Namespace) -> int:
    input_dir = Path(args.dir).resolve()
    out_dir = Path(args.out_dir).resolve() if args.out_dir else input_dir / "scan_output"
    extensions = parse_extensions(args.ext)
    max_pages = None if args.max_pages == 0 else args.max_pages

    if not input_dir.exists() or not input_dir.is_dir():
        print(f"错误：目录不存在：{input_dir}", file=sys.stderr)
        return 1

    files: list[Path] = []
    for ext in sorted(extensions):
        iterator = input_dir.rglob(f"*{ext}") if args.recursive else input_dir.glob(f"*{ext}")
        files.extend(path for path in iterator if path.is_file())
    files = sorted(set(files))

    if not files:
        print(f"错误：目录下没有找到 {', '.join(sorted(extensions))} 文件", file=sys.stderr)
        return 1

    print(f"扫描目录：{input_dir}")
    print(f"文件数：{len(files)}")
    print(f"解析器：{args.parser}")
    print(f"输出目录：{out_dir}")

    rows: list[dict[str, object]] = []
    for file_path in files:
        print(f"扫描: {file_path.name}...", end=" ", flush=True)
        row = scan_one_document(
            file_path,
            args.parser,
            args.start_page,
            max_pages,
            args.min_quality,
            args.fail_on_bad,
        )
        rows.append(row)
        print(f"{row['status']} {row['quality_label']} {row['quality_score']}")

    json_path, md_path, csv_path = write_scan_reports(
        rows,
        out_dir,
        input_dir,
        args.parser,
        extensions,
        args.recursive,
        args.min_quality,
        args.fail_on_bad,
    )
    print(f"扫描报告：{md_path}")
    print(f"机器可读报告：{json_path}")
    print(f"CSV：{csv_path}")

    gate_failed = [row for row in rows if not row["gate_passed"]]
    if args.fail_on_bad and gate_failed:
        print(f"质量门禁未通过：{len(gate_failed)} 个文件", file=sys.stderr)
        return 2
    return 0


# ---------------------------------------------------------------------------
# 子命令：tables（表格提取）
# ---------------------------------------------------------------------------

def cmd_tables(args: argparse.Namespace) -> int:
    file_path = args.file.resolve()

    if not file_path.exists():
        print(f"错误：文件不存在：{file_path}", file=sys.stderr)
        return 1

    if file_path.suffix.lower() != ".pdf":
        print("错误：表格提取仅支持 PDF 文件", file=sys.stderr)
        return 1

    # 获取页数
    page_count = get_pdf_page_count(file_path)
    print(f"PDF 总页数：{page_count}")

    # 解析页码
    if args.pages:
        page_indices = parse_page_spec(args.pages, page_count)
    else:
        page_indices = list(range(page_count))

    print(f"提取范围：第 {page_indices[0] + 1} 页 ~ 第 {page_indices[-1] + 1} 页（共 {len(page_indices)} 页）")

    # 提取
    start = time.perf_counter()
    try:
        tables = extract_tables_from_pdf(file_path, page_indices)
    except ImportError as exc:
        print(f"错误：{exc}", file=sys.stderr)
        return 1
    except Exception as exc:
        print(f"错误：提取失败：{exc}", file=sys.stderr)
        return 1

    elapsed = time.perf_counter() - start
    print(f"提取到 {len(tables)} 个表格，耗时 {elapsed:.3f}s")

    if not tables:
        print("未找到任何表格")
        return 0

    # 输出
    fmt = args.format or "md"
    if fmt == "all":
        out_dir = Path(args.output).resolve() if args.output else file_path.parent / f"{file_path.stem}_tables"
        out_dir.mkdir(parents=True, exist_ok=True)
        outputs = {
            "md": out_dir / f"{file_path.stem}_tables.md",
            "csv": out_dir / f"{file_path.stem}_tables.csv",
            "json": out_dir / f"{file_path.stem}_tables.json",
        }
    else:
        suffix = fmt
        if args.output:
            outputs = {suffix: Path(args.output).resolve()}
        else:
            outputs = {suffix: file_path.parent / f"{file_path.stem}_tables.{suffix}"}

    for suffix, out_path in outputs.items():
        if suffix == "csv":
            content = tables_to_csv(tables)
        elif suffix == "json":
            content = tables_to_json(tables)
        else:
            content = tables_to_markdown(tables)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        out_path.write_text(content, encoding="utf-8")
        print(f"输出：{out_path}")
    return 0


def cmd_table_vote(args: argparse.Namespace) -> int:
    """Run multiple table parsers and vote on the best table extraction."""
    file_path = args.file.resolve()
    if not file_path.exists():
        print(f"错误：文件不存在：{file_path}", file=sys.stderr)
        return 1
    if file_path.suffix.lower() != ".pdf":
        print("错误：table-vote 仅支持 PDF 文件", file=sys.stderr)
        return 1
    page_count = get_pdf_page_count(file_path)
    page_indices = resolve_page_indices(args.pages, page_count, default_all=True)
    methods = [item.strip() for item in args.methods.split(",") if item.strip()]
    if not methods:
        print("错误：没有表格解析器", file=sys.stderr)
        return 1
    out_dir = Path(args.out_dir).resolve() if args.out_dir else file_path.parent / f"{file_path.stem}_table_vote"
    print("PDF 表格投票")
    print(f"输入：{file_path}")
    print(f"页数：{len(page_indices)}")
    print(f"表格解析器：{', '.join(methods)}")
    try:
        payload = build_table_vote_payload(file_path, page_indices, methods)
        written = write_table_vote_outputs(payload, out_dir, args.format)
    except Exception as exc:
        print(f"错误：table-vote 失败：{exc}", file=sys.stderr)
        return 1
    winner = payload.get("winner") if isinstance(payload.get("winner"), dict) else None
    if winner:
        print(f"表格胜出：{winner.get('method')}，得分：{winner.get('vote_score')}，表格数：{winner.get('table_count')}")
    else:
        print("未抽取到可交付表格")
    print(f"输出：{', '.join(str(path) for path in written)}")
    return 0 if winner else 2


# ---------------------------------------------------------------------------
# 子命令：doctor（依赖自检）
# ---------------------------------------------------------------------------

def cmd_doctor(args: argparse.Namespace) -> int:
    rows = parser_dependency_rows(args.format)
    missing_packages = sorted(
        {
            str(row["package"])
            for row in rows
            if not row["available"] and row.get("kind") == "python"
        }
    )
    missing_commands = sorted(
        {
            str(row["package"])
            for row in rows
            if not row["available"] and row.get("kind") == "command"
        }
    )
    missing_ocr_parser_packages = sorted(
        {
            str(row["package"])
            for row in rows
            if not row["available"] and row.get("kind") == "ocr"
        }
    )
    ocr_rows = OCR_DEPENDENCIES if getattr(args, "ocr", False) else []
    missing_ocr_python = [
        str(row["package"])
        for row in ocr_rows
        if not row["available"] and row["kind"] == "python"
    ]
    missing_ocr_system = [
        str(row["package"])
        for row in ocr_rows
        if not row["available"] and row["kind"] == "system"
    ]
    opendataloader_rows = OPENLOADER_DEPENDENCIES if getattr(args, "opendataloader", False) else []
    missing_opendataloader_python = [
        str(row["package"])
        for row in opendataloader_rows
        if not row["available"] and row["kind"] == "python"
    ]
    missing_opendataloader_system = [
        str(row["package"])
        for row in opendataloader_rows
        if not row["available"] and row["kind"] == "system"
    ]

    if args.json:
        payload = {
            "version": __version__,
            "python": sys.version.split()[0],
            "format": args.format,
            "rows": rows,
            "missing_packages": missing_packages,
            "missing_commands": missing_commands,
            "missing_ocr_parser_packages": missing_ocr_parser_packages,
            "ocr_rows": ocr_rows,
            "missing_ocr_python_packages": missing_ocr_python,
            "missing_ocr_system_packages": missing_ocr_system,
            "opendataloader_rows": opendataloader_rows,
            "missing_opendataloader_python_packages": missing_opendataloader_python,
            "missing_opendataloader_system_packages": missing_opendataloader_system,
            "install_command": f"pip install {' '.join(missing_packages)}" if missing_packages else None,
            "command_install_hint": command_install_hint(missing_commands),
            "ocr_install_hint": ocr_install_hint(missing_ocr_python, missing_ocr_system),
            "opendataloader_install_hint": opendataloader_install_hint(missing_opendataloader_python, missing_opendataloader_system),
        }
        print(json.dumps(payload, ensure_ascii=False, indent=2))
        strict_missing = (
            missing_packages
            or missing_commands
            or missing_ocr_parser_packages
            or missing_ocr_python
            or missing_ocr_system
            or missing_opendataloader_python
            or missing_opendataloader_system
        )
        return 1 if strict_missing and args.strict else 0

    print(f"parse_pdf_compare {__version__}")
    print(f"Python: {sys.version.split()[0]}")
    print("")
    print("| 格式 | 解析器 | 状态 | 模块 | 安装包 |")
    print("|---|---|---|---|---|")
    for row in rows:
        modules = ", ".join(str(module) for module in row["modules"]) or ", ".join(str(command) for command in row.get("commands", []))
        print(f"| {row['format']} | {row['parser']} | {row['status']} | {modules} | {row['package']} |")

    if missing_packages:
        print("")
        print("缺失依赖安装建议：")
        print(f"pip install {' '.join(missing_packages)}")

    if missing_commands:
        print("")
        print("外部命令安装建议：")
        print(command_install_hint(missing_commands))

    if missing_ocr_parser_packages:
        print("")
        print("OCR parser 安装建议：")
        print("pip install pymupdf pytesseract pillow")
        print("并安装系统 Tesseract OCR、中文语言包（如 chi_sim），确保 tesseract 在 PATH 中。")

    if ocr_rows:
        print("")
        print("| OCR 依赖 | 类型 | 状态 | 安装包 |")
        print("|---|---|---|---|")
        for row in ocr_rows:
            status = "ok" if row["available"] else "missing"
            print(f"| {row['name']} | {row['kind']} | {status} | {row['package']} |")
        if missing_ocr_python or missing_ocr_system:
            print("")
            print("OCR 依赖安装建议：")
            if missing_ocr_python:
                print(f"pip install {' '.join(missing_ocr_python)}")
            if missing_ocr_system:
                print("安装系统 Tesseract OCR，并确保 tesseract 在 PATH 中；中文 OCR 还需要 chi_sim 等语言数据。")

    if opendataloader_rows:
        print("")
        print("| opendataloader-pdf 依赖 | 类型 | 状态 | 安装包 |")
        print("|---|---|---|---|")
        for row in opendataloader_rows:
            status = "ok" if row["available"] else "missing"
            print(f"| {row['name']} | {row['kind']} | {status} | {row['package']} |")
        if missing_opendataloader_python or missing_opendataloader_system:
            print("")
            print("opendataloader-pdf 依赖安装建议：")
            if missing_opendataloader_python:
                print(f"pip install {' '.join(missing_opendataloader_python)}")
            if missing_opendataloader_system:
                print("需要系统安装 Java (JDK/JRE)，并确保 java 命令在 PATH 中。")

    if missing_packages or missing_commands or missing_ocr_parser_packages or missing_ocr_python or missing_ocr_system or missing_opendataloader_python or missing_opendataloader_system:
        return 1 if args.strict else 0

    print("")
    print("所有检查范围内的解析器依赖均可导入。")
    return 0


# ---------------------------------------------------------------------------
# 子命令：metadata（元数据提取）
# ---------------------------------------------------------------------------

def cmd_metadata(args: argparse.Namespace) -> int:
    file_path = args.file.resolve()
    if not file_path.exists():
        print(f"错误：文件不存在：{file_path}", file=sys.stderr)
        return 1

    try:
        metadata = file_metadata(file_path)
    except Exception as exc:
        print(f"错误：元数据读取失败：{exc}", file=sys.stderr)
        return 1

    output_format = args.format
    if args.output:
        out_path = Path(args.output).resolve()
    else:
        out_path = file_path.parent / f"{file_path.stem}_metadata.{output_format}"

    out_path.parent.mkdir(parents=True, exist_ok=True)
    if output_format == "json":
        out_path.write_text(json.dumps(metadata, ensure_ascii=False, indent=2), encoding="utf-8")
    else:
        out_path.write_text(metadata_to_markdown(metadata), encoding="utf-8")
    print(f"输出：{out_path}")
    return 0


# ---------------------------------------------------------------------------
# 子命令：chunk（分块输出）
# ---------------------------------------------------------------------------

def cmd_chunk(args: argparse.Namespace) -> int:
    file_path = args.file.resolve()
    parser_name = resolve_parser_name(args.parser)

    if not file_path.exists():
        print(f"错误：文件不存在：{file_path}", file=sys.stderr)
        return 1

    fmt = detect_format(file_path)
    if fmt is None:
        supported = ", ".join(sorted(EXTENSION_FORMAT_MAP.keys()))
        print(f"错误：不支持的文件格式：{file_path.suffix}。支持：{supported}", file=sys.stderr)
        return 1

    if parser_name not in get_parsers_for_format(fmt):
        print(f"错误：{fmt} 格式不支持解析器 {parser_name}", file=sys.stderr)
        return 1

    extract_func = get_extractor(fmt, parser_name)
    if extract_func is None:
        print(f"错误：未找到解析器 {parser_name}", file=sys.stderr)
        return 1

    max_pages = None if args.max_pages == 0 else args.max_pages
    start_page = args.start_page
    temp_dir_obj: tempfile.TemporaryDirectory[str] | None = None

    if args.ocr_fallback:
        print("OCR fallback requested; use the `ocr` subcommand to run Tesseract OCR.", file=sys.stderr)

    try:
        target_path = file_path
        if fmt == "pdf" and parser_name in PDF_SUBSET_TARGET_PARSERS:
            temp_dir_obj = tempfile.TemporaryDirectory(prefix="parse_chunk_")
            target_path, skip_reason = prepare_pdf_subset_target(
                file_path, start_page, max_pages, Path(temp_dir_obj.name)
            )
            if skip_reason:
                print(f"错误：{skip_reason}", file=sys.stderr)
                return 1

        if args.chunk_by == "page":
            if fmt != "pdf":
                print("错误：--chunk-by page 仅支持 PDF", file=sys.stderr)
                return 1
            page_extractor = get_pdf_page_extractor(parser_name)
            if page_extractor is None:
                supported = ", ".join(sorted(PAGE_CHUNK_PARSERS))
                print(f"错误：解析器 {parser_name} 不支持分页分块。可用：{supported}", file=sys.stderr)
                return 1
            pages = page_extractor(file_path, start_page, max_pages)
            chunks = split_pages_into_chunks(pages)
        else:
            raw_text = extract_func(target_path, start_page, max_pages)
            chunks = split_text_into_chunks(raw_text, args.chunk_size, args.overlap)
    except Exception as exc:
        print(f"错误：分块失败：{exc}", file=sys.stderr)
        return 1
    finally:
        if temp_dir_obj:
            temp_dir_obj.cleanup()

    if args.output:
        out_path = Path(args.output).resolve()
    else:
        suffix = "jsonl" if args.format == "jsonl" else args.format
        out_path = file_path.parent / f"{file_path.stem}_chunks.{suffix}"
        if args.format == "all":
            out_path = file_path.parent / f"{file_path.stem}_chunks"

    meta = {
        "version": __version__,
        "source_file": str(file_path),
        "parser": parser_name,
        "format": fmt,
        "chunk_by": args.chunk_by,
        "chunk_size": args.chunk_size,
        "overlap": args.overlap,
        "ocr_fallback_requested": args.ocr_fallback,
        "ocr_fallback_available": False,
    }
    written = write_chunks(chunks, out_path, args.format, meta)
    print(f"分块数：{len(chunks)}")
    for path in written:
        print(f"输出：{path}")
    return 0


# ---------------------------------------------------------------------------
# 子命令：render-pages（PDF 页面截图）
# ---------------------------------------------------------------------------

def cmd_render_pages(args: argparse.Namespace) -> int:
    file_path = args.file.resolve()
    if not file_path.exists():
        print(f"错误：文件不存在：{file_path}", file=sys.stderr)
        return 1
    if file_path.suffix.lower() != ".pdf":
        print("错误：render-pages 仅支持 PDF 文件", file=sys.stderr)
        return 1

    try:
        page_count = get_pdf_page_count(file_path)
        page_indices = resolve_page_indices(args.pages, page_count, default_all=False)
        out_dir = args.out_dir.resolve() if args.out_dir else file_path.parent / f"{file_path.stem}_pages"
        rendered = render_pdf_pages(file_path, page_indices, out_dir, args.dpi)
    except Exception as exc:
        print(f"错误：页面截图失败：{exc}", file=sys.stderr)
        return 1

    manifest = {
        "version": __version__,
        "source_file": str(file_path),
        "page_count": page_count,
        "dpi": args.dpi,
        "pages": rendered,
    }
    manifest_path = out_dir / "render_pages.json"
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"输出目录：{out_dir}")
    print(f"页面数：{len(rendered)}")
    for item in rendered:
        print(f"输出：{item['file']}")
    print(f"清单：{manifest_path}")
    return 0


# ---------------------------------------------------------------------------
# 子命令：ocr（真实 OCR，可选依赖）
# ---------------------------------------------------------------------------

def cmd_ocr(args: argparse.Namespace) -> int:
    file_path = args.file.resolve()
    if not file_path.exists():
        print(f"错误：文件不存在：{file_path}", file=sys.stderr)
        return 1
    if file_path.suffix.lower() != ".pdf":
        print("错误：ocr 当前仅支持 PDF 文件", file=sys.stderr)
        return 1

    try:
        page_count = get_pdf_page_count(file_path)
        page_indices = resolve_page_indices(args.pages, page_count, default_all=False)
        pages = ocr_pdf_pages(file_path, page_indices, args.lang, args.dpi)
    except Exception as exc:
        print(f"错误：OCR 失败：{exc}", file=sys.stderr)
        print("依赖提示：pip install pymupdf pytesseract pillow；并安装系统 Tesseract OCR。", file=sys.stderr)
        return 1

    if args.output:
        out_path = Path(args.output).resolve()
    else:
        out_path = file_path.parent / f"{file_path.stem}_ocr.{args.format}"
    written = write_ocr_output(out_path, file_path, pages, args.format, args.lang, args.dpi)
    print(f"OCR 页数：{len(pages)}")
    print(f"输出：{written}")
    return 0


# ---------------------------------------------------------------------------
# 子命令：auto（智能流水线）
# ---------------------------------------------------------------------------

def cmd_auto(args: argparse.Namespace) -> int:
    file_path = args.file.resolve()
    if not file_path.exists():
        print(f"错误：文件不存在：{file_path}", file=sys.stderr)
        return 1

    fmt = detect_format(file_path)
    if fmt is None:
        supported = ", ".join(sorted(EXTENSION_FORMAT_MAP.keys()))
        print(f"错误：不支持的文件格式：{file_path.suffix}。支持：{supported}", file=sys.stderr)
        return 1

    out_dir = args.out_dir.resolve() if args.out_dir else file_path.parent / f"{file_path.stem}_auto_output"
    out_dir.mkdir(parents=True, exist_ok=True)
    max_pages = None if args.max_pages == 0 else args.max_pages

    try:
        metadata = file_metadata(file_path)
        text, result = extract_best_text(file_path, args.parser, args.start_page, max_pages)
        text, result, ocr_fallback = maybe_ocr_fallback(
            file_path,
            text,
            result,
            args.auto_ocr,
            args.min_quality,
            args.ocr_pages,
            args.ocr_lang,
            args.ocr_dpi,
        )
        classification = classify_document(text, metadata)
        gate = quality_gate_status(result, args.min_quality, args.fail_on_bad)
    except Exception as exc:
        print(f"错误：auto 流水线失败：{exc}", file=sys.stderr)
        return 1

    best_md = out_dir / "best.md"
    best_txt = out_dir / "best.txt"
    best_json = out_dir / "best.json"
    best_md.write_text(text + "\n", encoding="utf-8")
    best_txt.write_text(text + "\n", encoding="utf-8")
    best_json.write_text(
        json.dumps(
            {
                "version": __version__,
                "source_file": str(file_path),
                "parser": result.parser,
                "quality": asdict(result),
                "text": text,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )

    metadata_path = out_dir / "metadata.json"
    metadata_path.write_text(json.dumps(metadata, ensure_ascii=False, indent=2), encoding="utf-8")

    layout_path: Path | None = None
    page_map_path: Path | None = None
    if fmt == "pdf" and args.layout:
        try:
            layout = pdf_layout_json(file_path, args.start_page, max_pages)
            layout_path = write_json_file(out_dir / "layout.json", layout)
            page_map_path = write_json_file(out_dir / "page_map.json", layout_to_page_map(layout))
        except Exception as exc:
            write_json_file(out_dir / "layout_error.json", {"error": repr(exc)})

    chunks: list[dict[str, object]]
    if fmt == "pdf" and args.chunk_by == "page" and result.parser in PAGE_CHUNK_PARSERS:
        page_extractor = get_pdf_page_extractor(result.parser)
        chunks = split_pages_into_chunks(page_extractor(file_path, args.start_page, max_pages) if page_extractor else [])
    else:
        chunks = split_text_into_chunks(text, args.chunk_size, args.overlap)
    chunk_meta = {
        "version": __version__,
        "source_file": str(file_path),
        "parser": result.parser,
        "format": fmt,
        "chunk_by": args.chunk_by if fmt == "pdf" else "char",
    }
    chunk_paths = write_chunks(chunks, out_dir / "chunks", "all", chunk_meta)

    fields_path: str | None = None
    fields_payload: dict[str, object] | None = None
    if args.profile == "invoice" or (args.profile == "auto" and classification["profile"] == "invoice"):
        fields_payload = extract_invoice_fields_from_text(text, file_path, result.parser, result, strict=True)
        fields_path_obj = write_invoice_fields(fields_payload, out_dir / "fields.json", "json")
        write_invoice_fields(fields_payload, out_dir / "fields.md", "md")
        fields_path = str(fields_path_obj)

    decision = {
        "version": __version__,
        "source_file": str(file_path),
        "format": fmt,
        "parser": result.parser,
        "quality_score": result.quality_score,
        "quality_label": result.quality_label,
        "classification": classification,
        "quality_gate": gate,
        "needs_ocr": result.quality_label in QUALITY_BAD_LABELS or result.quality_score < args.min_quality,
        "ocr_fallback": ocr_fallback,
        "outputs": {
            "best_md": str(best_md),
            "best_txt": str(best_txt),
            "best_json": str(best_json),
            "metadata": str(metadata_path),
            "layout": str(layout_path) if layout_path else None,
            "page_map": str(page_map_path) if page_map_path else None,
            "chunks": [str(path) for path in chunk_paths],
            "fields": fields_path,
        },
    }
    decision_path = out_dir / "auto_report.json"
    decision_path.write_text(json.dumps(decision, ensure_ascii=False, indent=2), encoding="utf-8")

    report_lines = [
        "# Auto Parse Report",
        "",
        f"- Source: `{file_path}`",
        f"- Format: `{fmt}`",
        f"- Parser: `{result.parser}`",
        f"- Quality: `{result.quality_score}` (`{result.quality_label}`)",
        f"- Classification: `{classification['profile']}` (`{classification['confidence']}`)",
        f"- Needs OCR: `{decision['needs_ocr']}`",
        f"- OCR fallback: `{ocr_fallback}`",
        f"- Output directory: `{out_dir}`",
        "",
        "## Outputs",
        "",
        f"- Best Markdown: `{best_md}`",
        f"- Best TXT: `{best_txt}`",
        f"- Best JSON: `{best_json}`",
        f"- Metadata: `{metadata_path}`",
        f"- Auto JSON: `{decision_path}`",
    ]
    if fields_path:
        report_lines.append(f"- Fields: `{fields_path}`")
    if layout_path:
        report_lines.append(f"- Layout: `{layout_path}`")
    report_path = out_dir / "auto_report.md"
    report_path.write_text("\n".join(report_lines), encoding="utf-8")

    print(f"自动解析完成：{out_dir}")
    print(f"推荐解析器：{result.parser}，质量：{result.quality_score} {result.quality_label}")
    print(f"报告：{report_path}")
    if gate["enabled"] and not gate["passed"]:
        print(f"质量门禁未通过：{gate['reason']}", file=sys.stderr)
        return 2
    return 0


# ---------------------------------------------------------------------------
# 子命令：extract-fields（结构化字段抽取）
# ---------------------------------------------------------------------------

def cmd_extract_fields(args: argparse.Namespace) -> int:
    file_path = args.file.resolve()
    if not file_path.exists():
        print(f"错误：文件不存在：{file_path}", file=sys.stderr)
        return 1
    if args.profile != "invoice":
        print(f"错误：当前仅支持 profile=invoice，收到：{args.profile}", file=sys.stderr)
        return 1

    max_pages = None if args.max_pages == 0 else args.max_pages
    try:
        fields = extract_invoice_fields_from_file(file_path, args.parser, args.start_page, max_pages)
    except Exception as exc:
        print(f"错误：字段抽取失败：{exc}", file=sys.stderr)
        return 1

    if args.output:
        out_path = Path(args.output).resolve()
    else:
        out_path = file_path.parent / f"{file_path.stem}_fields.{args.format}"
    written = write_invoice_fields(fields, out_path, args.format)
    print(f"字段输出：{written}")
    validation = fields.get("validation")
    if isinstance(validation, dict):
        print(f"校验状态：{validation.get('status')}")
    return 0


# ---------------------------------------------------------------------------
# 子命令：export-xlsx（结构化批量导出）
# ---------------------------------------------------------------------------

def cmd_export_xlsx(args: argparse.Namespace) -> int:
    input_path = Path(args.path).resolve()
    if args.profile != "invoice":
        print(f"错误：当前仅支持 profile=invoice，收到：{args.profile}", file=sys.stderr)
        return 1

    extensions = parse_extensions(args.ext)
    max_pages = None if args.max_pages == 0 else args.max_pages
    try:
        files = collect_input_files(input_path, extensions, args.recursive)
    except Exception as exc:
        print(f"错误：收集文件失败：{exc}", file=sys.stderr)
        return 1

    if not files:
        print(f"错误：未找到输入文件，扩展名：{', '.join(sorted(extensions))}", file=sys.stderr)
        return 1

    rows: list[dict[str, object]] = []
    failures: list[dict[str, object]] = []
    for file_path in files:
        print(f"抽取: {file_path.name}...", end=" ", flush=True)
        try:
            fields = extract_invoice_fields_from_file(file_path, args.parser, args.start_page, max_pages)
            rows.append(fields)
            validation = fields.get("validation") if isinstance(fields.get("validation"), dict) else {}
            print(validation.get("status", "ok"))
        except Exception as exc:
            failures.append({"source_file": str(file_path), "error": repr(exc)})
            print(f"FAIL {exc}")

    if args.output:
        out_path = Path(args.output).resolve()
    else:
        if input_path.is_dir():
            out_path = input_path / "invoice_summary.xlsx"
        else:
            out_path = input_path.parent / f"{input_path.stem}_summary.xlsx"

    try:
        written = write_invoice_xlsx(rows, out_path)
    except Exception as exc:
        print(f"错误：XLSX 导出失败：{exc}", file=sys.stderr)
        return 1

    manifest_path = written.with_suffix(".json")
    manifest_path.write_text(
        json.dumps(
            {
                "version": __version__,
                "profile": args.profile,
                "input": str(input_path),
                "file_count": len(files),
                "success_count": len(rows),
                "failure_count": len(failures),
                "failures": failures,
                "output": str(written),
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    print(f"导出：{written}")
    print(f"清单：{manifest_path}")
    return 0 if not failures else 1


# ---------------------------------------------------------------------------
# 子命令：layout-json（版面结构输出）
# ---------------------------------------------------------------------------

def cmd_layout_json(args: argparse.Namespace) -> int:
    file_path = args.file.resolve()
    if not file_path.exists():
        print(f"错误：文件不存在：{file_path}", file=sys.stderr)
        return 1
    if file_path.suffix.lower() != ".pdf":
        print("错误：layout-json 当前仅支持 PDF", file=sys.stderr)
        return 1
    max_pages = None if args.max_pages == 0 else args.max_pages
    try:
        layout = pdf_layout_json(file_path, args.start_page, max_pages)
    except Exception as exc:
        print(f"错误：layout-json 失败：{exc}", file=sys.stderr)
        return 1
    out_path = Path(args.output).resolve() if args.output else file_path.parent / f"{file_path.stem}_layout.json"
    write_json_file(out_path, layout)
    print(f"输出：{out_path}")
    return 0


# ---------------------------------------------------------------------------
# 子命令：verify-fields（字段校验）
# ---------------------------------------------------------------------------

def cmd_verify_fields(args: argparse.Namespace) -> int:
    path = Path(args.path).resolve()
    if not path.exists():
        print(f"错误：文件不存在：{path}", file=sys.stderr)
        return 1
    if args.profile != "invoice":
        print(f"错误：当前仅支持 profile=invoice，收到：{args.profile}", file=sys.stderr)
        return 1

    try:
        if path.suffix.lower() == ".json":
            fields = json.loads(path.read_text(encoding="utf-8"))
            fields["validation"] = validate_invoice_fields(fields, strict=args.strict)
        else:
            max_pages = None if args.max_pages == 0 else args.max_pages
            fields = extract_invoice_fields_from_file(path, args.parser, args.start_page, max_pages, strict=args.strict)
    except Exception as exc:
        print(f"错误：字段校验失败：{exc}", file=sys.stderr)
        return 1

    validation = fields.get("validation") if isinstance(fields.get("validation"), dict) else {}
    out_path = Path(args.output).resolve() if args.output else path.parent / f"{path.stem}_verify.json"
    write_json_file(out_path, fields)
    print(f"校验状态：{validation.get('status')}")
    print(f"输出：{out_path}")
    return 0 if validation.get("status") == "ok" else 2


# ---------------------------------------------------------------------------
# 子命令：classify（自动分类和策略）
# ---------------------------------------------------------------------------

def cmd_classify(args: argparse.Namespace) -> int:
    file_path = args.file.resolve()
    if not file_path.exists():
        print(f"错误：文件不存在：{file_path}", file=sys.stderr)
        return 1
    max_pages = None if args.max_pages == 0 else args.max_pages
    try:
        metadata = file_metadata(file_path)
        text, result = extract_best_text(file_path, args.parser, args.start_page, max_pages)
        classification = classify_document(text, metadata)
        payload = {
            "version": __version__,
            "source_file": str(file_path),
            "parser": result.parser,
            "quality": asdict(result),
            "classification": classification,
        }
    except Exception as exc:
        print(f"错误：分类失败：{exc}", file=sys.stderr)
        return 1

    if args.output:
        out_path = Path(args.output).resolve()
        write_json_file(out_path, payload)
        print(f"输出：{out_path}")
    else:
        print(json.dumps(payload, ensure_ascii=False, indent=2))
    return 0


# ---------------------------------------------------------------------------
# 子命令：knowledge-pack（可追溯 RAG 包）
# ---------------------------------------------------------------------------

def cmd_knowledge_pack(args: argparse.Namespace) -> int:
    file_path = args.file.resolve()
    if not file_path.exists():
        print(f"错误：文件不存在：{file_path}", file=sys.stderr)
        return 1
    max_pages = None if args.max_pages == 0 else args.max_pages
    out_dir = args.out_dir.resolve() if args.out_dir else file_path.parent / f"{file_path.stem}_knowledge_pack"
    try:
        manifest = write_knowledge_pack(
            file_path,
            out_dir,
            args.parser,
            args.start_page,
            max_pages,
            args.chunk_by,
            args.chunk_size,
            args.overlap,
            args.min_quality,
        )
    except Exception as exc:
        print(f"错误：knowledge-pack 失败：{exc}", file=sys.stderr)
        return 1
    print(f"知识包：{out_dir}")
    print(f"manifest：{manifest['outputs']['manifest']}")
    return 0


# ---------------------------------------------------------------------------
# 子命令：batch-knowledge（批量知识库生产）
# ---------------------------------------------------------------------------

def cmd_batch_knowledge(args: argparse.Namespace) -> int:
    input_dir = Path(args.dir).resolve()
    extensions = parse_extensions(args.ext)
    max_pages = None if args.max_pages == 0 else args.max_pages
    try:
        files = collect_input_files(input_dir, extensions, args.recursive)
    except Exception as exc:
        print(f"错误：收集文件失败：{exc}", file=sys.stderr)
        return 1
    if not files:
        print(f"错误：没有找到 {', '.join(sorted(extensions))} 文件", file=sys.stderr)
        return 1

    out_dir = args.out_dir.resolve() if args.out_dir else input_dir / "knowledge_packs"
    out_dir.mkdir(parents=True, exist_ok=True)
    manifests: list[dict[str, object]] = []
    failures: list[dict[str, object]] = []
    for file_path in files:
        print(f"知识包: {file_path.name}...", end=" ", flush=True)
        try:
            pack_dir = out_dir / file_path.stem
            manifest = write_knowledge_pack(
                file_path,
                pack_dir,
                args.parser,
                args.start_page,
                max_pages,
                args.chunk_by,
                args.chunk_size,
                args.overlap,
                args.min_quality,
            )
            manifests.append(manifest)
            print("ok")
        except Exception as exc:
            failures.append({"source_file": str(file_path), "error": repr(exc)})
            print(f"FAIL {exc}")

    index = {
        "version": __version__,
        "type": "batch-knowledge",
        "input_dir": str(input_dir),
        "file_count": len(files),
        "success_count": len(manifests),
        "failure_count": len(failures),
        "manifests": [manifest.get("outputs", {}).get("manifest") for manifest in manifests],
        "failures": failures,
    }
    index_path = write_json_file(out_dir / "index.json", index)
    print(f"索引：{index_path}")
    return 0 if not failures else 1


# ---------------------------------------------------------------------------
# 子命令：qa（本地抽取式问答）
# ---------------------------------------------------------------------------

def cmd_qa(args: argparse.Namespace) -> int:
    source = args.source.resolve()
    if not source.exists():
        print(f"错误：输入不存在：{source}", file=sys.stderr)
        return 1

    max_pages = None if args.max_pages == 0 else args.max_pages
    if args.output:
        out_path = Path(args.output).resolve()
    elif args.format == "all":
        out_path = source.parent / f"{source.stem}_qa"
    else:
        out_path = source.parent / f"{source.stem}_qa.{args.format}"

    try:
        payload = build_qa_payload(
            source,
            args.question,
            args.parser,
            args.start_page,
            max_pages,
            args.chunk_by,
            args.chunk_size,
            args.overlap,
            args.top_k,
            args.answer_sentences,
        )
        written = write_qa_output(payload, out_path, args.format)
    except Exception as exc:
        print(f"错误：qa 失败：{exc}", file=sys.stderr)
        return 1

    answer = payload.get("answer") if isinstance(payload.get("answer"), dict) else {}
    print(f"状态：{answer.get('status')}")
    print(f"输出：{', '.join(str(path) for path in written)}")
    return 0 if answer.get("status") == "found" else 2


# ---------------------------------------------------------------------------
# 子命令：diff-docs（文档版本差异对比）
# ---------------------------------------------------------------------------

def cmd_diff_docs(args: argparse.Namespace) -> int:
    left_path = args.left.resolve()
    right_path = args.right.resolve()
    if not left_path.exists():
        print(f"错误：左侧文件不存在：{left_path}", file=sys.stderr)
        return 1
    if not right_path.exists():
        print(f"错误：右侧文件不存在：{right_path}", file=sys.stderr)
        return 1

    max_pages = None if args.max_pages == 0 else args.max_pages
    if args.output:
        out_path = Path(args.output).resolve()
    elif args.format == "all":
        out_path = left_path.parent / f"{left_path.stem}_vs_{right_path.stem}_diff"
    else:
        out_path = left_path.parent / f"{left_path.stem}_vs_{right_path.stem}_diff.{args.format}"

    try:
        payload = build_diff_payload(
            left_path,
            right_path,
            args.parser,
            args.start_page,
            max_pages,
            args.profile,
            args.max_diff_lines,
        )
        written = write_diff_output(payload, out_path, args.format)
    except Exception as exc:
        print(f"错误：diff-docs 失败：{exc}", file=sys.stderr)
        return 1

    line_diff = payload.get("line_diff") if isinstance(payload.get("line_diff"), dict) else {}
    similarity = payload.get("similarity") if isinstance(payload.get("similarity"), dict) else {}
    print(f"相似度：{similarity.get('full_text_ratio')}")
    print(f"变更块：{line_diff.get('changed_block_count')}")
    print(f"输出：{', '.join(str(path) for path in written)}")
    return 0


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def build_parser() -> argparse.ArgumentParser:
    all_parsers = list_all_parsers()

    parser = argparse.ArgumentParser(
        prog="parse_pdf_compare",
        description="多格式文档解析工具：对比、转换、批量处理、表格提取。",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("--version", action="version", version=f"%(prog)s {__version__}")

    subparsers = parser.add_subparsers(dest="command", help="子命令")

    # --- compare ---
    p_compare = subparsers.add_parser("compare", help="多解析器对比（默认子命令）",
                                      formatter_class=argparse.RawDescriptionHelpFormatter)
    p_compare.add_argument("file", type=Path, help="文件路径")
    p_compare.add_argument("--out-dir", type=Path, default=None, help="输出目录")
    p_compare.add_argument("--max-pages", type=int, default=30, help="（仅 PDF）前 N 页，0=全量。默认 30")
    p_compare.add_argument("--start-page", type=int, default=1, help="（仅 PDF）起始页。默认 1")
    p_compare.add_argument("--parsers", default=None,
                           help=f"逗号分隔解析器。默认全部：{','.join(all_parsers)}")
    p_compare.add_argument("--output-format", choices=sorted(TEXT_OUTPUT_FORMATS), default="md",
                           help="解析器输出格式：md/txt/json/all。默认 md")
    p_compare.add_argument("--similarity-chars", type=int, default=50000, help="相似度比较字符数")
    p_compare.add_argument("--min-quality", type=float, default=0.0,
                           help="最低质量评分门槛，0-1。低于该值返回退出码 2")
    p_compare.add_argument("--fail-on-bad", action="store_true",
                           help="当最佳结果标签为 empty/bad 时返回退出码 2")
    p_compare.add_argument("--ocr-fallback", action="store_true",
                           help="标记需要 OCR 回退；真实 OCR 请使用 ocr 子命令")
    p_compare.add_argument("--parallel", action="store_true", help="并行执行")

    # --- vote ---
    p_vote = subparsers.add_parser("vote", help="PDF 多解析器投票，输出客户可交付最佳文本")
    p_vote.add_argument("file", type=Path, help="PDF 文件路径")
    p_vote.add_argument("--out-dir", type=Path, default=None, help="输出目录（默认 <文件名>_vote_output）")
    p_vote.add_argument("--max-pages", type=int, default=30, help="前 N 页，0=全量。默认 30")
    p_vote.add_argument("--start-page", type=int, default=1, help="起始页。默认 1")
    p_vote.add_argument("--parsers", default=None,
                        help=f"逗号分隔解析器。默认全部 PDF 解析器：{','.join(get_parsers_for_format('pdf'))}")
    p_vote.add_argument("--format", choices=sorted(TEXT_OUTPUT_FORMATS), default="md",
                        help="客户交付 best 输出格式：md/txt/json/all。默认 md")
    profile_choices = ["auto", "invoice", "contract", "bank_statement", "quotation", "purchase_order", "report", "annual_report", "none"]
    p_vote.add_argument("--profile", choices=profile_choices, default="auto",
                        help="投票领域加权 profile。invoice 会启用字段校验加权；其它业务 profile 会做轻量结构化抽取。默认 auto")
    p_vote.add_argument("--customer", action="store_true",
                        help="额外输出 customer_best.md/txt/json；发票 profile 下输出结构化客户交付稿")
    p_vote.add_argument("--no-field-fusion", action="store_true",
                        help="关闭发票字段级投票融合，customer_best 只使用整份投票胜出解析器字段")
    p_vote.add_argument("--field-layout", action="store_true",
                        help="写 customer_best 时提取 PDF layout，并尽量给发票字段填充 page/bbox/location")
    p_vote.add_argument("--field-layout-max-pages", type=int, default=30,
                        help="--field-layout 的 layout 页数，0=全量。默认 30")
    p_vote.add_argument("--similarity-chars", type=int, default=50000, help="投票共识相似度比较字符数")
    p_vote.add_argument("--min-quality", type=float, default=0.5, help="最低质量评分门槛，默认 0.5")
    p_vote.add_argument("--fail-on-bad", action="store_true", help="胜出结果标签为 empty/bad 时返回退出码 2")
    p_vote.add_argument("--probe-before-vote", action="store_true",
                        help="正式投票前先用真实文件小样本体检解析器，只让 ready 解析器进入投票")
    p_vote.add_argument("--probe-max-pages", type=int, default=1,
                        help="--probe-before-vote 的体检页数，0=全量。默认 1")
    p_vote.add_argument("--probe-min-quality", type=float, default=0.5,
                        help="--probe-before-vote 的体检质量门槛。默认 0.5")
    p_vote.add_argument("--timeout", type=float, default=None,
                        help="单个解析器最大运行秒数；超时会杀掉子进程并标记 timeout。默认不限制")
    p_vote.add_argument("--parser-health-cache", action="store_true",
                        help="启用解析器健康缓存，跳过 24 小时内同文件同页段近期失败或超时的解析器")
    p_vote.add_argument("--health-cache", default=None,
                        help="解析器健康缓存 JSON 路径；默认输出目录下 .parser_health_cache.json")
    p_vote.add_argument("--result-cache", action="store_true",
                        help="启用解析结果缓存，同一 PDF 同一 parser/page 范围不重复重解析")
    p_vote.add_argument("--result-cache-dir", default=None,
                        help="解析结果缓存目录；默认输出目录下 .parser_result_cache")
    p_vote.add_argument("--field-evidence", action="store_true",
                        help="配合 --customer 生成字段截图证据包 field_evidence/")
    p_vote.add_argument("--review-html", action="store_true",
                        help="配合 --customer 生成本地可视化审阅页 review.html")
    p_vote.add_argument("--parallel", action="store_true", help="并行执行")

    # --- customer-pack ---
    p_pack_customer = subparsers.add_parser("customer-pack", help="复杂 PDF 客户交付包：最佳文本、表格、layout、metadata、manifest")
    p_pack_customer.add_argument("file", type=Path, help="PDF 文件路径")
    p_pack_customer.add_argument("--out-dir", type=Path, default=None, help="输出目录（默认 <文件名>_customer_pack）")
    p_pack_customer.add_argument("--max-pages", type=int, default=0, help="正式文本投票页数，0=全量。默认 0")
    p_pack_customer.add_argument("--start-page", type=int, default=1, help="起始页。默认 1")
    p_pack_customer.add_argument("--probe-max-pages", type=int, default=1, help="预检体检页数，0=全量。默认 1")
    p_pack_customer.add_argument("--table-pages", default="all", help="表格提取页码范围，如 1-5、1,3,5、all。默认 all")
    p_pack_customer.add_argument("--layout-max-pages", type=int, default=30, help="layout 输出页数，0=全量。默认 30")
    p_pack_customer.add_argument("--parsers", default=None,
                                 help=f"逗号分隔解析器。默认全部 PDF 解析器：{','.join(get_parsers_for_format('pdf'))}")
    p_pack_customer.add_argument("--profile", choices=profile_choices, default="auto",
                                 help="投票领域加权 profile。invoice 会启用字段校验加权；其它业务 profile 会做轻量结构化抽取。默认 auto")
    p_pack_customer.add_argument("--min-quality", type=float, default=0.5, help="最低质量评分门槛，默认 0.5")
    p_pack_customer.add_argument("--fail-on-bad", action="store_true", help="胜出结果标签为 empty/bad 时返回退出码 2")
    p_pack_customer.add_argument("--similarity-chars", type=int, default=50000, help="投票共识相似度比较字符数")
    p_pack_customer.add_argument("--timeout", type=float, default=None,
                                 help="单个解析器最大运行秒数；超时会杀掉子进程并标记 timeout。默认不限制")
    p_pack_customer.add_argument("--parser-health-cache", action="store_true",
                                 help="启用解析器健康缓存，跳过 24 小时内同文件同页段近期失败或超时的解析器")
    p_pack_customer.add_argument("--health-cache", default=None,
                                 help="解析器健康缓存 JSON 路径；默认输出目录下 .parser_health_cache.json")
    p_pack_customer.add_argument("--no-tables", action="store_true", help="不提取 PDF 表格")
    p_pack_customer.add_argument("--no-layout", action="store_true", help="不输出 layout.json/page_map.json")
    p_pack_customer.add_argument("--no-field-fusion", action="store_true",
                                 help="关闭发票字段级投票融合，customer_best 只使用整份投票胜出解析器字段")
    p_pack_customer.add_argument("--no-field-evidence", action="store_true", help="不生成字段截图证据包")
    p_pack_customer.add_argument("--no-review-html", action="store_true", help="不生成可视化 HTML 审阅页")
    p_pack_customer.add_argument("--result-cache", action="store_true",
                                 help="启用解析结果缓存，同一 PDF 同一 parser/page 范围不重复重解析")
    p_pack_customer.add_argument("--result-cache-dir", default=None,
                                 help="解析结果缓存目录；默认输出目录下 .parser_result_cache")
    p_pack_customer.add_argument("--parallel", action="store_true", help="并行执行解析器")

    # --- probe ---
    p_probe = subparsers.add_parser("probe", help="用真实文件逐个测试解析器运行状态")
    p_probe.add_argument("file", type=Path, help="文件路径")
    p_probe.add_argument("--out-dir", type=Path, default=None, help="输出目录（默认 <文件名>_probe_output）")
    p_probe.add_argument("--max-pages", type=int, default=1, help="（仅 PDF）体检前 N 页，0=全量。默认 1")
    p_probe.add_argument("--start-page", type=int, default=1, help="（仅 PDF）起始页。默认 1")
    p_probe.add_argument("--parsers", default=None,
                         help=f"逗号分隔解析器。默认当前格式全部注册解析器：{','.join(all_parsers)}")
    p_probe.add_argument("--min-quality", type=float, default=0.5, help="最低质量评分门槛，默认 0.5")
    p_probe.add_argument("--fail-on-bad", action="store_true",
                         help="没有任何解析器通过运行时体检时返回退出码 2")
    p_probe.add_argument("--keep-outputs", action="store_true",
                         help="保留每个成功解析器的 Markdown 样本文本；默认写 txt 样本")
    p_probe.add_argument("--timeout", type=float, default=None,
                         help="单个解析器最大运行秒数；超时会杀掉子进程并标记 timeout。默认不限制")
    p_probe.add_argument("--parser-health-cache", action="store_true",
                         help="启用解析器健康缓存，跳过 24 小时内同文件同页段近期失败或超时的解析器")
    p_probe.add_argument("--health-cache", default=None,
                         help="解析器健康缓存 JSON 路径；默认输出目录下 .parser_health_cache.json")
    p_probe.add_argument("--result-cache", action="store_true",
                         help="启用解析结果缓存，同一文件同一 parser/page 范围不重复重解析")
    p_probe.add_argument("--result-cache-dir", default=None,
                         help="解析结果缓存目录；默认输出目录下 .parser_result_cache")
    p_probe.add_argument("--parallel", action="store_true", help="并行执行")

    # --- convert ---
    p_convert = subparsers.add_parser("convert", help="文档转指定格式")
    p_convert.add_argument("file", type=Path, help="文件路径")
    p_convert.add_argument("--parser", default="markitdown", help="解析器（默认 markitdown）")
    p_convert.add_argument("--format", choices=sorted(TEXT_OUTPUT_FORMATS), default="md",
                           help="输出格式：md/txt/json/all。默认 md")
    p_convert.add_argument("-o", "--output", default=None,
                           help="输出文件路径；--format all 时为输出目录。默认按格式生成")
    p_convert.add_argument("--max-pages", type=int, default=30, help="（仅 PDF）前 N 页，0=全量")
    p_convert.add_argument("--start-page", type=int, default=1, help="（仅 PDF）起始页")
    p_convert.add_argument("--ocr-fallback", action="store_true",
                           help="标记需要 OCR 回退；真实 OCR 请使用 ocr 子命令")

    # --- batch ---
    p_batch = subparsers.add_parser("batch", help="批量转换目录下所有文档")
    p_batch.add_argument("dir", type=str, help="输入目录")
    p_batch.add_argument("--parser", default="markitdown", help="解析器（默认 markitdown）")
    p_batch.add_argument("--output-dir", default=None, help="输出目录（默认 <输入目录>/batch_output）")
    p_batch.add_argument("--ext", default=".pdf", help="文件扩展名，逗号分隔（默认 .pdf）")
    p_batch.add_argument("--format", choices=sorted(TEXT_OUTPUT_FORMATS), default="md",
                         help="输出格式：md/txt/json/all。默认 md")
    p_batch.add_argument("--parallel", action="store_true", help="并行处理")

    # --- scan-dir ---
    p_scan = subparsers.add_parser("scan-dir", help="批量扫描目录内文档解析质量")
    p_scan.add_argument("dir", type=str, help="输入目录")
    p_scan.add_argument("--parser", default="auto", help="解析器，默认 auto")
    p_scan.add_argument("--out-dir", type=Path, default=None, help="输出目录（默认 <输入目录>/scan_output）")
    p_scan.add_argument("--ext", default=".pdf", help="文件扩展名，逗号分隔（默认 .pdf）")
    p_scan.add_argument("--recursive", action="store_true", help="递归扫描子目录")
    p_scan.add_argument("--max-pages", type=int, default=3, help="（仅 PDF）抽样前 N 页，0=全量。默认 3")
    p_scan.add_argument("--start-page", type=int, default=1, help="（仅 PDF）起始页。默认 1")
    p_scan.add_argument("--min-quality", type=float, default=0.5, help="最低质量评分门槛，默认 0.5")
    p_scan.add_argument("--fail-on-bad", action="store_true", help="存在门禁失败文件时返回退出码 2")

    # --- tables ---
    p_tables = subparsers.add_parser("tables", help="从 PDF 提取表格")
    p_tables.add_argument("file", type=Path, help="PDF 文件路径")
    p_tables.add_argument("--pages", default=None, help="页码范围，如 1-5 或 1,3,5（默认全部）")
    p_tables.add_argument("--format", choices=["md", "csv", "json", "all"], default="md",
                          help="输出格式：md/csv/json/all。默认 md")
    p_tables.add_argument("-o", "--output", default=None,
                          help="输出文件路径；--format all 时为输出目录")

    # --- table-vote ---
    p_table_vote = subparsers.add_parser("table-vote", help="PDF 多表格解析器投票，输出最佳表格")
    p_table_vote.add_argument("file", type=Path, help="PDF 文件路径")
    p_table_vote.add_argument("--pages", default="all", help="页码范围，如 1-5、1,3,5、all。默认 all")
    p_table_vote.add_argument("--methods", default="pdfplumber,pdfplumber-text,pymupdf-text",
                              help="逗号分隔表格解析器：pdfplumber,pdfplumber-text,pymupdf-text")
    p_table_vote.add_argument("--format", choices=["md", "csv", "json", "all"], default="all",
                              help="最佳表格输出格式：md/csv/json/all。默认 all")
    p_table_vote.add_argument("--out-dir", type=Path, default=None,
                              help="输出目录（默认 <文件名>_table_vote）")

    # --- init-golden ---
    p_init_golden = subparsers.add_parser("init-golden", help="从 PDF 文件或目录初始化 Golden case 样本库")
    p_init_golden.add_argument("paths", nargs="+", type=Path, help="PDF 文件或目录，可传多个")
    p_init_golden.add_argument("--out-dir", type=Path, required=True, help="Golden 样本库输出目录")
    p_init_golden.add_argument("--recursive", action="store_true", help="递归扫描目录中的 PDF")
    p_init_golden.add_argument("--parsers", default=GOLDEN_DEFAULT_PARSERS,
                               help=f"用例默认解析器，逗号分隔。默认 {GOLDEN_DEFAULT_PARSERS}")
    p_init_golden.add_argument("--include-ocr", action="store_true",
                               help="把 ocr-tesseract 加入每个 case 的 parser 列表，适合混有扫描件/坏文本层 PDF 的样本库")
    p_init_golden.add_argument("--max-pages", type=int, default=3,
                               help="每个 case 默认评测页数；0=全量。默认 3")
    p_init_golden.add_argument("--min-quality", type=float, default=0.5,
                               help="每个 case 默认最低投票分。默认 0.5")
    p_init_golden.add_argument("--min-non-space-chars", type=int, default=120,
                               help="每个 case 默认 winner 最少非空白字符数。默认 120")
    p_init_golden.add_argument("--name-prefix", default=None, help="可选 case 名称前缀")
    p_init_golden.add_argument("--force", action="store_true", help="覆盖已存在 case JSON")

    # --- eval-golden ---
    p_eval = subparsers.add_parser("eval-golden", help="运行 golden PDF 用例，评测投票策略是否退步")
    p_eval.add_argument("path", type=Path, help="golden case JSON 文件或目录")
    p_eval.add_argument("--out-dir", type=Path, default=None, help="输出目录（默认 golden_eval_output）")
    p_eval.add_argument("--recursive", action="store_true", help="递归查找 case JSON")
    p_eval.add_argument("--max-pages", type=int, default=-1,
                        help="覆盖用例中的 max_pages；0=全量，-1=使用用例设置。默认 -1")
    p_eval.add_argument("--start-page", type=int, default=None, help="覆盖用例中的起始页。默认使用用例设置或 1")
    p_eval.add_argument("--parsers", default=None,
                        help="覆盖用例中的解析器，逗号分隔。默认使用用例设置或全部 PDF 解析器")
    p_eval.add_argument("--profile", choices=profile_choices, default=None,
                        help="覆盖用例中的 profile。默认使用用例设置")
    p_eval.add_argument("--min-quality", type=float, default=0.5, help="最低质量评分门槛。默认 0.5")
    p_eval.add_argument("--fail-on-bad", action="store_true", help="胜出结果标签为 empty/bad 时视为失败")
    p_eval.add_argument("--similarity-chars", type=int, default=50000, help="投票共识相似度比较字符数")
    p_eval.add_argument("--timeout", type=float, default=None,
                        help="单个解析器最大运行秒数；超时会杀掉子进程并标记 timeout。默认不限制")
    p_eval.add_argument("--parser-health-cache", action="store_true",
                        help="启用解析器健康缓存，跳过 24 小时内同文件同页段近期失败或超时的解析器")
    p_eval.add_argument("--health-cache", default=None,
                        help="解析器健康缓存 JSON 路径；默认输出目录下 .parser_health_cache.json")
    p_eval.add_argument("--result-cache", action="store_true",
                        help="启用解析结果缓存，同一 PDF 同一 parser/page 范围不重复重解析")
    p_eval.add_argument("--result-cache-dir", default=None,
                        help="解析结果缓存目录；默认输出目录下 .parser_result_cache")
    p_eval.add_argument("--parallel", action="store_true", help="并行执行解析器")

    # --- doctor ---
    p_doctor = subparsers.add_parser("doctor", help="检查解析器依赖是否可用")
    p_doctor.add_argument("--format", choices=sorted(FORMAT_PARSERS), default=None,
                          help="只检查某一类格式")
    p_doctor.add_argument("--json", action="store_true", help="输出 JSON")
    p_doctor.add_argument("--ocr", action="store_true", help="同时检查 OCR 可选依赖")
    p_doctor.add_argument("--opendataloader", action="store_true", help="同时检查 opendataloader-pdf 依赖（需要 Java）")
    p_doctor.add_argument("--strict", action="store_true",
                          help="存在缺失依赖时返回非零退出码")

    # --- metadata ---
    p_metadata = subparsers.add_parser("metadata", help="提取文件元数据")
    p_metadata.add_argument("file", type=Path, help="文件路径")
    p_metadata.add_argument("--format", choices=["md", "json"], default="json",
                            help="输出格式：md/json。默认 json")
    p_metadata.add_argument("-o", "--output", default=None, help="输出文件路径")

    # --- chunk ---
    p_chunk = subparsers.add_parser("chunk", help="分块输出文档文本")
    p_chunk.add_argument("file", type=Path, help="文件路径")
    p_chunk.add_argument("--parser", default="markitdown", help="解析器（默认 markitdown）")
    p_chunk.add_argument("--chunk-by", choices=["char", "page"], default="char",
                         help="分块方式：char/page。默认 char")
    p_chunk.add_argument("--format", choices=["jsonl", "json", "md", "txt", "all"], default="jsonl",
                         help="输出格式：jsonl/json/md/txt/all。默认 jsonl")
    p_chunk.add_argument("--chunk-size", type=int, default=2000, help="每块字符数。默认 2000")
    p_chunk.add_argument("--overlap", type=int, default=200, help="块间重叠字符数。默认 200")
    p_chunk.add_argument("--max-pages", type=int, default=30, help="（仅 PDF）前 N 页，0=全量")
    p_chunk.add_argument("--start-page", type=int, default=1, help="（仅 PDF）起始页")
    p_chunk.add_argument("--ocr-fallback", action="store_true",
                         help="标记需要 OCR 回退；真实 OCR 请使用 ocr 子命令")
    p_chunk.add_argument("-o", "--output", default=None,
                         help="输出文件路径；--format all 时为输出目录")

    # --- render-pages ---
    p_render = subparsers.add_parser("render-pages", help="把 PDF 页面渲染为 PNG 图片")
    p_render.add_argument("file", type=Path, help="PDF 文件路径")
    p_render.add_argument("--pages", default="1", help="页码范围，如 1-5、1,3,5、all。默认 1")
    p_render.add_argument("--dpi", type=int, default=200, help="渲染 DPI。默认 200")
    p_render.add_argument("--out-dir", type=Path, default=None, help="输出目录（默认 <文件名>_pages）")

    # --- ocr ---
    p_ocr = subparsers.add_parser("ocr", help="对 PDF 页面执行真实 OCR（Tesseract，可选依赖）")
    p_ocr.add_argument("file", type=Path, help="PDF 文件路径")
    p_ocr.add_argument("--pages", default="1", help="页码范围，如 1-5、1,3,5、all。默认 1")
    p_ocr.add_argument("--lang", default="chi_sim+eng", help="Tesseract 语言。默认 chi_sim+eng")
    p_ocr.add_argument("--dpi", type=int, default=200, help="渲染 DPI。默认 200")
    p_ocr.add_argument("--format", choices=["txt", "md", "json"], default="txt", help="输出格式。默认 txt")
    p_ocr.add_argument("-o", "--output", default=None, help="输出文件路径")

    # --- auto ---
    p_auto = subparsers.add_parser("auto", help="智能解析流水线：元数据、最佳文本、分块、字段和报告")
    p_auto.add_argument("file", type=Path, help="文件路径")
    p_auto.add_argument("--parser", default="auto", help="解析器，默认 auto")
    p_auto.add_argument("--profile", choices=["auto", "invoice"], default="auto", help="结构化 profile，默认 auto")
    p_auto.add_argument("--out-dir", type=Path, default=None, help="输出目录（默认 <文件名>_auto_output）")
    p_auto.add_argument("--max-pages", type=int, default=30, help="（仅 PDF）前 N 页，0=全量。默认 30")
    p_auto.add_argument("--start-page", type=int, default=1, help="（仅 PDF）起始页。默认 1")
    p_auto.add_argument("--min-quality", type=float, default=0.5, help="最低质量评分门槛，默认 0.5")
    p_auto.add_argument("--fail-on-bad", action="store_true", help="质量不达标时返回退出码 2")
    p_auto.add_argument("--chunk-by", choices=["char", "page"], default="char", help="分块方式，默认 char")
    p_auto.add_argument("--chunk-size", type=int, default=2000, help="字符分块大小。默认 2000")
    p_auto.add_argument("--overlap", type=int, default=200, help="字符分块重叠。默认 200")
    p_auto.add_argument("--auto-ocr", action="store_true", help="质量低于阈值时尝试 Tesseract OCR 回退")
    p_auto.add_argument("--ocr-pages", default="1", help="OCR 回退页码范围，如 1-3、all。默认 1")
    p_auto.add_argument("--ocr-lang", default="chi_sim+eng", help="OCR 语言。默认 chi_sim+eng")
    p_auto.add_argument("--ocr-dpi", type=int, default=200, help="OCR 渲染 DPI。默认 200")
    p_auto.add_argument("--layout", action="store_true", help="同时输出 layout.json 和 page_map.json")

    # --- extract-fields ---
    p_fields = subparsers.add_parser("extract-fields", help="从文档抽取结构化字段")
    p_fields.add_argument("file", type=Path, help="文件路径")
    p_fields.add_argument("--profile", choices=["invoice"], default="invoice", help="字段 profile。默认 invoice")
    p_fields.add_argument("--parser", default="auto", help="解析器，默认 auto")
    p_fields.add_argument("--format", choices=["json", "md"], default="json", help="输出格式。默认 json")
    p_fields.add_argument("--max-pages", type=int, default=3, help="（仅 PDF）前 N 页，0=全量。默认 3")
    p_fields.add_argument("--start-page", type=int, default=1, help="（仅 PDF）起始页。默认 1")
    p_fields.add_argument("-o", "--output", default=None, help="输出文件路径")

    # --- export-xlsx ---
    p_xlsx = subparsers.add_parser("export-xlsx", help="批量抽取结构化字段并导出 XLSX")
    p_xlsx.add_argument("path", type=str, help="文件或目录路径")
    p_xlsx.add_argument("--profile", choices=["invoice"], default="invoice", help="字段 profile。默认 invoice")
    p_xlsx.add_argument("--parser", default="auto", help="解析器，默认 auto")
    p_xlsx.add_argument("--ext", default=".pdf", help="目录输入时的扩展名，逗号分隔。默认 .pdf")
    p_xlsx.add_argument("--recursive", action="store_true", help="递归扫描目录")
    p_xlsx.add_argument("--max-pages", type=int, default=3, help="（仅 PDF）前 N 页，0=全量。默认 3")
    p_xlsx.add_argument("--start-page", type=int, default=1, help="（仅 PDF）起始页。默认 1")
    p_xlsx.add_argument("-o", "--output", default=None, help="输出 XLSX 路径")

    # --- layout-json ---
    p_layout = subparsers.add_parser("layout-json", help="输出 PDF 页面、块、行、span 和坐标 JSON")
    p_layout.add_argument("file", type=Path, help="PDF 文件路径")
    p_layout.add_argument("--max-pages", type=int, default=30, help="前 N 页，0=全量。默认 30")
    p_layout.add_argument("--start-page", type=int, default=1, help="起始页。默认 1")
    p_layout.add_argument("-o", "--output", default=None, help="输出 JSON 路径")

    # --- verify-fields ---
    p_verify = subparsers.add_parser("verify-fields", help="校验结构化字段")
    p_verify.add_argument("path", type=str, help="字段 JSON 或原始文档路径")
    p_verify.add_argument("--profile", choices=["invoice"], default="invoice", help="字段 profile。默认 invoice")
    p_verify.add_argument("--parser", default="auto", help="原始文档解析器，默认 auto")
    p_verify.add_argument("--strict", action="store_true", help="启用严格校验")
    p_verify.add_argument("--max-pages", type=int, default=3, help="原始 PDF 前 N 页，0=全量。默认 3")
    p_verify.add_argument("--start-page", type=int, default=1, help="原始 PDF 起始页。默认 1")
    p_verify.add_argument("-o", "--output", default=None, help="输出 JSON 路径")

    # --- classify ---
    p_classify = subparsers.add_parser("classify", help="自动分类文档并推荐处理策略")
    p_classify.add_argument("file", type=Path, help="文件路径")
    p_classify.add_argument("--parser", default="auto", help="解析器，默认 auto")
    p_classify.add_argument("--max-pages", type=int, default=3, help="前 N 页，0=全量。默认 3")
    p_classify.add_argument("--start-page", type=int, default=1, help="起始页。默认 1")
    p_classify.add_argument("-o", "--output", default=None, help="输出 JSON 路径；默认打印到 stdout")

    # --- knowledge-pack ---
    p_pack = subparsers.add_parser("knowledge-pack", help="生成可追溯 RAG/知识包")
    p_pack.add_argument("file", type=Path, help="文件路径")
    p_pack.add_argument("--parser", default="auto", help="解析器，默认 auto")
    p_pack.add_argument("--out-dir", type=Path, default=None, help="输出目录（默认 <文件名>_knowledge_pack）")
    p_pack.add_argument("--max-pages", type=int, default=30, help="前 N 页，0=全量。默认 30")
    p_pack.add_argument("--start-page", type=int, default=1, help="起始页。默认 1")
    p_pack.add_argument("--chunk-by", choices=["char", "page"], default="page", help="分块方式，默认 page")
    p_pack.add_argument("--chunk-size", type=int, default=2000, help="字符分块大小。默认 2000")
    p_pack.add_argument("--overlap", type=int, default=200, help="字符分块重叠。默认 200")
    p_pack.add_argument("--min-quality", type=float, default=0.5, help="最低质量评分门槛，默认 0.5")

    # --- batch-knowledge ---
    p_batch_pack = subparsers.add_parser("batch-knowledge", help="批量生成知识包")
    p_batch_pack.add_argument("dir", type=str, help="输入目录")
    p_batch_pack.add_argument("--parser", default="auto", help="解析器，默认 auto")
    p_batch_pack.add_argument("--out-dir", type=Path, default=None, help="输出目录（默认 <目录>/knowledge_packs）")
    p_batch_pack.add_argument("--ext", default=".pdf,.docx,.pptx,.xlsx,.html", help="扩展名，逗号分隔")
    p_batch_pack.add_argument("--recursive", action="store_true", help="递归扫描目录")
    p_batch_pack.add_argument("--max-pages", type=int, default=30, help="前 N 页，0=全量。默认 30")
    p_batch_pack.add_argument("--start-page", type=int, default=1, help="起始页。默认 1")
    p_batch_pack.add_argument("--chunk-by", choices=["char", "page"], default="page", help="分块方式，默认 page")
    p_batch_pack.add_argument("--chunk-size", type=int, default=2000, help="字符分块大小。默认 2000")
    p_batch_pack.add_argument("--overlap", type=int, default=200, help="字符分块重叠。默认 200")
    p_batch_pack.add_argument("--min-quality", type=float, default=0.5, help="最低质量评分门槛，默认 0.5")

    # --- batch-customer-pack ---
    p_batch_customer = subparsers.add_parser("batch-customer-pack", help="批量生成 PDF 客户交付包并写总索引")
    p_batch_customer.add_argument("dir", type=str, help="输入目录")
    p_batch_customer.add_argument("--out-dir", type=Path, default=None, help="输出目录（默认 <目录>/customer_packs）")
    p_batch_customer.add_argument("--ext", default=".pdf", help="扩展名，逗号分隔。默认 .pdf")
    p_batch_customer.add_argument("--recursive", action="store_true", help="递归扫描目录")
    p_batch_customer.add_argument("--max-pages", type=int, default=0, help="正式文本投票页数，0=全量。默认 0")
    p_batch_customer.add_argument("--start-page", type=int, default=1, help="起始页。默认 1")
    p_batch_customer.add_argument("--probe-max-pages", type=int, default=1, help="预检体检页数，0=全量。默认 1")
    p_batch_customer.add_argument("--table-pages", default="all", help="表格提取页码范围，如 1-5、all。默认 all")
    p_batch_customer.add_argument("--layout-max-pages", type=int, default=30, help="layout 输出页数，0=全量。默认 30")
    p_batch_customer.add_argument("--parsers", default=None,
                                  help=f"逗号分隔解析器。默认全部 PDF 解析器：{','.join(get_parsers_for_format('pdf'))}")
    p_batch_customer.add_argument("--profile", choices=profile_choices, default="auto",
                                  help="投票领域加权 profile。invoice 会校验加权；其它业务 profile 会做轻量结构化抽取。默认 auto")
    p_batch_customer.add_argument("--min-quality", type=float, default=0.5, help="最低质量评分门槛，默认 0.5")
    p_batch_customer.add_argument("--fail-on-bad", action="store_true", help="胜出结果标签为 empty/bad 时返回失败")
    p_batch_customer.add_argument("--similarity-chars", type=int, default=50000, help="投票共识相似度比较字符数")
    p_batch_customer.add_argument("--timeout", type=float, default=None,
                                  help="单个解析器最大运行秒数；超时会杀掉子进程并标记 timeout。默认不限制")
    p_batch_customer.add_argument("--parser-health-cache", action="store_true",
                                  help="启用解析器健康缓存，跳过 24 小时内同文件同页段近期失败或超时的解析器")
    p_batch_customer.add_argument("--health-cache", default=None,
                                  help="解析器健康缓存 JSON 路径；默认输出目录下 .parser_health_cache.json")
    p_batch_customer.add_argument("--no-tables", action="store_true", help="不提取 PDF 表格")
    p_batch_customer.add_argument("--no-layout", action="store_true", help="不输出 layout.json/page_map.json")
    p_batch_customer.add_argument("--no-field-fusion", action="store_true",
                                  help="关闭发票字段级投票融合，customer_best 只使用整份投票胜出解析器字段")
    p_batch_customer.add_argument("--no-field-evidence", action="store_true", help="不生成字段截图证据包")
    p_batch_customer.add_argument("--no-review-html", action="store_true", help="不生成可视化 HTML 审阅页")
    p_batch_customer.add_argument("--result-cache", action="store_true",
                                  help="启用解析结果缓存，同一 PDF 同一 parser/page 范围不重复重解析")
    p_batch_customer.add_argument("--result-cache-dir", default=None,
                                  help="解析结果缓存目录；默认输出目录下 .parser_result_cache")
    p_batch_customer.add_argument("--parallel", action="store_true", help="并行执行解析器")

    # --- qa ---
    p_qa = subparsers.add_parser("qa", help="基于知识包、chunks 或文档执行本地抽取式问答")
    p_qa.add_argument("source", type=Path, help="知识包目录、chunks.jsonl/json 或文档路径")
    p_qa.add_argument("question", help="问题")
    p_qa.add_argument("--parser", default="auto", help="直接输入文档时使用的解析器，默认 auto")
    p_qa.add_argument("--max-pages", type=int, default=30, help="直接输入文档时前 N 页，0=全量。默认 30")
    p_qa.add_argument("--start-page", type=int, default=1, help="直接输入 PDF 时起始页。默认 1")
    p_qa.add_argument("--chunk-by", choices=["char", "page"], default="page", help="直接输入 PDF 时分块方式，默认 page")
    p_qa.add_argument("--chunk-size", type=int, default=2000, help="字符分块大小。默认 2000")
    p_qa.add_argument("--overlap", type=int, default=200, help="字符分块重叠。默认 200")
    p_qa.add_argument("--top-k", type=int, default=5, help="检索候选块数量。默认 5")
    p_qa.add_argument("--answer-sentences", type=int, default=5, help="最多输出答案句数。默认 5")
    p_qa.add_argument("--format", choices=["json", "md", "all"], default="md", help="输出格式。默认 md")
    p_qa.add_argument("-o", "--output", default=None, help="输出文件路径；--format all 时为输出目录")

    # --- diff-docs ---
    p_diff = subparsers.add_parser("diff-docs", help="对比两个文档版本的文本、分类和可选发票字段差异")
    p_diff.add_argument("left", type=Path, help="左侧/旧版本文档")
    p_diff.add_argument("right", type=Path, help="右侧/新版本文档")
    p_diff.add_argument("--parser", default="auto", help="解析器，默认 auto")
    p_diff.add_argument("--max-pages", type=int, default=30, help="前 N 页，0=全量。默认 30")
    p_diff.add_argument("--start-page", type=int, default=1, help="起始页。默认 1")
    p_diff.add_argument("--profile", choices=["auto", "invoice", "none"], default="auto", help="结构化字段差异 profile。默认 auto")
    p_diff.add_argument("--max-diff-lines", type=int, default=200, help="最多输出 unified diff 行数。默认 200")
    p_diff.add_argument("--format", choices=["json", "md", "all"], default="md", help="输出格式。默认 md")
    p_diff.add_argument("-o", "--output", default=None, help="输出文件路径；--format all 时为输出目录")

    return parser


def main() -> int:
    if len(sys.argv) > 2 and sys.argv[1] == "__extract-one":
        return cmd_extract_one_internal(Path(sys.argv[2]))

    parser = build_parser()

    # 向后兼容：无子命令时等同于 compare
    if len(sys.argv) > 1 and sys.argv[1] not in (
        "compare", "vote", "customer-pack", "batch-customer-pack", "probe", "convert", "batch", "scan-dir", "tables", "table-vote", "init-golden", "eval-golden", "doctor", "metadata",
        "chunk", "render-pages", "ocr", "auto", "extract-fields", "export-xlsx",
        "layout-json", "verify-fields", "classify", "knowledge-pack", "batch-knowledge",
        "qa", "diff-docs",
        "--version", "-h", "--help", "__extract-one",
    ):
        sys.argv.insert(1, "compare")

    args = parser.parse_args()

    if args.command is None:
        parser.print_help()
        return 0

    dispatch = {
        "compare": cmd_compare,
        "vote": cmd_vote,
        "customer-pack": cmd_customer_pack,
        "batch-customer-pack": cmd_batch_customer_pack,
        "probe": cmd_probe,
        "convert": cmd_convert,
        "batch": cmd_batch,
        "scan-dir": cmd_scan_dir,
        "tables": cmd_tables,
        "table-vote": cmd_table_vote,
        "init-golden": cmd_init_golden,
        "eval-golden": cmd_eval_golden,
        "doctor": cmd_doctor,
        "metadata": cmd_metadata,
        "chunk": cmd_chunk,
        "render-pages": cmd_render_pages,
        "ocr": cmd_ocr,
        "auto": cmd_auto,
        "extract-fields": cmd_extract_fields,
        "export-xlsx": cmd_export_xlsx,
        "layout-json": cmd_layout_json,
        "verify-fields": cmd_verify_fields,
        "classify": cmd_classify,
        "knowledge-pack": cmd_knowledge_pack,
        "batch-knowledge": cmd_batch_knowledge,
        "qa": cmd_qa,
        "diff-docs": cmd_diff_docs,
    }

    handler = dispatch.get(args.command)
    if handler is None:
        parser.print_help()
        return 1

    return handler(args)


if __name__ == "__main__":
    raise SystemExit(main())
